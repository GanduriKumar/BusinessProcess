"""
Generate SIP MVP Technical Architecture PowerPoint Infographic.
Run: python generate_sip_pptx.py
Output: SIP_MVP_Architecture_Infographic.pptx (same folder)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colours ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG      = RGBColor(0x1E, 0x29, 0x3B)
CARD_BORDER  = RGBColor(0x33, 0x41, 0x55)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
DIM_GRAY     = RGBColor(0x64, 0x74, 0x8B)
BLUE         = RGBColor(0x3B, 0x82, 0xF6)
CYAN         = RGBColor(0x06, 0xB6, 0xD4)
GREEN        = RGBColor(0x10, 0xB9, 0x81)
PURPLE       = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE       = RGBColor(0xF5, 0x9E, 0x0B)
PINK         = RGBColor(0xEC, 0x48, 0x99)
RED          = RGBColor(0xEF, 0x44, 0x44)
TEAL         = RGBColor(0x14, 0xB8, 0xA6)

LIGHT_BLUE   = RGBColor(0x93, 0xC5, 0xFD)
LIGHT_CYAN   = RGBColor(0x67, 0xE8, 0xF9)
LIGHT_GREEN  = RGBColor(0x6E, 0xE7, 0xB7)
LIGHT_PURPLE = RGBColor(0xC4, 0xB5, 0xFD)
LIGHT_PINK   = RGBColor(0xF9, 0xA8, 0xD4)
LIGHT_RED    = RGBColor(0xFC, 0xA5, 0xA5)
LIGHT_ORANGE = RGBColor(0xFC, 0xD3, 0x4D)

# ── Helpers ──────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.08
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=10, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, alignment=PP_ALIGN.CENTER):
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, col, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = col
        p.font.bold = bld
        p.alignment = alignment
        p.space_after = Pt(1)
    return txBox

def add_arrow_connector(slide, x1, y1, x2, y2, color=MID_GRAY, width=Pt(1.5)):
    """Add a line with an arrowhead."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # End arrow
    connector.line.end_marker_style = 2  # Triangle
    return connector


# ── Slide dimensions (widescreen 16:9) ──────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

# ╔═══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 1 — Main Architecture Diagram                             ║
# ╚═══════════════════════════════════════════════════════════════════╝
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
set_slide_bg(slide1, BG_DARK)

# ── Title ──
add_text_box(slide1, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.45),
             "Support Intelligence Platform (SIP) — MVP Technical Architecture",
             font_size=20, color=LIGHT_CYAN, bold=True)
add_text_box(slide1, Inches(0.3), Inches(0.48), Inches(12.7), Inches(0.3),
             "Layered Service Architecture  |  Phase 1 Scope  |  React Frontend · Python Backend · REST APIs",
             font_size=10, color=MID_GRAY)

# ── Layer backgrounds ──
# Users
add_rounded_rect(slide1, Inches(0.2), Inches(0.85), Inches(12.9), Inches(0.7),
                 RGBColor(0x15, 0x20, 0x3A), BLUE, Pt(0.8))
add_text_box(slide1, Inches(0.35), Inches(0.87), Inches(1.8), Inches(0.22),
             "USERS & ROLES", font_size=8, color=DIM_GRAY, bold=True, alignment=PP_ALIGN.LEFT)

# User role cards
user_roles = [
    ("Support Engineer", "Investigate & Resolve"),
    ("Support Manager", "Monitor & Escalate"),
    ("Platform Admin", "Configure & Manage"),
    ("Knowledge Admin", "Curate & Validate"),
    ("Security Reviewer", "Audit & Comply"),
]
ux = Inches(0.5)
for role, desc in user_roles:
    r = add_rounded_rect(slide1, ux, Inches(1.07), Inches(2.2), Inches(0.38),
                         CARD_BG, BLUE, Pt(1))
    add_multi_text(slide1, ux + Inches(0.05), Inches(1.07), Inches(2.1), Inches(0.38), [
        (role, 9, LIGHT_BLUE, True),
        (desc, 7, DIM_GRAY, False),
    ])
    ux += Inches(2.45)

# ── Frontend layer ──
add_rounded_rect(slide1, Inches(0.2), Inches(1.7), Inches(12.9), Inches(1.05),
                 RGBColor(0x12, 0x22, 0x38), CYAN, Pt(0.8))
add_text_box(slide1, Inches(0.35), Inches(1.72), Inches(2), Inches(0.22),
             "FRONTEND — React SPA", font_size=8, color=DIM_GRAY, bold=True, alignment=PP_ALIGN.LEFT)

fe_components = [
    ("Ticket Workbench", "Summary, context\nannotations & actions"),
    ("Investigation Space", "Hypotheses, evidence\nnext-best actions"),
    ("Knowledge Panel", "AI recs, KB articles\nsimilar cases"),
    ("Generative UI", "Schema-driven\nrole-adaptive layout"),
    ("Admin Console", "Connectors, mappings\nprompts & policies"),
    ("Manager Dashboard", "SLA risk, queue\nadoption metrics"),
    ("Feedback & Approval", "Human-in-the-loop\nwrite-back approval"),
]
fx = Inches(0.35)
fw = Inches(1.7)
for name, desc in fe_components:
    add_rounded_rect(slide1, fx, Inches(1.98), fw, Inches(0.65),
                     CARD_BG, CYAN, Pt(0.8))
    add_multi_text(slide1, fx + Inches(0.03), Inches(1.98), fw - Inches(0.06), Inches(0.65), [
        (name, 8.5, LIGHT_CYAN, True),
        (desc, 6.5, DIM_GRAY, False),
    ])
    fx += fw + Inches(0.1)

# ── Downward arrows: Users → Frontend ──
for i in range(5):
    cx = Inches(0.5) + Inches(2.45) * i + Inches(1.1)
    add_arrow_connector(slide1, cx, Inches(1.45), cx, Inches(1.7), BLUE, Pt(1))

# ── API Gateway ──
add_rounded_rect(slide1, Inches(0.2), Inches(2.9), Inches(12.9), Inches(0.45),
                 RGBColor(0x1A, 0x1E, 0x2E), ORANGE, Pt(0.8))
add_text_box(slide1, Inches(0.35), Inches(2.93), Inches(12.5), Inches(0.2),
             "API GATEWAY  —  /api/v1/*   |   Auth · Routing · Validation · Rate-Limiting · Health/Readiness   |   16 REST Resource Areas",
             font_size=9, color=LIGHT_ORANGE, bold=True)
add_text_box(slide1, Inches(0.35), Inches(3.12), Inches(12.5), Inches(0.2),
             "/auth  /tickets  /tickets/{id}/recommendations  /connectors  /knowledge  /ui-composition  /audit  /prompts  /policies  /health",
             font_size=7, color=DIM_GRAY)

# Arrow: Frontend → API GW
add_arrow_connector(slide1, Inches(6.85), Inches(2.63), Inches(6.85), Inches(2.9), CYAN, Pt(1.5))
add_text_box(slide1, Inches(6.95), Inches(2.68), Inches(1.2), Inches(0.2),
             "REST / HTTPS", font_size=6.5, color=DIM_GRAY, alignment=PP_ALIGN.LEFT)

# ── Backend Services ──
add_rounded_rect(slide1, Inches(0.2), Inches(3.5), Inches(12.9), Inches(2.05),
                 RGBColor(0x14, 0x1A, 0x30), PURPLE, Pt(0.8))
add_text_box(slide1, Inches(0.35), Inches(3.52), Inches(2.5), Inches(0.22),
             "BACKEND SERVICES — Python", font_size=8, color=DIM_GRAY, bold=True, alignment=PP_ALIGN.LEFT)

# Row 1: Main services
svc_row1 = [
    ("Ticket Ingestion", "Connector abstraction\nPoll · Webhook · Import\n→ Canonical Model", PURPLE),
    ("Workflow Orchestration", "Ticket lifecycle stages\nSync & async coordination\nState transitions", PURPLE),
    ("Recommendation Service", "AI provider abstraction\nVersioned prompts\nSummarize · Classify · Rec", PURPLE),
    ("UI Composition", "Layout composition rules\nContext→component map\nGenUI payloads", PURPLE),
    ("Knowledge Service", "Ingest · Index · Retrieve\nKB, runbooks, change recs\nQuality feedback", PURPLE),
]
sx = Inches(0.4)
sw = Inches(2.35)
for name, desc, border in svc_row1:
    add_rounded_rect(slide1, sx, Inches(3.78), sw, Inches(0.78),
                     CARD_BG, border, Pt(1))
    add_multi_text(slide1, sx + Inches(0.05), Inches(3.78), sw - Inches(0.1), Inches(0.78), [
        (name, 8.5, LIGHT_PURPLE, True),
        (desc, 6.5, DIM_GRAY, False),
    ])
    sx += sw + Inches(0.12)

# Row 2: Supporting services
svc_row2 = [
    ("Connector Admin", "Config & field mapping\nSync policy & health", PURPLE),
    ("Audit & Governance", "Append-only audit logs\nDecision traceability", RED),
    ("Feedback Capture", "Human feedback on recs\nApproval & quality signals", PURPLE),
    ("Auth & Identity", "Enterprise SSO · RBAC\nSession & policy controls", RED),
    ("Configuration", "Prompt & workflow templates\nFeature flags · Policies", PURPLE),
]
sx = Inches(0.4)
for name, desc, border in svc_row2:
    add_rounded_rect(slide1, sx, Inches(4.65), sw, Inches(0.72),
                     CARD_BG, border, Pt(1))
    add_multi_text(slide1, sx + Inches(0.05), Inches(4.66), sw - Inches(0.1), Inches(0.72), [
        (name, 8.5, LIGHT_RED if border == RED else LIGHT_PURPLE, True),
        (desc, 6.5, DIM_GRAY, False),
    ])
    sx += sw + Inches(0.12)

# Inter-service arrows (horizontal)
for i in range(4):
    x_start = Inches(0.4) + (Inches(2.35) + Inches(0.12)) * (i + 1) - Inches(0.1)
    x_end = x_start + Inches(0.1)
    add_arrow_connector(slide1, x_start - Inches(0.02), Inches(4.17),
                        x_end + Inches(0.02), Inches(4.17), PURPLE, Pt(0.8))

# Arrows: API GW → Backend (fan out)
for i in range(5):
    gw_x = Inches(2) + Inches(2) * i
    svc_x = Inches(0.4) + (Inches(2.35) + Inches(0.12)) * i + Inches(1.17)
    add_arrow_connector(slide1, gw_x, Inches(3.35), svc_x, Inches(3.78), ORANGE, Pt(0.8))

# ── Data & Storage Layer ──
add_rounded_rect(slide1, Inches(0.2), Inches(5.6), Inches(12.9), Inches(0.85),
                 RGBColor(0x10, 0x1C, 0x28), GREEN, Pt(0.8))
add_text_box(slide1, Inches(0.35), Inches(5.62), Inches(2.5), Inches(0.22),
             "DATA & STORAGE LAYER", font_size=8, color=DIM_GRAY, bold=True, alignment=PP_ALIGN.LEFT)

data_stores = [
    ("Operational Database", "Tickets · Users · Roles\nWorkflow · Config", GREEN),
    ("Vector / Search Index", "Semantic retrieval\nKB & ticket index", GREEN),
    ("Audit Log Store", "Append-only, tamper-evident\nDecisions & events", GREEN),
    ("Secret Store", "API keys, creds\nManaged vault", GREEN),
    ("Cache Layer", "Session & ticket cache\nPrompt templates", GREEN),
]
dx = Inches(0.4)
dw = Inches(2.35)
for name, desc, border in data_stores:
    add_rounded_rect(slide1, dx, Inches(5.85), dw, Inches(0.5),
                     CARD_BG, border, Pt(1))
    add_multi_text(slide1, dx + Inches(0.05), Inches(5.84), dw - Inches(0.1), Inches(0.5), [
        (name, 8, LIGHT_GREEN, True),
        (desc, 6, DIM_GRAY, False),
    ])
    dx += dw + Inches(0.12)

# Arrows: Backend → Data
for i in range(5):
    cx = Inches(0.4) + (Inches(2.35) + Inches(0.12)) * i + Inches(1.17)
    add_arrow_connector(slide1, cx, Inches(5.37), cx, Inches(5.6), GREEN, Pt(0.8))

# ── External Systems (bottom) ──
add_rounded_rect(slide1, Inches(0.2), Inches(6.55), Inches(12.9), Inches(0.78),
                 RGBColor(0x18, 0x14, 0x22), PINK, Pt(0.8))
add_text_box(slide1, Inches(0.35), Inches(6.57), Inches(3), Inches(0.22),
             "EXTERNAL SYSTEMS & INTEGRATIONS", font_size=8, color=DIM_GRAY, bold=True, alignment=PP_ALIGN.LEFT)

ext_systems = [
    ("Enterprise Ticketing", "ServiceNow · Jira SM\nSalesforce SC · Custom\nBidirectional R/W"),
    ("AI / LLM Providers", "Provider abstraction\nModel-agnostic\nGoverned & auditable"),
    ("Identity Provider", "Enterprise SSO\nSAML / OIDC\nRole provisioning"),
    ("Knowledge Sources", "KB articles · Runbooks\nChange records\nEnvironment notes"),
    ("Observability Stack", "Logging · Metrics\nTracing · Dashboards\nAlerts"),
]
ex = Inches(0.4)
ew = Inches(2.35)
for name, desc in ext_systems:
    add_rounded_rect(slide1, ex, Inches(6.75), ew, Inches(0.5),
                     CARD_BG, PINK, Pt(1))
    add_multi_text(slide1, ex + Inches(0.05), Inches(6.74), ew - Inches(0.1), Inches(0.5), [
        (name, 8, LIGHT_PINK, True),
        (desc, 6, DIM_GRAY, False),
    ])
    ex += ew + Inches(0.12)

# Arrows: Data → External
for i in range(5):
    cx = Inches(0.4) + (Inches(2.35) + Inches(0.12)) * i + Inches(1.17)
    add_arrow_connector(slide1, cx, Inches(6.35), cx, Inches(6.55), PINK, Pt(0.8))


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 2 — Workflow & Data Flow                                   ║
# ╚═══════════════════════════════════════════════════════════════════╝
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

add_text_box(slide2, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.45),
             "SIP — Ticket Lifecycle Workflow & End-to-End Data Flow",
             font_size=20, color=LIGHT_CYAN, bold=True)
add_text_box(slide2, Inches(0.3), Inches(0.55), Inches(12.7), Inches(0.3),
             "From ticket intake through AI-assisted enrichment to human-approved write-back  |  Phase 1 MVP",
             font_size=10, color=MID_GRAY)

# Workflow steps
wf_steps = [
    ("1", "Ticket Intake", "Poll / Webhook / Import\nfrom Ticketing System", PINK),
    ("2", "Normalization", "Map to Canonical\nTicket Model", PURPLE),
    ("3", "Classification", "AI-assisted category\n& priority assignment", PURPLE),
    ("4", "Enrichment", "Context assembly\nKB + history + env data", CYAN),
    ("5", "Recommendation", "AI hypotheses &\nnext-best actions", ORANGE),
    ("6", "Human Review", "Approve / modify\nbefore write-back", BLUE),
    ("7", "Write-back", "Update source\nticketing system", GREEN),
    ("8", "Feedback", "Quality signals &\nknowledge capture", TEAL),
]

# Draw the workflow as a horizontal chain
wx_start = Inches(0.3)
wy = Inches(1.2)
ww = Inches(1.42)
wh = Inches(1.1)
wgap = Inches(0.18)

for i, (num, name, desc, color) in enumerate(wf_steps):
    wx = wx_start + i * (ww + wgap)

    # Number circle
    circ = slide2.shapes.add_shape(MSO_SHAPE.OVAL, wx + Inches(0.55), wy, Inches(0.32), Inches(0.32))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Card
    add_rounded_rect(slide2, wx, wy + Inches(0.4), ww, wh - Inches(0.4),
                     CARD_BG, color, Pt(1.2))
    add_multi_text(slide2, wx + Inches(0.04), wy + Inches(0.42), ww - Inches(0.08), wh - Inches(0.4), [
        (name, 10, WHITE, True),
        (desc, 7.5, DIM_GRAY, False),
    ])

    # Arrow between steps
    if i < len(wf_steps) - 1:
        ax1 = wx + ww
        ax2 = wx + ww + wgap
        ay = wy + Inches(0.85)
        add_arrow_connector(slide2, ax1, ay, ax2, ay, color, Pt(2))

# ── Detailed data flow section ──
add_text_box(slide2, Inches(0.3), Inches(2.6), Inches(12.7), Inches(0.3),
             "▾  Detailed Data Flow by Stage", font_size=12, color=LIGHT_CYAN, bold=True, alignment=PP_ALIGN.LEFT)

# Flow detail rows
detail_rows = [
    ("Intake", "External System → Connector → Polling/Webhook → Raw Ticket Payload → Ingestion Service → Canonical Ticket Model → Operational DB", PINK),
    ("Enrich", "Canonical Ticket → Workflow Orchestration → Knowledge Service (Vector Search) → Context Assembly → Enriched Ticket View", CYAN),
    ("Recommend", "Enriched Ticket → Recommendation Service → AI Provider (LLM) → Prompt Template → Hypotheses + Next-Best Actions → UI Composition", ORANGE),
    ("Review", "Generative UI → Ticket Workbench → Support Engineer reviews → Approves / Modifies → Feedback Capture Service → Audit Log", BLUE),
    ("Write-back", "Approved Update → Connector Admin → Write-back to Source Ticketing System → Status/Comment Sync → Audit Record", GREEN),
    ("Learn", "Feedback Signals → Knowledge Service → Index Update → Improved future retrieval & recommendations → Quality Loop", TEAL),
]

dy = Inches(3.0)
for label, flow, color in detail_rows:
    # Label box
    add_rounded_rect(slide2, Inches(0.3), dy, Inches(1.0), Inches(0.36),
                     CARD_BG, color, Pt(1))
    add_text_box(slide2, Inches(0.3), dy + Inches(0.02), Inches(1.0), Inches(0.32),
                 label, font_size=9, color=color, bold=True)
    # Arrow
    add_arrow_connector(slide2, Inches(1.35), dy + Inches(0.18),
                        Inches(1.6), dy + Inches(0.18), color, Pt(1))
    # Flow text
    add_text_box(slide2, Inches(1.65), dy, Inches(11.3), Inches(0.36),
                 flow, font_size=8, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)
    dy += Inches(0.45)

# ── Architecture Principles strip ──
add_text_box(slide2, Inches(0.3), Inches(5.85), Inches(12.7), Inches(0.3),
             "▾  Architecture Design Principles", font_size=12, color=LIGHT_CYAN, bold=True, alignment=PP_ALIGN.LEFT)

principles = [
    ("Separation of Concerns", "UI, orchestration, intelligence,\nintegration are modular"),
    ("Provider Abstraction", "Not coupled to one LLM,\nvector store, or ticketing system"),
    ("Configurability", "Connectors, prompts, mappings,\nworkflows — all configurable"),
    ("Human Accountability", "All AI outputs advisory;\napproval before write-back"),
    ("Enterprise Ready", "Security, audit, privacy,\nobservability first-class"),
    ("Incremental Delivery", "Phased rollout by module,\nconnector, and use case"),
]

px = Inches(0.3)
pw = Inches(2.0)
for name, desc in principles:
    add_rounded_rect(slide2, px, Inches(6.2), pw, Inches(0.65),
                     CARD_BG, CYAN, Pt(0.8))
    add_multi_text(slide2, px + Inches(0.06), Inches(6.2), pw - Inches(0.12), Inches(0.65), [
        (name, 8.5, LIGHT_CYAN, True),
        (desc, 6.5, DIM_GRAY, False),
    ])
    px += pw + Inches(0.13)

# ── Footer ──
add_text_box(slide2, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3),
             "SIP MVP Architecture  |  Generated from Technical Requirements Document  |  Phase 1",
             font_size=7, color=DIM_GRAY)


# ╔═══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 3 — API Endpoints & Service Map                           ║
# ╚═══════════════════════════════════════════════════════════════════╝
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

add_text_box(slide3, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.45),
             "SIP — REST API Map & Service Ownership",
             font_size=20, color=LIGHT_CYAN, bold=True)
add_text_box(slide3, Inches(0.3), Inches(0.55), Inches(12.7), Inches(0.3),
             "16 versioned API resource areas mapped to owning backend services  |  All endpoints under /api/v1/*",
             font_size=10, color=MID_GRAY)

# API → Service mapping table
api_map = [
    ("/api/v1/auth", "Auth & Identity", "Login, session, token management", BLUE),
    ("/api/v1/users", "Auth & Identity", "User profiles and preferences", BLUE),
    ("/api/v1/roles", "Auth & Identity", "Role definitions and assignments", BLUE),
    ("/api/v1/tickets", "Ticket Ingestion + Workflow", "CRUD, list, filter, sort tickets", PURPLE),
    ("/api/v1/tickets/{id}/context", "Workflow Orchestration", "Enriched context for a ticket", PURPLE),
    ("/api/v1/tickets/{id}/recommendations", "Recommendation Service", "AI-generated recommendations", ORANGE),
    ("/api/v1/tickets/{id}/feedback", "Feedback Capture", "Submit feedback on recommendations", TEAL),
    ("/api/v1/connectors", "Connector Admin", "Manage ticketing connectors", PINK),
    ("/api/v1/connectors/{id}/mappings", "Connector Admin", "Field mapping configuration", PINK),
    ("/api/v1/connectors/{id}/sync-jobs", "Connector Admin", "Sync job status and history", PINK),
    ("/api/v1/knowledge", "Knowledge Service", "Search and retrieve knowledge", GREEN),
    ("/api/v1/prompts", "Configuration Service", "Manage prompt templates", PURPLE),
    ("/api/v1/policies", "Audit & Governance", "Governance policy management", RED),
    ("/api/v1/ui-composition", "UI Composition Service", "Generative UI layout metadata", CYAN),
    ("/api/v1/audit", "Audit & Governance", "Audit trail queries", RED),
    ("/api/v1/health", "API Gateway", "Health and readiness checks", ORANGE),
]

# Table header
ty = Inches(1.05)
add_rounded_rect(slide3, Inches(0.3), ty, Inches(3.5), Inches(0.3), RGBColor(0x1A, 0x25, 0x3D))
add_text_box(slide3, Inches(0.4), ty + Inches(0.02), Inches(3.3), Inches(0.26),
             "API Endpoint", font_size=8.5, color=LIGHT_CYAN, bold=True, alignment=PP_ALIGN.LEFT)

add_rounded_rect(slide3, Inches(3.85), ty, Inches(3.2), Inches(0.3), RGBColor(0x1A, 0x25, 0x3D))
add_text_box(slide3, Inches(3.95), ty + Inches(0.02), Inches(3.0), Inches(0.26),
             "Owning Service", font_size=8.5, color=LIGHT_CYAN, bold=True, alignment=PP_ALIGN.LEFT)

add_rounded_rect(slide3, Inches(7.1), ty, Inches(6.1), Inches(0.3), RGBColor(0x1A, 0x25, 0x3D))
add_text_box(slide3, Inches(7.2), ty + Inches(0.02), Inches(5.9), Inches(0.26),
             "Description", font_size=8.5, color=LIGHT_CYAN, bold=True, alignment=PP_ALIGN.LEFT)

# Table rows
ty += Inches(0.34)
for endpoint, service, desc, color in api_map:
    row_bg = RGBColor(0x16, 0x1E, 0x33) if api_map.index((endpoint, service, desc, color)) % 2 == 0 else RGBColor(0x12, 0x1A, 0x2E)
    add_rounded_rect(slide3, Inches(0.3), ty, Inches(12.8), Inches(0.3), row_bg)

    add_text_box(slide3, Inches(0.4), ty + Inches(0.02), Inches(3.3), Inches(0.26),
                 endpoint, font_size=7.5, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Color dot for service
    dot = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.95), ty + Inches(0.1), Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_text_box(slide3, Inches(4.1), ty + Inches(0.02), Inches(2.9), Inches(0.26),
                 service, font_size=7.5, color=color, bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide3, Inches(7.2), ty + Inches(0.02), Inches(5.9), Inches(0.26),
                 desc, font_size=7.5, color=MID_GRAY, alignment=PP_ALIGN.LEFT)
    ty += Inches(0.3)

# Deployment strip at bottom
add_text_box(slide3, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.3),
             "▾  Deployment & Non-Functional Requirements", font_size=11, color=LIGHT_CYAN, bold=True, alignment=PP_ALIGN.LEFT)

deploy_items = [
    "Containerized\nDeployment",
    "CI / CD\nPipelines",
    "Independent\nScaling",
    "Env-Specific\nConfiguration",
    "Private\nNetworking",
    "Graceful\nDegradation",
    "Multi-Environment\n(Dev/Test/Stg/Prod)",
]
dpx = Inches(0.3)
dpw = Inches(1.73)
for item in deploy_items:
    add_rounded_rect(slide3, dpx, Inches(6.85), dpw, Inches(0.45),
                     CARD_BG, CARD_BORDER, Pt(0.8))
    add_text_box(slide3, dpx + Inches(0.04), Inches(6.86), dpw - Inches(0.08), Inches(0.43),
                 item, font_size=7.5, color=MID_GRAY, bold=True)
    dpx += dpw + Inches(0.1)

# Footer
add_text_box(slide3, Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3),
             "SIP MVP Architecture  |  REST API Map  |  Phase 1",
             font_size=7, color=DIM_GRAY)


# ── Save ─────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "SIP_MVP_Architecture_Infographic.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
