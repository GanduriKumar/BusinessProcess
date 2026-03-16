from pptx import Presentation

prs = Presentation(r"c:\Users\kumar.gn\HCLProjects\BusinessProcess\docs\input\TechVDU_Transformation_MasterSlide.pptx")
theme = prs.slide_masters[0].element
ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
clrScheme = theme.find(".//a:clrScheme", ns)
if clrScheme is not None:
    print("Color scheme:", clrScheme.attrib.get("name", ""))
    for child in clrScheme:
        tag = child.tag.split("}")[1] if "}" in child.tag else child.tag
        for ac in child:
            at = ac.tag.split("}")[1] if "}" in ac.tag else ac.tag
            val = ac.attrib.get("val", "") or ac.attrib.get("lastClr", "")
            print("  %s: %s=%s" % (tag, at, val))

master = prs.slide_masters[0]
print("\nMaster shapes:", len(master.shapes))
for sh in master.shapes:
    print("  %s name=%r l=%.2f t=%.2f w=%.2f h=%.2f" % (
        sh.shape_type, sh.name, sh.left/914400, sh.top/914400,
        sh.width/914400, sh.height/914400))

# Check a few layouts with background shapes
for idx in [41, 42, 46, 102, 107]:
    lay = prs.slide_layouts[idx]
    print("\nLayout %d (%s) shapes:" % (idx, lay.name))
    for sh in lay.shapes:
        is_ph = False
        try:
            _ = sh.placeholder_format
            is_ph = True
        except:
            pass
        print("  %s%s name=%r l=%.2f t=%.2f w=%.2f h=%.2f" % (
            "PH " if is_ph else "",
            sh.shape_type, sh.name, sh.left/914400, sh.top/914400,
            sh.width/914400, sh.height/914400))
