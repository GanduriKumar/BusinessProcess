"""
Generate GenAI Service Transformation PPTX — matching the HTML slide deck styling.
Splits content across 3 slides for legibility with larger fonts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── HCL Brand Colours ──
HCL_DARK     = RGBColor(0x00, 0x2B, 0x5C)
HCL_BLUE     = RGBColor(0x00, 0x60, 0xA9)
HCL_ACCENT   = RGBColor(0x00, 0xA3, 0xE0)
HCL_GREEN    = RGBColor(0x00, 0x87, 0x5A)
HCL_ORANGE   = RGBColor(0xE8, 0x77, 0x22)
HCL_PURPLE   = RGBColor(0x6C, 0x3F, 0xA5)
HCL_TEAL     = RGBColor(0x00, 0x85, 0x7C)
HCL_RED      = RGBColor(0xC9, 0x24, 0x3F)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50      = RGBColor(0xF8, 0xFA, 0xFC)
GRAY_100     = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_200     = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_400     = RGBColor(0x94, 0xA3, 0xB8)
GRAY_500     = RGBColor(0x64, 0x74, 0x8B)
GRAY_600     = RGBColor(0x47, 0x55, 0x69)
GRAY_800     = RGBColor(0x1E, 0x29, 0x3B)

# Category accent tints (light bg behind icon)
CAT_BG = {
    "quality": RGBColor(0xD9, 0xE8, 0xF5),
    "support": RGBColor(0xD5, 0xEC, 0xEA),
    "dev":     RGBColor(0xE6, 0xDD, 0xF2),
    "govern":  RGBColor(0xD5, 0xEC, 0xE2),
}
CAT_ACCENT = {
    "quality": HCL_BLUE,
    "support": HCL_TEAL,
    "dev":     HCL_PURPLE,
    "govern":  HCL_GREEN,
}

# ── Helpers ──
def set_bg(slide, color):
    bg = slide.background; f = bg.fill; f.solid(); f.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill, border=None, bw=Pt(0.75)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.06
    return s

def add_tb(slide, l, t, w, h, text, sz=10, color=GRAY_800, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tb

def add_ml(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    """lines = [(text, size, color, bold), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, sz, col, bld) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = col; p.font.bold = bld
        p.alignment = align; p.space_after = Pt(1)
    return tb

def add_dot(slide, l, t, color):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.09), Inches(0.09))
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()

def draw_gradient_bar(slide, l, t, w, h):
    """Draw the multi-colour accent bar across bottom of header."""
    seg_w = w // 4
    colors = [HCL_ACCENT, HCL_GREEN, HCL_ORANGE, HCL_PURPLE]
    for i, c in enumerate(colors):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l + seg_w * i, t, seg_w + (w - seg_w * 4 if i == 3 else 0), h)
        s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def draw_header(slide):
    """Draw the branded top bar — reused on every slide."""
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.95), HCL_DARK, None)
    draw_gradient_bar(slide, Inches(0), Inches(0.92), Inches(13.333), Inches(0.05))
    add_tb(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.45),
           "Accelerate Delivery with GenAI", sz=22, color=WHITE, bold=True)
    add_tb(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.35),
           "Battle-tested AI solutions that cut costs, raise quality, and compress timelines", sz=11, color=RGBColor(0xBB,0xDB,0xF0))
    add_tb(slide, Inches(11.2), Inches(0.25), Inches(1.8), Inches(0.4),
           "HCLTech", sz=16, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def draw_footer(slide, page_text=""):
    """Draw the bottom bar."""
    add_rect(slide, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45), GRAY_50, GRAY_200, Pt(0.5))
    add_tb(slide, Inches(0.5), Inches(7.12), Inches(8.5), Inches(0.3),
           "Enterprise-grade safety controls  |  Full audit trails  |  Human-approval gates  |  Zero-disruption deployment",
           sz=8.5, color=GRAY_600)
    # Legend
    lx = Inches(10.0)
    for label, color in [("Live", HCL_GREEN), ("Pilot", HCL_ORANGE), ("POC", HCL_ACCENT)]:
        add_dot(slide, lx, Inches(7.18), color)
        add_tb(slide, lx + Inches(0.12), Inches(7.12), Inches(0.6), Inches(0.3), label, sz=8, color=GRAY_500)
        lx += Inches(0.8)
    if page_text:
        add_tb(slide, Inches(12.3), Inches(7.12), Inches(0.9), Inches(0.3), page_text, sz=8, color=GRAY_400, align=PP_ALIGN.RIGHT)


def draw_card(slide, x, y, w, h, accent, title, desc, metric, status, status_color):
    """Draw a single solution card."""
    # Card background
    card = add_rect(slide, x, y, w, h, WHITE, GRAY_200, Pt(0.75))
    # Top accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.05))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = accent; stripe.line.fill.background()
    # Title
    add_tb(slide, x + Inches(0.14), y + Inches(0.1), w - Inches(0.28), Inches(0.32),
           title, sz=11.5, color=GRAY_800, bold=True)
    # Description
    add_tb(slide, x + Inches(0.14), y + Inches(0.4), w - Inches(0.28), Inches(0.7),
           desc, sz=9, color=GRAY_600)
    # Metric badge
    badge = add_rect(slide, x + Inches(0.14), y + h - Inches(0.36), w - Inches(0.28), Inches(0.26),
                     GRAY_50, GRAY_200, Pt(0.5))
    add_ml(slide, x + Inches(0.2), y + h - Inches(0.36), w - Inches(0.6), Inches(0.26), [
        ("▲ " + metric, 8, HCL_GREEN, True),
    ])
    # Status dot + label
    add_dot(slide, x + w - Inches(0.6), y + h - Inches(0.34), status_color)
    add_tb(slide, x + w - Inches(0.52), y + h - Inches(0.38), Inches(0.46), Inches(0.24),
           status, sz=7, color=GRAY_400)


# ═══════════════════════════════════════════════
# Build Presentation
# ═══════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════
# SLIDE 1 — Title + KPIs + Quality Eng + Support
# ═══════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, WHITE)
draw_header(s1)

# KPI strip
add_rect(s1, Inches(0), Inches(0.97), Inches(13.333), Inches(0.65), GRAY_50, GRAY_200, Pt(0.5))
kpis = [("15+", "Deployable\nSolutions"), ("25–30%", "QA Effort\nEliminated"), ("60%", "Faster\nMigrations"), ("~90%", "Bug-Fix\nAccuracy"), ("30%", "Compliance\nSavings")]
kx = Inches(0.7)
for num, label in kpis:
    add_tb(s1, kx, Inches(1.0), Inches(1.1), Inches(0.3), num, sz=18, color=HCL_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(s1, kx + Inches(1.1), Inches(1.0), Inches(1.2), Inches(0.55), label, sz=9, color=GRAY_600, align=PP_ALIGN.LEFT)
    kx += Inches(2.45)
    if kx < Inches(11):
        sep = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, kx - Inches(0.15), Inches(1.08), Inches(0.015), Inches(0.4))
        sep.fill.solid(); sep.fill.fore_color.rgb = GRAY_200; sep.line.fill.background()

# ── Category 1: Quality Engineering ──
cy = Inches(1.85)
add_rect(s1, Inches(0.4), cy, Inches(0.32), Inches(0.32), CAT_BG["quality"])
add_tb(s1, Inches(0.82), cy + Inches(0.02), Inches(4), Inches(0.3),
       "QUALITY ENGINEERING & TESTING", sz=10.5, color=GRAY_800, bold=True)
line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), cy + Inches(0.15), Inches(7.9), Inches(0.015))
line.fill.solid(); line.fill.fore_color.rgb = GRAY_200; line.line.fill.background()

cards_q = [
    ("Agentic QA Lifecycle",
     "Autonomous agents handle test creation, execution, and reporting end-to-end — plugs into Jira & Azure DevOps.",
     "25–30% productivity gain", "Piloted", HCL_ORANGE),
    ("Regression Testing Agent",
     "Scriptless AI automation that runs full regression suites, freeing teams to focus on new features.",
     "25–30% effort reduction", "Live", HCL_GREEN),
    ("Visual App Validator",
     "Computer-vision agent that validates apps across devices — eliminates fragile selectors, cuts maintenance.",
     "15–25% less validation effort", "Ready for pilot", HCL_ACCENT),
    ("Responsible AI Evaluator",
     "Automated pipeline that tests any AI agent for bias, safety, and compliance with full auditability.",
     "Built-in AI safety assurance", "Pilot", HCL_ORANGE),
]
cx = Inches(0.4); cw = Inches(3.05); ch = Inches(1.65)
for title, desc, metric, status, scolor in cards_q:
    draw_card(s1, cx, cy + Inches(0.45), cw, ch, HCL_BLUE, title, desc, metric, status, scolor)
    cx += cw + Inches(0.12)

# ── Category 2: Intelligent Support ──
cy2 = Inches(4.15)
add_rect(s1, Inches(0.4), cy2, Inches(0.32), Inches(0.32), CAT_BG["support"])
add_tb(s1, Inches(0.82), cy2 + Inches(0.02), Inches(4), Inches(0.3),
       "INTELLIGENT SUPPORT & OPERATIONS", sz=10.5, color=GRAY_800, bold=True)
line2 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), cy2 + Inches(0.15), Inches(7.7), Inches(0.015))
line2.fill.solid(); line2.fill.fore_color.rgb = GRAY_200; line2.line.fill.background()

cards_s = [
    ("Autonomous Bug Fixer",
     "11-step AI agent: triage to pull request — with safety controls, audit trails, and human approval gates.",
     "18–25% effort saved · ~90% accuracy", "Live (5 programs)", HCL_GREEN),
    ("Intelligent Support Agent",
     "Surfaces the right knowledge instantly — classifies tickets, triages incidents, guides engineers to resolution.",
     "Faster MTTR · fewer re-opens", "Live (3+ envs)", HCL_GREEN),
    ("Self-Healing Pipelines",
     "Monitors pipelines 24/7, detects anomalies, triggers remediation — before users notice an issue.",
     "30% fewer incidents in 3–4 mo", "Pilot", HCL_ORANGE),
    ("Proactive Risk Monitor",
     "AI agents in Teams track project health, flag delivery risks early, surface actionable insights.",
     "Risks flagged before impact", "Pilot", HCL_ORANGE),
]
cx = Inches(0.4)
for title, desc, metric, status, scolor in cards_s:
    draw_card(s1, cx, cy2 + Inches(0.45), cw, ch, HCL_TEAL, title, desc, metric, status, scolor)
    cx += cw + Inches(0.12)

draw_footer(s1, "1 / 3")


# ═══════════════════════════════════════════════
# SLIDE 2 — Dev Productivity + Governance
# ═══════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, WHITE)
draw_header(s2)

# ── Category 3: Developer Productivity ──
cy3 = Inches(1.2)
add_rect(s2, Inches(0.4), cy3, Inches(0.32), Inches(0.32), CAT_BG["dev"])
add_tb(s2, Inches(0.82), cy3 + Inches(0.02), Inches(5), Inches(0.3),
       "DEVELOPER PRODUCTIVITY & MODERNIZATION", sz=10.5, color=GRAY_800, bold=True)
line3 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), cy3 + Inches(0.15), Inches(7.1), Inches(0.015))
line3.fill.solid(); line3.fill.fore_color.rgb = GRAY_200; line3.line.fill.background()

cards_d = [
    ("Migration Accelerator",
     "AI-generated field mappings and scripts that move data and code to the cloud in a fraction of the usual time.",
     "Up to 60% effort reduction", "Live", HCL_GREEN),
    ("Figma-to-Code",
     "Turns Figma designs into production-ready UI components — designers hand off, developers get buildable code.",
     "Dramatically faster UI builds", "Piloted", HCL_ORANGE),
    ("AI-Assisted Development",
     "Structured Copilot adoption across coding, analytics, and SQL — quantifiable speed gains from day one.",
     "16% faster dev · 12% less QA", "Live", HCL_GREEN),
    ("Creative Asset Engine",
     "AI-driven image processing replacing a 3-day pipeline with sub-3-hour automation — 70% less designer effort.",
     "70% effort eliminated", "Live", HCL_GREEN),
]
cx = Inches(0.4); ch2 = Inches(1.75)
for title, desc, metric, status, scolor in cards_d:
    draw_card(s2, cx, cy3 + Inches(0.45), cw, ch2, HCL_PURPLE, title, desc, metric, status, scolor)
    cx += cw + Inches(0.12)

# ── Category 4: Governance ──
cy4 = Inches(3.65)
add_rect(s2, Inches(0.4), cy4, Inches(0.32), Inches(0.32), CAT_BG["govern"])
add_tb(s2, Inches(0.82), cy4 + Inches(0.02), Inches(5.5), Inches(0.3),
       "GOVERNANCE, COMPLIANCE & BUSINESS EFFICIENCY", sz=10.5, color=GRAY_800, bold=True)
line4 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), cy4 + Inches(0.15), Inches(6.5), Inches(0.015))
line4.fill.solid(); line4.fill.fore_color.rgb = GRAY_200; line4.line.fill.background()

cards_g = [
    ("RACE — Accessibility",
     "Automates WCAG compliance checks at scale — catch accessibility issues before release, not after complaints.",
     "25–30% effort savings", "Live", HCL_GREEN),
    ("Compliance Automation",
     "AI monitors regulatory changes, extracts contract metadata, scans vulnerabilities — audits on autopilot.",
     "30% compliance effort cut", "Live", HCL_GREEN),
    ("Smart Proposal Writer",
     "Generates tailored RFP responses and case studies from your knowledge base — better proposals, faster.",
     "Faster bids · better win rate", "Pilot", HCL_ORANGE),
    ("Live Multilingual Transcription",
     "Real-time meeting transcription in any language — instant subtitles for distributed global teams.",
     "Instant cross-language access", "Pilot", HCL_ORANGE),
]
cx = Inches(0.4)
for title, desc, metric, status, scolor in cards_g:
    draw_card(s2, cx, cy4 + Inches(0.45), cw, ch2, HCL_GREEN, title, desc, metric, status, scolor)
    cx += cw + Inches(0.12)

# ── Why HCLTech? (bottom of slide 2) ──
wy = Inches(5.7)
add_tb(s2, Inches(0.5), wy, Inches(12.5), Inches(0.35),
       "WHY HCLTech?", sz=11, color=HCL_DARK, bold=True)

why_items = [
    ("Enterprise Safety First", "SAFE_MODE, sandboxing, and human-approval gates built into every solution."),
    ("Platform Agnostic", "Works with Azure, AWS, GCP, and any Git-based toolchain — no vendor lock-in."),
    ("Proven at Scale", "Deployed across 15+ programs with measurable, auditable results."),
    ("Zero-Disruption Adoption", "Pilot-first rollout — integrates with existing tools, no process overhaul needed."),
]
wx = Inches(0.4)
ww = Inches(3.05)
for wtitle, wdesc in why_items:
    add_rect(s2, wx, wy + Inches(0.35), ww, Inches(0.85), GRAY_50, GRAY_200, Pt(0.5))
    add_ml(s2, wx + Inches(0.14), wy + Inches(0.4), ww - Inches(0.28), Inches(0.75), [
        (wtitle, 10, HCL_BLUE, True),
        (wdesc, 8.5, GRAY_600, False),
    ])
    wx += ww + Inches(0.12)

draw_footer(s2, "2 / 3")


# ═══════════════════════════════════════════════
# SLIDE 3 — Summary & CTA
# ═══════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, WHITE)
draw_header(s3)

# Large summary grid
add_tb(s3, Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.4),
       "SOLUTION LANDSCAPE AT A GLANCE", sz=13, color=HCL_DARK, bold=True)
add_tb(s3, Inches(0.5), Inches(1.55), Inches(12.5), Inches(0.3),
       "16 proven GenAI solutions across 4 categories — most already live in enterprise programs", sz=10, color=GRAY_500)

# Summary table
cols = [
    ("Quality Engineering\n& Testing", HCL_BLUE,
     ["Agentic QA Lifecycle", "Regression Testing Agent", "Visual App Validator", "Responsible AI Evaluator"]),
    ("Intelligent Support\n& Operations", HCL_TEAL,
     ["Autonomous Bug Fixer", "Intelligent Support Agent", "Self-Healing Pipelines", "Proactive Risk Monitor"]),
    ("Developer Productivity\n& Modernization", HCL_PURPLE,
     ["Migration Accelerator", "Figma-to-Code", "AI-Assisted Development", "Creative Asset Engine"]),
    ("Governance, Compliance\n& Business Efficiency", HCL_GREEN,
     ["RACE — Accessibility", "Compliance Automation", "Smart Proposal Writer", "Live Multilingual Transcription"]),
]

tx = Inches(0.5)
tw = Inches(3.0)
for cat_name, accent, items in cols:
    # Category header card
    add_rect(s3, tx, Inches(2.05), tw, Inches(0.65), accent)
    add_tb(s3, tx + Inches(0.12), Inches(2.08), tw - Inches(0.24), Inches(0.6),
           cat_name, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Solution list
    iy = Inches(2.8)
    for item in items:
        add_rect(s3, tx, iy, tw, Inches(0.42), WHITE, GRAY_200, Pt(0.5))
        # accent dot
        add_dot(s3, tx + Inches(0.12), iy + Inches(0.16), accent)
        add_tb(s3, tx + Inches(0.28), iy + Inches(0.06), tw - Inches(0.4), Inches(0.32),
               item, sz=10.5, color=GRAY_800)
        iy += Inches(0.48)
    tx += tw + Inches(0.12)

# Key outcomes row
oy = Inches(4.85)
add_tb(s3, Inches(0.5), oy, Inches(12.5), Inches(0.35),
       "KEY OUTCOMES DELIVERED", sz=11, color=HCL_DARK, bold=True)

outcomes = [
    ("25–30%", "QA effort eliminated\nacross testing lifecycle"),
    ("60%", "Faster cloud migrations\nwith AI-generated scripts"),
    ("~90%", "Bug-fix accuracy from\nautonomous resolution"),
    ("70%", "Creative production effort\nreplaced by automation"),
    ("30%", "Compliance workload\nreduced with AI monitoring"),
]
ox = Inches(0.5); ow = Inches(2.38)
for num, label in outcomes:
    add_rect(s3, ox, oy + Inches(0.4), ow, Inches(0.85), GRAY_50, GRAY_200, Pt(0.5))
    add_ml(s3, ox + Inches(0.1), oy + Inches(0.42), ow - Inches(0.2), Inches(0.8), [
        (num, 20, HCL_BLUE, True),
        (label, 9, GRAY_600, False),
    ], align=PP_ALIGN.CENTER)
    ox += ow + Inches(0.12)

# CTA
add_rect(s3, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.65), HCL_DARK)
add_tb(s3, Inches(0.7), Inches(6.28), Inches(11.93), Inches(0.5),
       "Ready to accelerate?  Let's identify which solutions deliver the fastest ROI for your delivery model.",
       sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

draw_footer(s3, "3 / 3")


# ── Save ──
out = os.path.join(os.path.dirname(__file__), "TechVDU_GenAI_Service_Transformation_Solutions.pptx")
prs.save(out)
print(f"Saved: {out}")
