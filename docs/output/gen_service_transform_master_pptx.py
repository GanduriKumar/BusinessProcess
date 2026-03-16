"""
Generate GenAI Service Transformation PPTX using HCLTech master slide template.
Uses 'Title w/t Subtitle' layout for category slides with manually-positioned card
boxes — ensures large, legible fonts and precise layout control.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──
TEMPLATE = os.path.join(os.path.dirname(__file__), "..", "input",
                        "TechVDU_Transformation_MasterSlide.pptx")
OUTPUT   = os.path.join(os.path.dirname(__file__),
                        "TechVDU_GenAI_Service_Transformation_Solutions.pptx")

# ── HCL Brand Colours ──
HCL_DARK   = RGBColor(0x00, 0x2B, 0x5C)
HCL_BLUE   = RGBColor(0x00, 0x60, 0xA9)
HCL_ACCENT = RGBColor(0x00, 0xA3, 0xE0)
HCL_GREEN  = RGBColor(0x00, 0x87, 0x5A)
HCL_ORANGE = RGBColor(0xE8, 0x77, 0x22)
HCL_PURPLE = RGBColor(0x6C, 0x3F, 0xA5)
HCL_TEAL   = RGBColor(0x00, 0x85, 0x7C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50    = RGBColor(0xF8, 0xFA, 0xFC)
GRAY_200   = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_500   = RGBColor(0x64, 0x74, 0x8B)
GRAY_600   = RGBColor(0x47, 0x55, 0x69)
GRAY_800   = RGBColor(0x1E, 0x29, 0x3B)

# ── Layout indices in the master template ──
LY_COVER     = 0    # Cover – The Beam (Intense)
LY_TITLE_SUB = 42   # Title w/t Subtitle — gives branded header + full canvas
LY_CTA       = 72   # Premium layout – Main message (Dark)


# ═══════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════
def delete_slide(prs, index):
    """Remove existing slide at index."""
    rId = prs.slides._sldIdLst[index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def set_ph_text(slide, ph_idx, text, font_size=None, bold=None, color=None,
                align=None):
    """Set simple single-run text in a placeholder."""
    ph = slide.placeholders[ph_idx]
    ph.text = text
    for p in ph.text_frame.paragraphs:
        if align:
            p.alignment = align
        for r in p.runs:
            if font_size:
                r.font.size = Pt(font_size)
            if bold is not None:
                r.font.bold = bold
            if color:
                r.font.color.rgb = color


def tb(slide, l, t, w, h, text, sz=12, color=GRAY_800, bold=False,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box with a single run."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def rich_tb(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """Add a text box with multiple styled paragraphs.
    lines = [(text, size_pt, color, bold, align), ...]
    """
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (txt, sz, col, bld, al) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        if al:
            p.alignment = al
        run = p.add_run()
        run.text = txt
        if sz:
            run.font.size = Pt(sz)
        if col:
            run.font.color.rgb = col
        if bld is not None:
            run.font.bold = bld
    return box


def stripe(slide, x, y, w, color, h=Inches(0.06)):
    """Thin accent stripe."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def dot(slide, x, y, color, size=Inches(0.12)):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background()


def rounded_box(slide, l, t, w, h, fill=WHITE, border=GRAY_200):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(0.75)
    s.adjustments[0] = 0.04
    return s


def set_footer(slide, number):
    for ph in slide.placeholders:
        if ph.name and "Slide Number" in ph.name:
            ph.text = str(number)
        elif ph.name and "Footer" in ph.name:
            ph.text = "Copyright \u00A9 2025 HCLTech | Confidential"


# ═══════════════════════════════════════════════
# Content data
# ═══════════════════════════════════════════════
KPIS = [
    ("15+",         "Deployable Solutions"),
    ("25\u201330%", "QA Effort Eliminated"),
    ("60%",         "Faster Migrations"),
    ("\u223C90%",   "Bug-Fix Accuracy"),
    ("30%",         "Compliance Savings"),
]

CATEGORIES = [
    {
        "title": "Quality Engineering & Testing",
        "subtitle": "AI-driven testing that finds defects faster and eliminates manual QA overhead",
        "accent": HCL_BLUE,
        "cards": [
            ("Agentic QA Lifecycle",
             "Autonomous agents handle test creation, execution, and reporting end-to-end \u2014 plugs into Jira & Azure DevOps.",
             "\u25B2 25\u201330% productivity gain", "Piloted", HCL_ORANGE),
            ("Regression Testing Agent",
             "Scriptless AI automation that runs full regression suites, freeing teams to focus on new features.",
             "\u25B2 25\u201330% effort reduction", "Live", HCL_GREEN),
            ("Visual App Validator",
             "Computer-vision agent that validates apps across devices \u2014 eliminates fragile selectors, cuts maintenance.",
             "\u25B2 15\u201325% less validation effort", "Ready for pilot", HCL_ACCENT),
            ("Responsible AI Evaluator",
             "Automated pipeline that tests any AI agent for bias, safety, and compliance with full auditability.",
             "\u25B2 Built-in AI safety assurance", "Pilot", HCL_ORANGE),
        ],
    },
    {
        "title": "Intelligent Support & Operations",
        "subtitle": "AI agents that resolve issues autonomously and keep production environments healthy",
        "accent": HCL_TEAL,
        "cards": [
            ("Autonomous Bug Fixer",
             "11-step AI agent: triage to pull request \u2014 with safety controls, audit trails, and human approval gates.",
             "\u25B2 18\u201325% effort saved \u00B7 ~90% accuracy", "Live in 5 programs", HCL_GREEN),
            ("Intelligent Support Agent",
             "Surfaces the right knowledge instantly \u2014 classifies tickets, triages incidents, guides engineers to resolution.",
             "\u25B2 Faster MTTR \u00B7 fewer re-opens", "Live in 3+ envs", HCL_GREEN),
            ("Self-Healing Pipelines",
             "Monitors pipelines 24/7, detects anomalies, triggers remediation \u2014 before users notice an issue.",
             "\u25B2 30% fewer incidents in 3\u20134 mo", "Pilot", HCL_ORANGE),
            ("Proactive Risk Monitor",
             "AI agents in Teams track project health, flag delivery risks early, surface actionable insights.",
             "\u25B2 Risks flagged before impact", "Pilot", HCL_ORANGE),
        ],
    },
    {
        "title": "Developer Productivity & Modernization",
        "subtitle": "Accelerate migrations, automate UI builds, and amplify developer output with AI tooling",
        "accent": HCL_PURPLE,
        "cards": [
            ("Migration Accelerator",
             "AI-generated field mappings and scripts that move data and code to the cloud in a fraction of the usual time.",
             "\u25B2 Up to 60% effort reduction", "Live", HCL_GREEN),
            ("Figma-to-Code",
             "Turns Figma designs into production-ready UI components \u2014 designers hand off, developers get buildable code.",
             "\u25B2 Dramatically faster UI builds", "Piloted", HCL_ORANGE),
            ("AI-Assisted Development",
             "Structured Copilot adoption across coding, analytics, and SQL \u2014 quantifiable speed gains from day one.",
             "\u25B2 16% faster dev \u00B7 12% less QA", "Live", HCL_GREEN),
            ("Creative Asset Engine",
             "AI-driven image processing replacing a 3-day pipeline with sub-3-hour automation \u2014 70% less designer effort.",
             "\u25B2 70% effort eliminated", "Live", HCL_GREEN),
        ],
    },
    {
        "title": "Governance, Compliance & Business Efficiency",
        "subtitle": "Automate audits, ensure accessibility, and generate proposals faster with AI-powered workflows",
        "accent": HCL_GREEN,
        "cards": [
            ("RACE \u2014 Accessibility",
             "Automates WCAG compliance checks at scale \u2014 catch accessibility issues before release, not after complaints.",
             "\u25B2 25\u201330% effort savings", "Live", HCL_GREEN),
            ("Compliance Automation",
             "AI monitors regulatory changes, extracts contract metadata, scans vulnerabilities \u2014 audits on autopilot.",
             "\u25B2 30% compliance effort cut", "Live", HCL_GREEN),
            ("Smart Proposal Writer",
             "Generates tailored RFP responses and case studies from your knowledge base \u2014 better proposals, faster.",
             "\u25B2 Faster bids \u00B7 better win rate", "Pilot", HCL_ORANGE),
            ("Live Multilingual Transcription",
             "Real-time meeting transcription in any language \u2014 instant subtitles for distributed global teams.",
             "\u25B2 Instant cross-language access", "Pilot", HCL_ORANGE),
        ],
    },
]

WHY_ITEMS = [
    ("Enterprise Safety First",
     "SAFE_MODE, sandboxing, and human-approval gates built into every solution."),
    ("Platform Agnostic",
     "Works with Azure, AWS, GCP, and any Git-based toolchain \u2014 no vendor lock-in."),
    ("Proven at Scale",
     "Deployed across 15+ programs with measurable, auditable results."),
    ("Zero-Disruption Adoption",
     "Pilot-first rollout \u2014 integrates with existing tools, no process overhaul needed."),
]


# ═══════════════════════════════════════════════
# Build Presentation
# ═══════════════════════════════════════════════
prs = Presentation(TEMPLATE)
delete_slide(prs, 0)

slide_num = 0

# ─────────────────────────────────────────────
# SLIDE 1 — Cover (uses branded cover layout)
# ─────────────────────────────────────────────
slide_num += 1
s1 = prs.slides.add_slide(prs.slide_layouts[LY_COVER])
set_ph_text(s1, 0, "Accelerate Delivery with GenAI", font_size=32, bold=True)
set_ph_text(s1, 1,
            "Battle-tested AI solutions that cut costs, raise quality, "
            "and compress timelines",
            font_size=16)

# ─────────────────────────────────────────────
# SLIDE 2 — Key Metrics (Title w/t Subtitle + manual KPI boxes)
# ─────────────────────────────────────────────
slide_num += 1
s2 = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_SUB])
set_ph_text(s2, 0, "Impact at a Glance", font_size=28, bold=True, color=HCL_DARK)
set_ph_text(s2, 12, "Proven, measurable results across 15+ enterprise programs",
            font_size=14, color=GRAY_600)

# 5 KPI boxes in a row, vertically centred on the slide
kpi_w = Inches(2.2)
kpi_h = Inches(2.6)
kpi_gap = Inches(0.2)
total_kpi_w = 5 * kpi_w + 4 * kpi_gap
kpi_x0 = (Inches(13.333) - total_kpi_w) // 2
kpi_y = Inches(2.2)

for i, (num, label) in enumerate(KPIS):
    x = kpi_x0 + i * (kpi_w + kpi_gap)
    rounded_box(s2, x, kpi_y, kpi_w, kpi_h, GRAY_50, GRAY_200)
    # Big number
    tb(s2, x, kpi_y + Inches(0.35), kpi_w, Inches(0.8),
       num, sz=40, color=HCL_BLUE, bold=True, align=PP_ALIGN.CENTER)
    # Label below
    tb(s2, x, kpi_y + Inches(1.3), kpi_w, Inches(0.8),
       label, sz=16, color=GRAY_600, align=PP_ALIGN.CENTER)
    # Accent stripe at bottom of box
    stripe(s2, x, kpi_y + kpi_h - Inches(0.06), kpi_w, HCL_ACCENT, Inches(0.06))

set_footer(s2, slide_num)

# ─────────────────────────────────────────────
# SLIDES 3–6 — One category per slide (manual card layout)
# ─────────────────────────────────────────────
CARD_W    = Inches(2.95)
CARD_H    = Inches(4.5)
CARD_GAP  = Inches(0.13)
CARDS_X0  = (Inches(13.333) - 4 * CARD_W - 3 * CARD_GAP) // 2
CARD_Y    = Inches(1.75)

for cat in CATEGORIES:
    slide_num += 1
    s = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_SUB])

    # Title & subtitle via placeholders (inherits template font/position)
    set_ph_text(s, 0, cat["title"], font_size=26, bold=True, color=HCL_DARK)
    set_ph_text(s, 12, cat["subtitle"], font_size=14, color=GRAY_600)

    for ci, (title, desc, metric, status, status_color) in enumerate(cat["cards"]):
        x = CARDS_X0 + ci * (CARD_W + CARD_GAP)

        # Card background box
        rounded_box(s, x, CARD_Y, CARD_W, CARD_H, WHITE, GRAY_200)

        # Top accent stripe
        stripe(s, x, CARD_Y, CARD_W, cat["accent"], Inches(0.06))

        # Card title (large, bold)
        tb(s, x + Inches(0.15), CARD_Y + Inches(0.18), CARD_W - Inches(0.3), Inches(0.55),
           title, sz=15, color=GRAY_800, bold=True)

        # Description (readable size)
        tb(s, x + Inches(0.15), CARD_Y + Inches(0.75), CARD_W - Inches(0.3), Inches(1.9),
           desc, sz=12, color=GRAY_600)

        # Metric badge (green, bold)
        metric_y = CARD_Y + CARD_H - Inches(1.15)
        rounded_box(s, x + Inches(0.1), metric_y, CARD_W - Inches(0.2), Inches(0.5),
                    GRAY_50, GRAY_200)
        tb(s, x + Inches(0.18), metric_y + Inches(0.06),
           CARD_W - Inches(0.36), Inches(0.4),
           metric, sz=11, color=HCL_GREEN, bold=True, align=PP_ALIGN.CENTER)

        # Status dot + label at bottom
        status_y = CARD_Y + CARD_H - Inches(0.45)
        dot(s, x + Inches(0.15), status_y + Inches(0.05), status_color, Inches(0.12))
        tb(s, x + Inches(0.32), status_y, CARD_W - Inches(0.45), Inches(0.3),
           status, sz=11, color=GRAY_500)

    set_footer(s, slide_num)

# ─────────────────────────────────────────────
# SLIDE 7 — Why HCLTech? (manual 4-box layout)
# ─────────────────────────────────────────────
slide_num += 1
s7 = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_SUB])

set_ph_text(s7, 0, "Why HCLTech?", font_size=28, bold=True, color=HCL_DARK)
set_ph_text(s7, 12,
            "Enterprise-grade safety controls  |  Full audit trails  |  "
            "Human-approval gates  |  Zero-disruption deployment",
            font_size=13, color=GRAY_600)

why_w = Inches(2.95)
why_h = Inches(3.6)
why_gap = Inches(0.13)
why_x0 = (Inches(13.333) - 4 * why_w - 3 * why_gap) // 2
why_y = Inches(1.85)
why_icons = ["\u26A1", "\u2601", "\u2714", "\u21BB"]

for i, (wtitle, wdesc) in enumerate(WHY_ITEMS):
    x = why_x0 + i * (why_w + why_gap)
    rounded_box(s7, x, why_y, why_w, why_h, WHITE, GRAY_200)
    stripe(s7, x, why_y, why_w, HCL_BLUE, Inches(0.06))

    # Icon
    tb(s7, x, why_y + Inches(0.2), why_w, Inches(0.6),
       why_icons[i], sz=32, color=HCL_BLUE, align=PP_ALIGN.CENTER)

    # Title
    tb(s7, x + Inches(0.2), why_y + Inches(0.9), why_w - Inches(0.4), Inches(0.5),
       wtitle, sz=16, color=HCL_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Description
    tb(s7, x + Inches(0.2), why_y + Inches(1.5), why_w - Inches(0.4), Inches(1.6),
       wdesc, sz=13, color=GRAY_600, align=PP_ALIGN.CENTER)

# Legend bar at bottom
legend_y = Inches(5.7)
tb(s7, Inches(0.7), legend_y, Inches(12), Inches(0.35),
   "Enterprise-grade safety controls  \u00B7  Full audit trails  \u00B7  "
   "Human-approval gates  \u00B7  Zero-disruption deployment",
   sz=12, color=GRAY_500, align=PP_ALIGN.CENTER)

# Legend dots
lx = Inches(3.8)
for lbl, col in [("Live", HCL_GREEN), ("Pilot", HCL_ORANGE), ("POC", HCL_ACCENT)]:
    dot(s7, lx, legend_y + Inches(0.42), col, Inches(0.1))
    tb(s7, lx + Inches(0.15), legend_y + Inches(0.37), Inches(0.6), Inches(0.25),
       lbl, sz=11, color=GRAY_500)
    lx += Inches(1.5)

set_footer(s7, slide_num)

# ─────────────────────────────────────────────
# SLIDE 8 — CTA (Premium – Main message, Dark)
# ─────────────────────────────────────────────
slide_num += 1
s8 = prs.slides.add_slide(prs.slide_layouts[LY_CTA])

set_ph_text(s8, 14, "GenAI for Service Transformation", font_size=16,
            color=RGBColor(0xBB, 0xDB, 0xF0))
set_ph_text(s8, 0,
            "Ready to accelerate?\n"
            "Let\u2019s identify which solutions deliver the fastest ROI "
            "for your delivery model.",
            font_size=28, bold=True)

set_footer(s8, slide_num)

# ── Save ──
prs.save(OUTPUT)
print("Saved:", OUTPUT)
