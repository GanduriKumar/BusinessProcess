from collections import defaultdict
from pathlib import Path
import re
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


WORKBOOK = Path(r"C:\users\kumar.gn\HClProjects\BusinessProcess\CUNB-2 Mar'26 V.1.xlsx")
OUTPUT = Path(r"C:\users\kumar.gn\HClProjects\BusinessProcess\CUNB Exec Review - L3 Account Movement Mar 2026.pptx")

NS = {
    "m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

BG = RGBColor(247, 247, 245)
TEXT = RGBColor(32, 37, 41)
MUTED = RGBColor(92, 99, 106)
RED = RGBColor(192, 57, 43)
GREEN = RGBColor(26, 127, 55)
GOLD = RGBColor(184, 134, 11)
BLUE = RGBColor(21, 101, 192)
LIGHT_RED = RGBColor(244, 199, 195)
LIGHT_GREEN = RGBColor(198, 239, 206)
LIGHT_BLUE = RGBColor(207, 226, 243)
CARD = RGBColor(255, 255, 255)
GRID = RGBColor(214, 219, 223)


def load_workbook_rows(sheet_name):
    with zipfile.ZipFile(WORKBOOK) as z:
        wb = ET.fromstring(z.read("xl/workbook.xml"))
        rels = ET.fromstring(z.read("xl/_rels/workbook.xml.rels"))
        relmap = {r.attrib["Id"]: "xl/" + r.attrib["Target"] for r in rels}

        shared = []
        sst = ET.fromstring(z.read("xl/sharedStrings.xml"))
        for si in sst.findall("m:si", NS):
            t = si.find("m:t", NS)
            if t is not None:
                shared.append(t.text or "")
            else:
                parts = []
                for run in si.findall("m:r", NS):
                    rt = run.find("m:t", NS)
                    if rt is not None:
                        parts.append(rt.text or "")
                shared.append("".join(parts))

        target = None
        for s in wb.find("m:sheets", NS):
            if s.attrib["name"] == sheet_name:
                rel_id = s.attrib["{%s}id" % NS["r"]]
                target = relmap[rel_id]
                break

        if not target:
            raise ValueError(f"Sheet not found: {sheet_name}")

        sh = ET.fromstring(z.read(target))
        rows = []
        headers = {}
        for row in sh.findall("m:sheetData/m:row", NS):
            rnum = int(row.attrib["r"])
            values = {}
            for c in row.findall("m:c", NS):
                ref = c.attrib["r"]
                col = re.match(r"([A-Z]+)", ref).group(1)
                v = c.find("m:v", NS)
                if v is None:
                    continue
                raw = v.text or ""
                if c.attrib.get("t") == "s":
                    raw = shared[int(raw)]
                values[col] = raw
            if rnum == 3:
                headers = values
            elif rnum > 3 and values:
                rows.append(values)
        return headers, rows


def fte(row):
    return float(row.get("BK", "0") or 0)


def short_l3(name):
    return name.replace("ERS VDU-TECH-", "")


ACCOUNT_SHORT = {
    "MICROSOFT CORPORATION": "Microsoft",
    "Google LLC": "Google",
    "HCL Technologies  Ltd.": "HCL Tech",
    "SAMSUNG DATA SYSTEMS INDIA PVT LTD": "Samsung SDS",
    "Pearson Technology": "Pearson",
    "Sirius XM Radio Inc.": "SiriusXM",
    "Meta Platforms, Inc.": "Meta",
    "PayPal Inc.,": "PayPal",
    "Western Union Inc": "Western Union",
    "Dassault Systemes SE": "Dassault",
    "Spatial Corporation": "Spatial",
    "Apple Inc": "Apple",
    "J.Crew Group": "J.Crew",
}


def short_account(name):
    return ACCOUNT_SHORT.get(name, name[:18] + ("..." if len(name) > 18 else ""))


def aggregate():
    _, b2u_rows = load_workbook_rows("Billed to Unbilled")
    _, u2b_rows = load_workbook_rows("Got Billed")

    by_l3 = defaultdict(lambda: {"b2u": 0.0, "u2b": 0.0})
    by_account = defaultdict(lambda: {"b2u": 0.0, "u2b": 0.0})
    by_pair = defaultdict(lambda: {"b2u": 0.0, "u2b": 0.0})
    b2u_reason = defaultdict(float)
    u2b_reason = defaultdict(float)

    for row in b2u_rows:
        val = fte(row)
        l3 = row.get("AQ", "Unknown")
        acct = row.get("Y", "Unknown")
        by_l3[l3]["b2u"] += val
        by_account[acct]["b2u"] += val
        by_pair[(l3, acct)]["b2u"] += val
        b2u_reason[row.get("DI", row.get("S", "Unknown"))] += val

    for row in u2b_rows:
        val = fte(row)
        l3 = row.get("AQ", "Unknown")
        acct = row.get("Y", "Unknown")
        by_l3[l3]["u2b"] += val
        by_account[acct]["u2b"] += val
        by_pair[(l3, acct)]["u2b"] += val
        u2b_reason[row.get("S", "Unknown")] += val

    l3_summary = []
    for l3, vals in by_l3.items():
        l3_summary.append(
            {
                "label": short_l3(l3),
                "full": l3,
                "b2u": vals["b2u"],
                "u2b": vals["u2b"],
                "net": vals["u2b"] - vals["b2u"],
            }
        )

    account_summary = []
    for acct, vals in by_account.items():
        account_summary.append(
            {
                "label": short_account(acct),
                "full": acct,
                "b2u": vals["b2u"],
                "u2b": vals["u2b"],
                "net": vals["u2b"] - vals["b2u"],
            }
        )

    pair_summary = []
    for (l3, acct), vals in by_pair.items():
        pair_summary.append(
            {
                "l3": short_l3(l3),
                "account": short_account(acct),
                "full_l3": l3,
                "full_account": acct,
                "b2u": vals["b2u"],
                "u2b": vals["u2b"],
                "net": vals["u2b"] - vals["b2u"],
            }
        )

    return {
        "total_b2u": sum(fte(r) for r in b2u_rows),
        "total_u2b": sum(fte(r) for r in u2b_rows),
        "total_net": sum(fte(r) for r in u2b_rows) - sum(fte(r) for r in b2u_rows),
        "b2u_rows": len(b2u_rows),
        "u2b_rows": len(u2b_rows),
        "l3": sorted(l3_summary, key=lambda x: (-abs(x["net"]), x["label"])),
        "accounts": sorted(account_summary, key=lambda x: (-abs(x["net"]), x["label"])),
        "pairs_neg": sorted([p for p in pair_summary if p["net"] < 0], key=lambda x: (x["net"], x["l3"], x["account"])),
        "pairs_pos": sorted([p for p in pair_summary if p["net"] > 0], key=lambda x: (-x["net"], x["l3"], x["account"])),
        "b2u_reason": sorted(b2u_reason.items(), key=lambda x: (-x[1], x[0])),
        "u2b_reason": sorted(u2b_reason.items(), key=lambda x: (-x[1], x[0])),
    }


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    return box


def add_bullet_list(slide, left, top, width, height, bullets, size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.bullet = True
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, title, value, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = GRID
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.09), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(slide, left + Inches(0.18), top + Inches(0.18), width - Inches(0.3), Inches(0.28), title, size=13, color=MUTED)
    add_text(slide, left + Inches(0.18), top + Inches(0.52), width - Inches(0.3), Inches(0.5), value, size=26, bold=True, color=TEXT)


def style_chart(chart, legend=True):
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = GRID
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(10)
    try:
        chart.plot_area.format.fill.solid()
        chart.plot_area.format.fill.fore_color.rgb = CARD
    except Exception:
        pass


def add_clustered_bar_chart(slide, left, top, width, height, categories, series_map, title=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, values in series_map:
        chart_data.add_series(series_name, values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data)
    chart = chart_shape.chart
    style_chart(chart, legend=True)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = max(max(v) for _, v in series_map) * 1.2 if categories else 1
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    for idx, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.line.fill.background()
        series.format.fill.fore_color.rgb = RED if idx == 0 else GREEN
        series.has_data_labels = True
        series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        series.data_labels.font.size = Pt(9)
    return chart


def add_column_chart(slide, left, top, width, height, categories, values, color, title=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("FTE", values)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
    chart = chart_shape.chart
    style_chart(chart, legend=False)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = max(values) * 1.25 if values else 1
    chart.category_axis.tick_labels.font.size = Pt(11)
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.fill.background()
    series.has_data_labels = True
    series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    series.data_labels.font.size = Pt(10)
    return chart


def add_table_like_box(slide, left, top, width, title, headers, rows, accent):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(2.7))
    outer.fill.solid()
    outer.fill.fore_color.rgb = CARD
    outer.line.color.rgb = GRID
    outer.line.width = Pt(1)

    add_text(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.25), title, size=14, bold=True, color=accent)
    add_text(slide, left + Inches(0.15), top + Inches(0.38), width - Inches(0.3), Inches(0.18), " | ".join(headers), size=10, bold=True, color=MUTED)

    y = top + Inches(0.62)
    for row in rows:
        line = " | ".join(str(x) for x in row)
        add_text(slide, left + Inches(0.15), y, width - Inches(0.3), Inches(0.24), line, size=11, color=TEXT)
        y += Inches(0.28)


def build_deck(data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, Inches(0.45), Inches(0.25), Inches(8), Inches(0.5), "CUNB Movement Executive Review", size=24, bold=True)
    add_text(slide, Inches(0.45), Inches(0.68), Inches(8), Inches(0.3), "Movement between 23 Feb 2026 and 2 Mar 2026 | Source: 'Billed to Unbilled' and 'Got Billed' sheets", size=11, color=MUTED)

    add_card(slide, Inches(0.45), Inches(1.1), Inches(2.0), Inches(1.1), "Billed to Unbilled", f"{data['total_b2u']:.0f} FTE", RED)
    add_card(slide, Inches(2.65), Inches(1.1), Inches(2.0), Inches(1.1), "Unbilled to Billed", f"{data['total_u2b']:.0f} FTE", GREEN)
    add_card(slide, Inches(4.85), Inches(1.1), Inches(2.0), Inches(1.1), "Net Movement", f"{data['total_net']:.0f} FTE", GOLD)

    add_column_chart(
        slide,
        Inches(0.45),
        Inches(2.45),
        Inches(4.2),
        Inches(3.2),
        ["Billed to Unbilled", "Unbilled to Billed"],
        [data["total_b2u"], data["total_u2b"]],
        BLUE,
        "Overall movement volume",
    )

    insight_bullets = [
        "Net deterioration is 124 FTE: 167 moved out of billed status while only 43 were recovered into billed status.",
        "Pressure is concentrated in ALPHABET, MS CORE ENG, ECP, META&CSG and SERVICE VERTICALS.",
        "The only L3 with positive net recovery is MS COMPUTE TECHNOLOGY at +5 FTE.",
        "Largest negative accounts are Google (-32), HCL Tech (-27) and Microsoft (-23). Largest positive recovery is PayPal (+8).",
    ]
    add_text(slide, Inches(5.1), Inches(2.45), Inches(3.6), Inches(0.3), "Executive takeaways", size=15, bold=True)
    add_bullet_list(slide, Inches(5.1), Inches(2.8), Inches(7.6), Inches(2.6), insight_bullets, size=15)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, Inches(0.45), Inches(0.28), Inches(7.5), Inches(0.4), "L3 movement by direction", size=22, bold=True)
    add_text(slide, Inches(0.45), Inches(0.64), Inches(8), Inches(0.25), "Billed to Unbilled vs Unbilled to Billed, FTE", size=11, color=MUTED)

    l3_top = data["l3"][:6]
    cats = [x["label"] for x in l3_top]
    b2u_vals = [x["b2u"] for x in l3_top]
    u2b_vals = [x["u2b"] for x in l3_top]
    add_clustered_bar_chart(
        slide,
        Inches(0.45),
        Inches(1.0),
        Inches(7.2),
        Inches(5.6),
        cats,
        [("Billed to Unbilled", b2u_vals), ("Unbilled to Billed", u2b_vals)],
        "Where the movement sits by L3",
    )

    add_text(slide, Inches(8.0), Inches(1.0), Inches(4.7), Inches(0.3), "Net effect by L3", size=15, bold=True)
    y = Inches(1.42)
    for item in sorted(l3_top, key=lambda x: x["net"]):
        color = GREEN if item["net"] > 0 else RED
        txt = f"{item['label']}: {'+' if item['net'] > 0 else ''}{item['net']:.0f} FTE"
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.0), y, Inches(4.4), Inches(0.46))
        chip.fill.solid()
        chip.fill.fore_color.rgb = LIGHT_GREEN if item["net"] > 0 else LIGHT_RED
        chip.line.fill.background()
        add_text(slide, Inches(8.15), y + Inches(0.07), Inches(4.0), Inches(0.2), txt, size=14, bold=True, color=color)
        y += Inches(0.58)

    add_text(slide, Inches(8.0), Inches(5.15), Inches(4.5), Inches(0.6), "Alphabet and MS Core drive half of the deterioration. Service Verticals has the largest gross recovery, but still ends net negative.", size=13, color=MUTED)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, Inches(0.45), Inches(0.28), Inches(8), Inches(0.4), "Account-level movement", size=22, bold=True)
    add_text(slide, Inches(0.45), Inches(0.64), Inches(8), Inches(0.25), "Largest deterioration and largest recovery accounts, net FTE", size=11, color=MUTED)

    neg_accounts = [x for x in data["accounts"] if x["net"] < 0][:6]
    pos_accounts = [x for x in data["accounts"] if x["net"] > 0][:6]
    add_column_chart(
        slide,
        Inches(0.45),
        Inches(1.0),
        Inches(5.9),
        Inches(4.3),
        [x["label"] for x in neg_accounts],
        [abs(x["net"]) for x in neg_accounts],
        RED,
        "Top deterioration accounts",
    )
    add_column_chart(
        slide,
        Inches(6.75),
        Inches(1.0),
        Inches(5.9),
        Inches(4.3),
        [x["label"] for x in pos_accounts],
        [x["net"] for x in pos_accounts],
        GREEN,
        "Top recovery accounts",
    )

    bullets = [
        "Google is the largest single negative account at -32 FTE, entirely within ALPHABET.",
        "HCL Tech and Microsoft follow at -27 and -23 FTE respectively.",
        "PayPal is the strongest positive account at +8 FTE, driven from SERVICE VERTICALS-FS2.",
        "The positive side is narrow; most improvements are +1 or +2 FTE while deterioration is concentrated and larger.",
    ]
    add_bullet_list(slide, Inches(0.65), Inches(5.6), Inches(12.0), Inches(1.3), bullets, size=14)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, Inches(0.45), Inches(0.28), Inches(8), Inches(0.4), "Hotspots and movement drivers", size=22, bold=True)
    add_text(slide, Inches(0.45), Inches(0.64), Inches(8), Inches(0.25), "Largest L3-account swings and dominant reasons", size=11, color=MUTED)

    neg_rows = [[p["l3"], p["account"], f"{p['net']:.0f}"] for p in data["pairs_neg"][:6]]
    pos_rows = [[p["l3"], p["account"], f"+{p['net']:.0f}"] for p in data["pairs_pos"][:6]]
    add_table_like_box(slide, Inches(0.45), Inches(1.0), Inches(5.9), "Top negative hotspots", ["L3", "Account", "Net"], neg_rows, RED)
    add_table_like_box(slide, Inches(6.75), Inches(1.0), Inches(5.9), "Top positive hotspots", ["L3", "Account", "Net"], pos_rows, GREEN)

    b2u_driver = [f"{name} ({val:.0f})" for name, val in data["b2u_reason"][:5]]
    u2b_driver = [f"{name} ({val:.0f})" for name, val in data["u2b_reason"][:5]]
    add_text(slide, Inches(0.55), Inches(4.2), Inches(2.5), Inches(0.25), "Billed -> Unbilled drivers", size=14, bold=True, color=RED)
    add_bullet_list(slide, Inches(0.55), Inches(4.48), Inches(5.8), Inches(2.0), b2u_driver, size=13)
    add_text(slide, Inches(6.85), Inches(4.2), Inches(2.5), Inches(0.25), "Unbilled -> Billed drivers", size=14, bold=True, color=GREEN)
    add_bullet_list(slide, Inches(6.85), Inches(4.48), Inches(5.8), Inches(2.0), u2b_driver, size=13)

    # Footer on all slides
    for s in prs.slides:
        add_text(s, Inches(10.45), Inches(7.0), Inches(2.4), Inches(0.2), "Source: CUNB-2 Mar'26 V.1.xlsx", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_deck(aggregate())
    print(OUTPUT)
