from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(__file__).resolve().parents[3] / "docs" / "output" / "Upcoming Review - Exec Summary.pptx"

BG = RGBColor(247, 245, 240)
INK = RGBColor(21, 38, 62)
MUTED = RGBColor(93, 107, 128)
TEAL = RGBColor(0, 142, 163)
CORAL = RGBColor(233, 99, 76)
GOLD = RGBColor(240, 181, 68)
SLATE = RGBColor(224, 229, 236)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(188, 52, 61)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def panel(slide, left, top, width, height, color, rounded=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def textbox(slide, left, top, width, height, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box


def bullet_list(slide, left, top, width, height, items, color=INK, size=14, bullet_color=TEAL):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.bullet = True
    return box


def header(slide, step, title, subtitle):
    textbox(slide, Inches(0.68), Inches(0.42), Inches(1.2), Inches(0.3), step, 11, TEAL, True)
    textbox(slide, Inches(0.68), Inches(0.78), Inches(8.7), Inches(0.55), title, 26, INK, True)
    textbox(slide, Inches(0.68), Inches(1.4), Inches(10.8), Inches(0.35), subtitle, 11, MUTED)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    panel(slide, Inches(0), Inches(0), Inches(13.33), Inches(7.5), BG)
    panel(slide, Inches(0), Inches(0), Inches(0.24), Inches(7.5), TEAL)
    panel(slide, Inches(8.9), Inches(0), Inches(4.43), Inches(7.5), INK)
    panel(slide, Inches(8.25), Inches(0.75), Inches(3.95), Inches(1.2), CORAL, True)
    panel(slide, Inches(9.1), Inches(2.2), Inches(2.7), Inches(2.7), GOLD, True)
    panel(slide, Inches(8.5), Inches(5.25), Inches(3.4), Inches(0.65), TEAL, True)
    textbox(slide, Inches(0.75), Inches(1.15), Inches(5.0), Inches(0.25), "UPCOMING REVIEW", 12, TEAL, True)
    textbox(slide, Inches(0.75), Inches(1.7), Inches(7.0), Inches(1.3), "Executive Summary\nfrom Delivery Leadership Demos", 28, INK, True)
    textbox(slide, Inches(0.75), Inches(3.45), Inches(6.9), Inches(0.8), "Condensed view of what leaders responded to, what still needs proof, and what should be tightened before the next review.", 15, MUTED)
    textbox(slide, Inches(9.05), Inches(6.37), Inches(2.3), Inches(0.24), "Prepared from Teams transcript", 12, WHITE, True, PP_ALIGN.CENTER)


def summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "01", "Executive Takeaway", "Leadership interest is real, but the next review needs tighter proof, sharper scope boundaries, and clearer differentiation.")
    panel(slide, Inches(0.68), Inches(1.95), Inches(7.05), Inches(4.75), WHITE, True)
    textbox(slide, Inches(0.98), Inches(2.3), Inches(6.2), Inches(0.25), "Headline", 11, TEAL, True)
    textbox(slide, Inches(0.98), Inches(2.65), Inches(6.2), Inches(0.7), "The demos landed as credible innovation, but leadership repeatedly pulled the conversation back to delivery practicality.", 23, INK, True)
    bullet_list(
        slide,
        Inches(0.98),
        Inches(3.65),
        Inches(6.0),
        Inches(2.2),
        [
            "Strongest signals: workflow integration, measurable productivity, and fast onboarding potential.",
            "Main gaps: workflow explainability, scope clarity beyond current frameworks, and stronger evidence packaging.",
            "Next review should be built around proof, not feature narration.",
        ],
        size=14,
    )
    cards = [
        ("Critical asks", "03", CORAL, "Workflow proof\nScope boundaries\nBusiness evidence"),
        ("Open actions", "05", TEAL, "Deck tightening\nDemo reframing\nEvidence pack"),
        ("At-risk themes", "03", GOLD, "Scope drift\nReadability\nPositioning"),
    ]
    for idx, (label, number, color, note) in enumerate(cards):
        top = 2.05 + idx * 1.45
        panel(slide, Inches(8.0), Inches(top), Inches(4.55), Inches(1.15), WHITE, True)
        panel(slide, Inches(8.0), Inches(top), Inches(0.18), Inches(1.15), color)
        textbox(slide, Inches(8.35), Inches(top + 0.12), Inches(2.5), Inches(0.18), label, 11, MUTED, True)
        textbox(slide, Inches(8.35), Inches(top + 0.38), Inches(1.0), Inches(0.28), number, 24, INK, True)
        textbox(slide, Inches(9.5), Inches(top + 0.38), Inches(2.6), Inches(0.44), note, 10, MUTED)


def themes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "02", "Leadership Themes to Carry Forward", "The transcript shows a consistent pattern across the design-to-code, CUA testing, and GEO demos.")
    themes = [
        ("Workflow first", TEAL, [
            "Leaders repeatedly asked what the end-to-end flow is from requirement to delivery artifact.",
            "Standalone capability is not enough; the differentiator must be integrated workflow plus review guardrails.",
        ]),
        ("Scope discipline", GOLD, [
            "Questions surfaced on mobile/native support, framework coverage, and where the current solution is or is not production-ready.",
            "The next review should separate current-state support from roadmap experimentation.",
        ]),
        ("Proof over promise", CORAL, [
            "The strongest moments were quantified: ~21% frontend effort reduction, 40+ test cases already running, and known customization effort.",
            "Any claim without evidence was immediately tested by the audience.",
        ]),
    ]
    for idx, (title, color, bullets) in enumerate(themes):
        left = 0.68 + idx * 4.15
        panel(slide, Inches(left), Inches(2.1), Inches(3.75), Inches(4.45), WHITE, True)
        panel(slide, Inches(left), Inches(2.1), Inches(3.75), Inches(0.18), color)
        textbox(slide, Inches(left + 0.22), Inches(2.42), Inches(2.8), Inches(0.28), title, 16, INK, True)
        bullet_list(slide, Inches(left + 0.22), Inches(2.9), Inches(3.2), Inches(2.9), bullets, size=12)


def action_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "03", "Recommended Actions Before the Review", "Keep the review short, operational, and evidence-led.")
    panel(slide, Inches(0.68), Inches(2.0), Inches(6.0), Inches(4.8), WHITE, True)
    panel(slide, Inches(6.9), Inches(2.0), Inches(5.75), Inches(4.8), WHITE, True)
    textbox(slide, Inches(1.0), Inches(2.28), Inches(2.8), Inches(0.22), "What to show", 15, INK, True)
    bullet_list(
        slide,
        Inches(1.0),
        Inches(2.75),
        Inches(5.2),
        Inches(3.4),
        [
            "One clean workflow visual per solution: input, orchestration, guardrails, output.",
            "Current-state support table: frameworks, environments, dependencies, and exclusions.",
            "Evidence slide with measured benefit, implementation effort, and current maturity.",
            "One-slide differentiation versus manual delivery and external point tools.",
        ],
        size=13,
    )
    textbox(slide, Inches(7.2), Inches(2.28), Inches(2.8), Inches(0.22), "What to avoid", 15, INK, True)
    bullet_list(
        slide,
        Inches(7.2),
        Inches(2.75),
        Inches(4.9),
        Inches(3.4),
        [
            "Live demo navigation without visual clarity.",
            "Broad platform claims before support boundaries are explicit.",
            "Long feature explanations without a delivery outcome.",
            "Unqualified GEO positioning that triggers credibility or policy questions.",
        ],
        size=13,
    )
    panel(slide, Inches(0.95), Inches(6.0), Inches(11.3), Inches(0.5), INK, True)
    textbox(slide, Inches(1.25), Inches(6.13), Inches(10.6), Inches(0.2), "Review ask: align every demo to delivery challenge, proof point, adoption motion, and next decision needed from leadership.", 12, WHITE, True)


def next_review_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(slide, "04", "Next Review Talking Points", "Use these prompts to keep the discussion anchored on adoption and scale.")
    rows = [
        ("Design-to-code agent", "Show workflow from Jira + Figma to PR with review gates and current framework support.", "Differentiation versus SaaS and rollout feasibility"),
        ("CUA testing", "Bring customization effort, break-even logic, and report artifacts.", "Where it is best applied versus where it is not recommended"),
        ("OptiFlow GEO", "Tighten positioning, proof method, and what is advisory versus measurable output.", "Credibility of business model fit and implementation path"),
    ]
    panel(slide, Inches(0.68), Inches(2.0), Inches(12.0), Inches(4.9), WHITE, True)
    panel(slide, Inches(0.95), Inches(2.25), Inches(11.45), Inches(0.45), INK)
    headers = [("Solution", 1.08, 1.9), ("Lead with", 4.0, 3.4), ("Leadership concern to resolve", 8.0, 3.1)]
    for label, left, width in headers:
        textbox(slide, Inches(left), Inches(2.36), Inches(width), Inches(0.18), label, 11, WHITE, True)
    for idx, row in enumerate(rows):
        y = 2.8 + idx * 1.25
        fill = WHITE if idx % 2 == 0 else RGBColor(242, 244, 247)
        panel(slide, Inches(0.95), Inches(y), Inches(11.45), Inches(1.0), fill)
        textbox(slide, Inches(1.08), Inches(y + 0.12), Inches(2.6), Inches(0.55), row[0], 13, INK, True)
        textbox(slide, Inches(4.0), Inches(y + 0.12), Inches(3.5), Inches(0.58), row[1], 11, MUTED)
        textbox(slide, Inches(8.0), Inches(y + 0.12), Inches(3.0), Inches(0.58), row[2], 11, INK)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    summary_slide(prs)
    themes_slide(prs)
    action_slide(prs)
    next_review_slide(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Created {OUT_PATH}")


if __name__ == "__main__":
    main()
