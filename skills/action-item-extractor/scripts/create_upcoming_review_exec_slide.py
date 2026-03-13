from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(__file__).resolve().parents[3] / "docs" / "output" / "Upcoming Review - Exec Summary Slide.pptx"

BG = RGBColor(245, 245, 245)
HCL_BLUE = RGBColor(0, 107, 182)
HCL_BLUE_DARK = RGBColor(0, 76, 129)
HCL_BLUE_SOFT = RGBColor(226, 239, 249)
BLACK = RGBColor(13, 13, 13)
MUTED = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)
SLATE = RGBColor(230, 234, 238)
STEEL = RGBColor(74, 122, 163)


THEMES = [
    {
        "title": "Workflow Differentiation",
        "questions": [
            "What exactly is the workflow versus SaaS or manual delivery?",
            "What sits in Jira versus Figma, and how are non-visual requirements handled?",
            "Where are the human review and guardrail checkpoints before code reaches delivery?",
        ],
        "bring": "One clean flow from intake to PR with review gates and requirement handling beyond design artifacts.",
        "color": HCL_BLUE,
    },
    {
        "title": "Scope and Platform Support",
        "questions": [
            "Is the current scope only React/web, or does it truly extend to mobile and responsive cases?",
            "What is supported now versus only planned for Vue, Angular, React Native, Android, or iOS?",
            "Which use cases are production-ready and which are still experimental?",
        ],
        "bring": "A current-state scope matrix with explicit supported platforms, exclusions, and roadmap items.",
        "color": HCL_BLUE_DARK,
    },
    {
        "title": "Maturity and Rollout Readiness",
        "questions": [
            "Is this plug-and-play in practice, or does adoption require setup and project screening?",
            "Interactive today or autonomous tomorrow: what is the real maturity level?",
            "What dependencies are mandatory before rollout: APIs, MCP, access, repos, prompts?",
        ],
        "bring": "A rollout-readiness view: maturity stage, prerequisites, pilot profile, and operational ownership.",
        "color": STEEL,
    },
    {
        "title": "ROI and Benefit Proof",
        "questions": [
            "How is the claimed productivity benefit derived: effort, elapsed time, or both?",
            "Is the 21% saving repeatable, and under what delivery conditions?",
            "What is the baseline and sample size behind the benefit claims?",
        ],
        "bring": "A proof slide with baseline, sample size, before/after logic, and where the benefit does not hold.",
        "color": HCL_BLUE,
    },
    {
        "title": "Quality, Testing, and Guardrails",
        "questions": [
            "Does the solution recommend what tests to run, or only generate output?",
            "How do we prevent quality drift or careless behavior if the agent is doing more of the work?",
            "What audit evidence, traceability, and human approval remain in the loop?",
        ],
        "bring": "A guardrail model: recommended tests, review checkpoints, traceability, and accountability model.",
        "color": HCL_BLUE_DARK,
    },
    {
        "title": "CUA Reliability and Cost",
        "questions": [
            "How deterministic is the CUA flow and how do DOM, vision, and fallback actually work?",
            "What is the token/runtime cost story and where does caching reduce it?",
            "What implementation effort is required for onboarding, customization, and scaling?",
        ],
        "bring": "A practical reliability-cost view covering failover, best-fit scenarios, customization effort, and cost envelope.",
        "color": STEEL,
    },
    {
        "title": "GEO Credibility and Evidence",
        "questions": [
            "What makes the GEO approach credible beyond a demo narrative?",
            "How do impression metrics connect to business outcomes, not just optimization scores?",
            "What artifacts prove the output is actionable for delivery or clients?",
        ],
        "bring": "A clearer measurement chain from method to artifact to business relevance, with scope and proof boundaries.",
        "color": HCL_BLUE,
    },
]


def add_panel(slide, left, top, width, height, color, rounded=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_text(slide, left, top, width, height, text, size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=11, color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.bullet = True
    return box


def add_theme_card(slide, theme: dict, left: float, top: float, width: float, height: float) -> None:
    add_panel(slide, Inches(left), Inches(top), Inches(width), Inches(height), WHITE, True)
    add_panel(slide, Inches(left), Inches(top), Inches(width), Inches(0.16), theme["color"], True)
    add_text(slide, Inches(left + 0.18), Inches(top + 0.22), Inches(width - 0.36), Inches(0.22), theme["title"], 13, BLACK, True)
    add_text(slide, Inches(left + 0.18), Inches(top + 0.53), Inches(width - 0.36), Inches(0.18), "Pointed leadership questions", 9, MUTED, True)
    add_bullets(slide, Inches(left + 0.18), Inches(top + 0.76), Inches(width - 0.36), Inches(0.9), theme["questions"], 10, BLACK)
    add_panel(slide, Inches(left + 0.14), Inches(top + height - 0.44), Inches(width - 0.28), Inches(0.3), HCL_BLUE_SOFT, True)
    add_text(slide, Inches(left + 0.24), Inches(top + height - 0.37), Inches(width - 0.48), Inches(0.14), f"Bring: {theme['bring']}", 8, HCL_BLUE_DARK, True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    add_panel(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.18), HCL_BLUE)
    add_text(slide, Inches(0.56), Inches(0.42), Inches(1.4), Inches(0.2), "UPCOMING REVIEW", 10, HCL_BLUE, True)
    add_text(slide, Inches(0.56), Inches(0.76), Inches(7.8), Inches(0.42), "Leadership Question Themes from the Transcript", 24, BLACK, True)
    add_text(
        slide,
        Inches(0.56),
        Inches(1.18),
        Inches(8.6),
        Inches(0.28),
        "The transcript surfaces a broader set of pointed questions than the earlier summary. These should anchor the next review discussion.",
        11,
        MUTED,
    )

    add_panel(slide, Inches(9.0), Inches(0.48), Inches(3.78), Inches(0.96), HCL_BLUE, True)
    add_text(slide, Inches(9.25), Inches(0.7), Inches(3.3), Inches(0.18), "Leadership readout", 10, WHITE, True)
    add_text(slide, Inches(9.25), Inches(0.95), Inches(3.15), Inches(0.24), "The meeting generated specific questions on workflow, scope, proof, quality, reliability, and GEO credibility.", 10, WHITE)

    # 3x2 grid plus a bottom full-width GEO card
    card_w = 4.06
    card_h = 1.7
    x_positions = [0.56, 4.63, 8.7]
    y_positions = [1.72, 3.52]

    idx = 0
    for y in y_positions:
        for x in x_positions:
            add_theme_card(slide, THEMES[idx], x, y, card_w, card_h)
            idx += 1

    add_theme_card(slide, THEMES[6], 0.56, 5.32, 12.2, 1.5)

    add_panel(slide, Inches(0.56), Inches(6.98), Inches(12.2), Inches(0.22), HCL_BLUE_DARK, True)
    add_text(slide, Inches(0.82), Inches(7.0), Inches(11.7), Inches(0.14), "HCLTech styling: HCL Blue #006BB6, black #0D0D0D, white #FFFFFF, Aptos typography for office-safe delivery.", 8, WHITE)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Created {OUT_PATH}")


if __name__ == "__main__":
    main()
