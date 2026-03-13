from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(__file__).resolve().parents[1] / "templates" / "Executive Review Action Item Template.pptx"

# HCLTech-inspired palette using the brand's red-led identity with dark navy support.
BG = RGBColor(246, 246, 248)
NAVY = RGBColor(17, 25, 39)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(99, 115, 129)
HCL_RED = RGBColor(227, 24, 55)
HCL_RED_DARK = RGBColor(191, 18, 48)
STEEL = RGBColor(70, 86, 105)
SLATE = RGBColor(228, 232, 238)
WHITE = RGBColor(255, 255, 255)
SOFT_RED = RGBColor(252, 234, 238)
SOFT_STEEL = RGBColor(236, 241, 246)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_panel(slide, left, top, width, height, color, rounded=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_textbox(slide, left, top, width, height, text, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=12, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.bullet = True
    return box


def add_chip(slide, left, top, width, label, value, color):
    add_panel(slide, left, top, width, Inches(0.82), WHITE, True)
    add_panel(slide, left, top, Inches(0.16), Inches(0.82), color, True)
    add_textbox(slide, left + Inches(0.28), top + Inches(0.12), width - Inches(0.45), Inches(0.16), label, 10, MUTED, True)
    add_textbox(slide, left + Inches(0.28), top + Inches(0.34), Inches(0.8), Inches(0.24), value, 22, NAVY, True)


def add_status_pill(slide, left, top, label, color, dark_text=False):
    add_panel(slide, left, top, Inches(0.84), Inches(0.24), color, True)
    add_textbox(slide, left, top + Inches(0.01), Inches(0.84), Inches(0.16), label, 8, NAVY if dark_text else WHITE, True, PP_ALIGN.CENTER)


def build_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_panel(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.2), HCL_RED)
    add_panel(slide, Inches(0.55), Inches(0.48), Inches(1.72), Inches(0.24), HCL_RED, True)
    add_textbox(slide, Inches(0.72), Inches(0.5), Inches(1.1), Inches(0.18), "EXEC REVIEW", 10, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.62), Inches(0.86), Inches(7.8), Inches(0.5), "Leadership Questions, Intent, and Follow-up Actions", 24, NAVY, True)
    add_textbox(
        slide,
        Inches(0.62),
        Inches(1.35),
        Inches(7.9),
        Inches(0.38),
        "Single-slide format for turning transcript extraction into an HCLTech-style leadership review frame.",
        11,
        MUTED,
    )

    add_chip(slide, Inches(8.72), Inches(0.78), Inches(1.2), "Critical asks", "03", HCL_RED)
    add_chip(slide, Inches(10.03), Inches(0.78), Inches(1.2), "Open items", "04", STEEL)
    add_chip(slide, Inches(11.34), Inches(0.78), Inches(1.2), "Risks", "02", HCL_RED_DARK)

    add_panel(slide, Inches(0.62), Inches(1.95), Inches(2.65), Inches(4.95), NAVY, True)
    add_textbox(slide, Inches(0.9), Inches(2.22), Inches(1.9), Inches(0.18), "Meeting Signal", 11, HCL_RED, True)
    add_textbox(slide, Inches(0.9), Inches(2.52), Inches(2.02), Inches(0.92), "What leaders are really asking for", 22, WHITE, True)
    add_bullets(
        slide,
        Inches(0.9),
        Inches(3.75),
        Inches(2.0),
        Inches(2.1),
        [
            "Workflow clarity over feature narration",
            "Scope boundaries before roadmap claims",
            "Evidence strong enough for operational adoption",
        ],
        size=12,
        color=WHITE,
    )
    add_panel(slide, Inches(0.92), Inches(5.98), Inches(1.95), Inches(0.45), HCL_RED, True)
    add_textbox(slide, Inches(1.0), Inches(6.11), Inches(1.78), Inches(0.14), "Replace with meeting headline", 8, WHITE, True, PP_ALIGN.CENTER)

    add_panel(slide, Inches(3.45), Inches(1.95), Inches(9.22), Inches(4.95), WHITE, True)
    add_panel(slide, Inches(3.7), Inches(2.25), Inches(8.72), Inches(0.48), NAVY, True)

    columns = [
        ("Question / input heard", 3.83, 2.15),
        ("Underlying intent", 6.12, 1.65),
        ("Suggested follow-up", 7.92, 2.35),
        ("Owner", 10.42, 0.68),
        ("Due", 11.15, 0.5),
        ("Status", 11.78, 0.66),
    ]
    for label, left, width in columns:
        add_textbox(slide, Inches(left), Inches(2.38), Inches(width), Inches(0.16), label, 9, WHITE, True)

    rows = [
        (
            "Is it applicable for mobile as well, or just web development?",
            "Scope expansion and roadmap clarity",
            "State current support, exclusions, and next experiment path.",
            "Owner",
            "Date",
            ("Open", HCL_RED, False),
        ),
        (
            "For SaaS solutions in the market, what is the workflow?",
            "Differentiate integrated workflow from point tooling",
            "Show Jira to Figma to code to PR with review guardrails.",
            "Owner",
            "Date",
            ("Ready", STEEL, False),
        ),
        (
            "What is the productivity benefit? How is the percentage derived?",
            "Proof that the value is measurable and defensible",
            "Bring baseline, method, and before/after metrics on one slide.",
            "Owner",
            "Date",
            ("Proof", HCL_RED_DARK, False),
        ),
    ]

    for idx, row in enumerate(rows):
        top = 2.88 + idx * 1.08
        fill = WHITE if idx % 2 == 0 else SLATE
        add_panel(slide, Inches(3.7), Inches(top), Inches(8.72), Inches(0.92), fill, True)
        add_textbox(slide, Inches(3.85), Inches(top + 0.1), Inches(2.1), Inches(0.58), row[0], 10, INK, idx == 0)
        add_textbox(slide, Inches(6.12), Inches(top + 0.1), Inches(1.55), Inches(0.58), row[1], 10, MUTED)
        add_textbox(slide, Inches(7.92), Inches(top + 0.1), Inches(2.3), Inches(0.58), row[2], 10, INK)
        add_textbox(slide, Inches(10.42), Inches(top + 0.14), Inches(0.62), Inches(0.2), row[3], 9, MUTED)
        add_textbox(slide, Inches(11.15), Inches(top + 0.14), Inches(0.42), Inches(0.2), row[4], 9, MUTED)
        add_status_pill(slide, Inches(11.8), Inches(top + 0.2), row[5][0], row[5][1], row[5][2])

    add_panel(slide, Inches(3.7), Inches(6.22), Inches(4.18), Inches(0.46), SOFT_STEEL, True)
    add_textbox(slide, Inches(3.88), Inches(6.34), Inches(3.78), Inches(0.14), "Best use: limit the review to the 3 questions that should change the next meeting.", 9, NAVY, True)

    add_panel(slide, Inches(8.12), Inches(6.22), Inches(4.3), Inches(0.46), SOFT_RED, True)
    add_textbox(slide, Inches(8.3), Inches(6.34), Inches(3.9), Inches(0.14), "Rewrite raw transcript wording into short business language before using this slide.", 9, NAVY, True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_slide(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Created template at {OUT_PATH}")


if __name__ == "__main__":
    main()
