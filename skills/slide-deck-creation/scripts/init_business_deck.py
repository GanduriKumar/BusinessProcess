from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BLUE = RGBColor(21, 101, 192)
RED = RGBColor(198, 40, 40)
YELLOW = RGBColor(245, 180, 0)
GREEN = RGBColor(46, 125, 50)
TEXT = RGBColor(32, 33, 36)
MUTED = RGBColor(95, 99, 104)
BG = RGBColor(248, 249, 250)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(232, 234, 237)


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_top_band(slide):
    colors = [BLUE, RED, YELLOW, GREEN]
    x = 0
    for color in colors:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), 0, Inches(3.33325), Inches(0.16))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        x += 3.33325


def add_cover(prs: Presentation, spec: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CARD
    add_top_band(slide)
    add_text(slide, Inches(0.75), Inches(0.65), Inches(12.0), Inches(0.3), spec.get("kicker", "Business deck"), size=12, bold=True, color=BLUE)
    add_text(slide, Inches(0.75), Inches(1.15), Inches(11.5), Inches(1.0), spec["title"], size=28, bold=True)
    if spec.get("subtitle"):
        add_text(slide, Inches(0.75), Inches(2.2), Inches(11.0), Inches(0.6), spec["subtitle"], size=14, color=MUTED)


def add_bullet_slide(prs: Presentation, slide_spec: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_top_band(slide)
    add_text(slide, Inches(0.65), Inches(0.45), Inches(11.8), Inches(0.65), slide_spec["title"], size=24, bold=True)
    if slide_spec.get("subtitle"):
        add_text(slide, Inches(0.65), Inches(1.05), Inches(11.2), Inches(0.35), slide_spec["subtitle"], size=12, color=MUTED)
    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.75), Inches(11.2), Inches(4.8))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(slide_spec.get("bullets", [])):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.bullet = True
        p.level = 0
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(20)
            run.font.color.rgb = TEXT


def add_card_slide(prs: Presentation, slide_spec: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_top_band(slide)
    add_text(slide, Inches(0.65), Inches(0.45), Inches(11.8), Inches(0.65), slide_spec["title"], size=24, bold=True)
    cards = slide_spec.get("cards", [])
    width = Inches(3.85)
    for idx, card in enumerate(cards[:3]):
        left = Inches(0.75 + (idx * 4.15))
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.7), width, Inches(3.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD
        shape.line.color.rgb = LINE
        add_text(slide, left + Inches(0.2), Inches(1.95), width - Inches(0.3), Inches(0.35), card["title"], size=14, bold=True)
        add_text(slide, left + Inches(0.2), Inches(2.35), width - Inches(0.3), Inches(2.7), card["body"], size=13, color=MUTED)


def build_deck(spec: dict, output: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_cover(prs, spec)
    for slide_spec in spec.get("slides", []):
        kind = slide_spec.get("type", "bullets")
        if kind == "cards":
            add_card_slide(prs, slide_spec)
        else:
            add_bullet_slide(prs, slide_spec)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    parser = argparse.ArgumentParser(description="Create a simple business PPTX from a JSON slide spec.")
    parser.add_argument("--spec-json", required=True, help="Path to slide spec JSON")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    args = parser.parse_args()

    spec = json.loads(Path(args.spec_json).read_text(encoding="utf-8"))
    build_deck(spec, Path(args.output))
    print(f"Created {args.output}")


if __name__ == "__main__":
    main()
