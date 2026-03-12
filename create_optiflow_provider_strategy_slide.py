from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(
    r"C:\Users\kumar.gn\HCLProjects\BusinessProcess\docs\OptiFlow - Provider Mapping and Buyer Split Slide.pptx"
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

GOOGLE_BLUE = RGBColor(66, 133, 244)
GOOGLE_RED = RGBColor(234, 67, 53)
GOOGLE_YELLOW = RGBColor(251, 188, 5)
GOOGLE_GREEN = RGBColor(52, 168, 83)
TEXT = RGBColor(32, 33, 36)
MUTED = RGBColor(95, 99, 104)
LINE = RGBColor(232, 234, 237)
CARD = RGBColor(255, 255, 255)
BG = RGBColor(248, 249, 250)
GOOD = RGBColor(46, 125, 50)
WARN = RGBColor(245, 124, 0)
BAD = RGBColor(198, 40, 40)
FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def add_text(slide, left, top, width, height, text, size=16, color=TEXT, bold=False, align=PP_ALIGN.LEFT, font=FONT_BODY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_top_band(slide):
    x = 0
    for color in (GOOGLE_BLUE, GOOGLE_RED, GOOGLE_YELLOW, GOOGLE_GREEN):
        seg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), 0, Inches(3.33325), Inches(0.18))
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.line.fill.background()
        x += 3.33325


def add_card(slide, left, top, width, height, title, body, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.1), height)
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.3), Inches(0.25), title, size=13, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.4), width - Inches(0.3), height - Inches(0.5), body, size=11, color=MUTED)


def add_provider_table(slide, rows):
    left = Inches(0.45)
    top = Inches(1.2)
    widths = [Inches(2.05), Inches(0.8), Inches(2.4), Inches(3.0), Inches(4.15)]
    headers = ["Provider", "Fit", "Buyer motion", "Where OptiFlow fits", "Required change"]
    row_h = Inches(0.42)
    fit_fill = {"Green": RGBColor(232, 245, 233), "Yellow": RGBColor(255, 243, 224), "Red": RGBColor(252, 235, 233)}
    fit_text = {"Green": GOOD, "Yellow": WARN, "Red": BAD}

    x = left
    for idx, header in enumerate(headers):
        cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, widths[idx], row_h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GOOGLE_BLUE
        cell.line.color.rgb = BG
        add_text(slide, x + Inches(0.06), top + Inches(0.09), widths[idx] - Inches(0.12), Inches(0.2), header, size=9, color=CARD, bold=True)
        x += widths[idx]

    for r_idx, row in enumerate(rows, start=1):
        y = top + Inches(0.05) + row_h * r_idx
        vals = [row["provider"], row["fit"], row["motion"], row["fitment"], row["change"]]
        x = left
        for c_idx, val in enumerate(vals):
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, widths[c_idx], row_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fit_fill[row["fit"]] if c_idx == 1 else CARD
            cell.line.color.rgb = LINE
            add_text(
                slide,
                x + Inches(0.06),
                y + Inches(0.07),
                widths[c_idx] - Inches(0.12),
                Inches(0.27),
                val,
                size=8.5,
                color=fit_text[row["fit"]] if c_idx == 1 else TEXT,
                bold=c_idx in (0, 1),
            )
            x += widths[c_idx]


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_top_band(slide)

    add_text(slide, Inches(0.5), Inches(0.32), Inches(12.2), Inches(0.28), "Platform map", size=11, color=GOOGLE_BLUE, bold=True)
    add_text(slide, Inches(0.5), Inches(0.62), Inches(12.0), Inches(0.55), "Provider fit, buyer split, and go-to-market strategy for OptiFlow", size=24, bold=True, font=FONT_HEAD)

    rows = [
        {"provider": "Microsoft Bing / Copilot", "fit": "Green", "motion": "Platform + brands", "fitment": "Best fit for AI citations, Webmaster diagnostics, and advertiser optimization.", "change": "Add Bing AI Performance, IndexNow, Copilot citation, and ad diagnostics."},
        {"provider": "Google Search / AI Mode", "fit": "Yellow", "motion": "Platform + merchants", "fitment": "Strong fit for merchant readiness, feed-page consistency, and landing-page quality.", "change": "Drop llms.txt; focus on schema, crawlability, Merchant Center, snippets, and LP quality."},
        {"provider": "OpenAI ChatGPT", "fit": "Green", "motion": "Platform + sellers", "fitment": "Strong commerce fit for product feeds, retailer pages, and checkout readiness.", "change": "Build merchant-feed connectors, shopping readiness, and PDP quality scoring."},
        {"provider": "Perplexity", "fit": "Yellow", "motion": "Publishers + brands", "fitment": "Useful for citation readiness, answerability, and publisher trust signals.", "change": "Add citation-quality and source-trust diagnostics; keep claims advisory."},
        {"provider": "Brave / Yahoo", "fit": "Yellow", "motion": "Channel extension", "fitment": "Secondary fit via content clarity and Bing-linked search ecosystems.", "change": "Treat as indirect channel plays, not standalone platform motions."},
        {"provider": "DuckDuckGo", "fit": "Red", "motion": "Indirect only", "fitment": "Low direct fit; mostly routes through Microsoft search and ads.", "change": "Pursue only via Microsoft-led positioning."},
        {"provider": "Amazon Rufus / shopping AI", "fit": "Yellow", "motion": "Seller-first", "fitment": "Relevant for commerce content, catalog quality, and product attributes.", "change": "Extend from GEO into retail catalog, product attributes, and PDP conversion quality."},
    ]
    add_provider_table(slide, rows)

    add_card(
        slide,
        Inches(0.5),
        Inches(4.9),
        Inches(4.0),
        Inches(1.4),
        "Buyer split",
        "Platform edition: sell to search providers as ecosystem-quality, advertiser-readiness, and merchant-health tooling.\nBrand edition: sell to sellers, brands, publishers, and merchants as AI-search and commerce-readiness software.",
        GOOGLE_BLUE,
    )
    add_card(
        slide,
        Inches(4.68),
        Inches(4.9),
        Inches(3.95),
        Inches(1.4),
        "Best near-term buyer",
        "Brands, sellers, merchants, and publishers. The pain is immediate, budgets already exist, and the product can be sold without platform approval.",
        GOOGLE_GREEN,
    )
    add_card(
        slide,
        Inches(8.82),
        Inches(4.9),
        Inches(4.0),
        Inches(1.4),
        "Platform strategy",
        "Lead with Microsoft. Build a commerce-led variant for Google and OpenAI. Treat Perplexity as a narrower publisher play and Brave/Yahoo as secondary channels.",
        GOOGLE_YELLOW,
    )

    strategy = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.65))
    strategy.fill.solid()
    strategy.fill.fore_color.rgb = RGBColor(232, 240, 254)
    strategy.line.color.rgb = GOOGLE_BLUE
    add_text(
        slide,
        Inches(0.72),
        Inches(6.63),
        Inches(11.9),
        Inches(0.22),
        "Recommended strategy: sell Brand Edition first to prove ROI, then approach Microsoft first for Platform Edition, followed by Google/OpenAI with commerce-specific positioning.",
        size=13,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )

    footer = "Green = best fit. Yellow = viable with repositioning. Red = low direct fit. Strategy should separate Platform Edition from Brand Edition."
    add_text(slide, Inches(0.5), Inches(7.12), Inches(12.2), Inches(0.18), footer, size=8, color=MUTED, align=PP_ALIGN.RIGHT)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    build()
