from __future__ import annotations

import csv
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation


ROOT = Path(r"C:\Users\kumar.gn\HCLProjects\BusinessProcess")
PPTX_PATH = ROOT / "docs" / "input" / "Transformation_Slides_oneslider_v1.pptx"
XLSX_PATH = ROOT / "docs" / "input" / "genai_initiatives_table.xlsx"
OUT_MD = ROOT / "docs" / "output" / "Transformation_Programs_GenAI_Mapping.md"
OUT_CSV = ROOT / "docs" / "output" / "Transformation_Programs_GenAI_Mapping.csv"


MANUAL_MAPPING = {
    "Design to Code Generation - Digital": ("Figma-to-Code Plugin (Proposed)", "High"),
    "AI Code Review Bot": ("Code Change Analysis Tool (Proposed)", "High"),
    "Regulatory Report Processing": ("Legal Compliance Checks (LPM)", "Medium"),
    "Prompt Library": ("Reusable Prompt Template Library", "High"),
    "Tibco Document Generator": ("No clear match identified", "Low"),
    "Intelligent Compliance Email Automation": ("No clear match identified", "Low"),
    "Run deck Automation": ("No clear match identified", "Low"),
    "Agentic Self Healing of WIMM Issues": ("Auto-Fix Framework for Data Support", "Medium"),
    "DILO Analysis": ("No clear match identified", "Low"),
    "Certificate Management": ("No clear match identified", "Low"),
    "Optibot NeoX Mobile App": ("No clear match identified", "Low"),
    "Service now Chat and Service now Virtual agent": ("Knowledge Navigator", "High"),
    "WorkBlaze": ("Auto-Fix Framework for Data Support", "Low"),
    "Vulnerability Assessment and Resolution": ("No clear match identified", "Low"),
    "Gen AI Based Unit Testing - Digital": ("No clear match identified", "Low"),
    "Fixathon": ("No clear match identified", "Low"),
    "R3 Feature Extraction": ("No clear match identified", "Low"),
}


def extract_programs_from_pptx(path: Path) -> list[str]:
    prs = Presentation(path)
    programs: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            headers = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
            try:
                prog_idx = headers.index("Programs")
            except ValueError:
                continue
            for r in range(1, len(table.rows)):
                value = table.cell(r, prog_idx).text.strip()
                if value:
                    programs.append(value)
    return programs


def read_shared_strings(zf: zipfile.ZipFile) -> list[str]:
    ns = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    shared = []
    sst = ET.fromstring(zf.read("xl/sharedStrings.xml"))
    for si in sst.findall("m:si", ns):
        shared.append("".join(t.text or "" for t in si.findall(".//m:t", ns)))
    return shared


def extract_initiatives_from_xlsx(path: Path) -> dict[str, dict[str, str]]:
    ns = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    with zipfile.ZipFile(path) as zf:
        shared = read_shared_strings(zf)
        sh = ET.fromstring(zf.read("xl/worksheets/sheet1.xml"))
        initiatives: dict[str, dict[str, str]] = {}
        for row in sh.findall("m:sheetData/m:row", ns)[1:]:
            vals: dict[str, str] = {}
            for cell in row.findall("m:c", ns):
                col = re.match(r"([A-Z]+)", cell.attrib["r"]).group(1)
                v = cell.find("m:v", ns)
                if v is None:
                    continue
                value = v.text or ""
                if cell.attrib.get("t") == "s":
                    value = shared[int(value)]
                vals[col] = value
            name = vals.get("B", "").strip()
            desc = vals.get("E", "").strip()
            if name:
                initiatives[name] = {"name": name, "description": desc}
        return initiatives


def build_rows(programs: list[str], initiatives: dict[str, dict[str, str]]) -> list[dict[str, str]]:
    rows = []
    for idx, program in enumerate(programs, start=1):
        mapped_name = MANUAL_MAPPING.get(program, "No clear match identified")
        if isinstance(mapped_name, tuple):
            mapped_name, confidence = mapped_name
        else:
            confidence = "Low" if mapped_name == "No clear match identified" else "Medium"
        if mapped_name in initiatives:
            solution = f"{mapped_name} - {initiatives[mapped_name]['description']}"
        else:
            solution = mapped_name
        rows.append(
            {
                "Sl.No": str(idx),
                "Program Name - from pptx file": program,
                "GenAI solution - from Excel sheet": solution,
                "Match Confidence": confidence,
            }
        )
    return rows


def write_csv(rows: list[dict[str, str]], path: Path) -> None:
    with path.open("w", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(
            fh,
            fieldnames=["Sl.No", "Program Name - from pptx file", "GenAI solution - from Excel sheet", "Match Confidence"],
        )
        writer.writeheader()
        writer.writerows(rows)


def write_markdown(rows: list[dict[str, str]], path: Path) -> None:
    lines = [
        "# Transformation Programs to GenAI Mapping",
        "",
        "Best-effort mapping from the `Programs` column in `Transformation_Slides_oneslider_v1.pptx` to the Excel columns `Name of the GenAI Initiative` and `Short Description` in `genai_initiatives_table.xlsx`.",
        "",
        "| Sl.No | Program Name - from pptx file | GenAI solution - from Excel sheet | Match Confidence |",
        "|---:|---|---|---|",
    ]
    for row in rows:
        program = row["Program Name - from pptx file"].replace("|", "\\|")
        solution = row["GenAI solution - from Excel sheet"].replace("|", "\\|")
        lines.append(f"| {row['Sl.No']} | {program} | {solution} | {row['Match Confidence']} |")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> None:
    programs = extract_programs_from_pptx(PPTX_PATH)
    initiatives = extract_initiatives_from_xlsx(XLSX_PATH)
    rows = build_rows(programs, initiatives)
    OUT_MD.parent.mkdir(parents=True, exist_ok=True)
    write_markdown(rows, OUT_MD)
    write_csv(rows, OUT_CSV)
    print(f"Created {OUT_MD}")
    print(f"Created {OUT_CSV}")


if __name__ == "__main__":
    main()
