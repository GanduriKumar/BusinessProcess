from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\kumar.gn\HCLProjects\BusinessProcess")
OUT_DIR = ROOT / "docs" / "output"
PPT_OUT = OUT_DIR / "Customer_Service_Transformation_GenAI_Portfolio.pptx"
HTML_OUT = OUT_DIR / "Customer_Service_Transformation_GenAI_Portfolio.html"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HCL_BLUE = RGBColor(0, 107, 182)
HCL_BLUE_DARK = RGBColor(0, 76, 129)
HCL_BLUE_SOFT = RGBColor(226, 239, 249)
INK = RGBColor(13, 13, 13)
MUTED = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(246, 248, 250)
LINE = RGBColor(214, 221, 228)
GREEN = RGBColor(38, 132, 85)
AMBER = RGBColor(218, 145, 0)
SOFT_GREEN = RGBColor(233, 246, 239)
SOFT_AMBER = RGBColor(255, 246, 222)
SOFT_STEEL = RGBColor(239, 243, 246)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"

CARDS = [
    {
        "title": "AI Support Copilot",
        "theme": "Help support teams find the right answer faster and reduce repeated troubleshooting effort.",
        "examples": "Support Assistant, DocuJi",
        "value": "Knowledge retrieval, guided resolution, and faster engineer ramp-up for complex support cases.",
        "accent": HCL_BLUE,
        "fill": HCL_BLUE_SOFT,
    },
    {
        "title": "Bug Resolution Agent",
        "theme": "Automate issue analysis, fix generation, validation, and pull-request preparation.",
        "examples": "BugFixer 2.0",
        "value": "Shorter bug-resolution cycles, lower context switching, and safer enterprise bug fixing workflows.",
        "accent": GREEN,
        "fill": SOFT_GREEN,
    },
    {
        "title": "Web Test Agent",
        "theme": "Use computer-using agents to automate UI testing with less script maintenance.",
        "examples": "Web Test Agent, Regression Testing Agent",
        "value": "Faster regression execution, broader coverage, and more resilient UI automation for changing applications.",
        "accent": HCL_BLUE_DARK,
        "fill": SOFT_STEEL,
    },
    {
        "title": "Design-to-Code AI Agent",
        "theme": "Convert design intent into delivery-ready UI components with review and validation built in.",
        "examples": "Figma Design-to-Code AI Agent",
        "value": "Accelerated UI delivery, better design consistency, and reduced repetitive development effort.",
        "accent": AMBER,
        "fill": SOFT_AMBER,
    },
    {
        "title": "Composable Agent Platform",
        "theme": "Compose reusable agents into customer-specific workflows without locking into one stack.",
        "examples": "Composable Agent Platform",
        "value": "Faster solution assembly, multi-agent interoperability, and easier integration with existing experiences.",
        "accent": HCL_BLUE,
        "fill": HCL_BLUE_SOFT,
    },
    {
        "title": "Operations and Compliance Automation",
        "theme": "Improve run operations, alert handling, control checks, and service governance.",
        "examples": "Auto-Resolution / Alert Classification, AutoFix.AI, LPM Consolidator, RACE",
        "value": "Lower manual intervention, better compliance readiness, and stronger operational resilience.",
        "accent": GREEN,
        "fill": SOFT_GREEN,
    },
]


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, fill, line=LINE, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    add_panel(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), HCL_BLUE, HCL_BLUE, rounded=False)
    add_text(slide, Inches(0.58), Inches(0.38), Inches(2.2), Inches(0.18), "CUSTOMER CONVERSATION", size=10, bold=True, color=HCL_BLUE)
    add_text(slide, Inches(0.58), Inches(0.72), Inches(11.7), Inches(0.42), "GenAI Solutions Available for Service Transformation", size=24, bold=True, font=FONT_HEAD)
    add_text(slide, Inches(0.58), Inches(1.12), Inches(11.7), Inches(0.24), "A customer-facing portfolio view of practical GenAI solutions that can improve support, engineering, operations, and service productivity.", size=11, color=MUTED)

    add_panel(slide, Inches(0.58), Inches(1.55), Inches(12.1), Inches(0.62), WHITE)
    add_text(slide, Inches(0.82), Inches(1.78), Inches(11.6), Inches(0.16), "These solution areas are drawn from the internal GenAI solution inventory and offsite solution frameworks, and are presented here as practical transformation options for customer engagements.", size=10, color=INK)

    positions = [
        (0.58, 2.35), (4.42, 2.35), (8.26, 2.35),
        (0.58, 4.78), (4.42, 4.78), (8.26, 4.78),
    ]
    for card, (left, top) in zip(CARDS, positions):
        add_panel(slide, Inches(left), Inches(top), Inches(3.52), Inches(1.92), WHITE)
        add_panel(slide, Inches(left), Inches(top), Inches(3.52), Inches(0.16), card["accent"], card["accent"], rounded=False)
        add_text(slide, Inches(left + 0.16), Inches(top + 0.24), Inches(3.0), Inches(0.28), card["title"], size=13, bold=True)
        add_text(slide, Inches(left + 0.16), Inches(top + 0.56), Inches(3.05), Inches(0.42), card["theme"], size=10, color=MUTED)
        add_panel(slide, Inches(left + 0.16), Inches(top + 1.06), Inches(3.0), Inches(0.28), card["fill"], card["fill"])
        add_text(slide, Inches(left + 0.28), Inches(top + 1.15), Inches(2.7), Inches(0.1), f"Examples: {card['examples']}", size=8, bold=True, color=card["accent"] if card["accent"] != AMBER else INK)
        add_text(slide, Inches(left + 0.16), Inches(top + 1.42), Inches(3.03), Inches(0.3), card["value"], size=9)

    add_panel(slide, Inches(0.58), Inches(6.98), Inches(12.1), Inches(0.26), HCL_BLUE_DARK, HCL_BLUE_DARK)
    add_text(slide, Inches(0.82), Inches(7.02), Inches(11.6), Inches(0.14), "Customer takeaway: start with one or two high-friction service workflows, prove measurable value, and then scale the model across support, engineering, and operations.", size=8, color=WHITE)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_OUT)


def build_html() -> None:
    cards_html = []
    accent_map = {
        "rgb(0, 107, 182)": "var(--hcl-blue)",
    }
    for card in CARDS:
        accent = "#006bb6" if card["accent"] == HCL_BLUE else "#004c81" if card["accent"] == HCL_BLUE_DARK else "#268455" if card["accent"] == GREEN else "#da9100"
        fill = "#e2eff9" if card["fill"] == HCL_BLUE_SOFT else "#e9f6ef" if card["fill"] == SOFT_GREEN else "#fff6de" if card["fill"] == SOFT_AMBER else "#eff3f6"
        cards_html.append(
            f"""
            <article class="card" style="--accent:{accent}; --chip:{fill};">
              <h3>{card['title']}</h3>
              <p class="theme">{card['theme']}</p>
              <div class="chip">Examples: {card['examples']}</div>
              <p class="value">{card['value']}</p>
            </article>
            """
        )

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>GenAI Solutions Available for Service Transformation</title>
  <style>
    :root {{
      --hcl-blue:#006bb6;
      --hcl-blue-dark:#004c81;
      --ink:#0d0d0d;
      --muted:#5c5c5c;
      --line:#d6dde4;
      --bg:#f6f8fa;
      --card:#ffffff;
    }}
    * {{ box-sizing:border-box; }}
    body {{
      margin:0;
      font-family:"Aptos","Segoe UI",sans-serif;
      background:radial-gradient(circle at top left, rgba(0,107,182,0.12), transparent 26%), linear-gradient(180deg, #eef4f8 0%, #f8fafc 100%);
      color:var(--ink);
    }}
    .slide {{
      width:min(1280px, calc(100vw - 32px));
      margin:18px auto;
      background:var(--card);
      border:1px solid var(--line);
      border-radius:26px;
      box-shadow:0 18px 40px rgba(13,13,13,0.08);
      overflow:hidden;
    }}
    .topbar {{ height:10px; background:linear-gradient(90deg, var(--hcl-blue), #2b86c5 42%, #5ca9d8 100%); }}
    .wrap {{ padding:28px 30px 22px; }}
    .eyebrow {{ color:var(--hcl-blue); text-transform:uppercase; letter-spacing:.12em; font-size:.78rem; font-weight:700; }}
    h1 {{ margin:10px 0 8px; font-family:"Aptos Display","Aptos","Segoe UI",sans-serif; font-size:2rem; line-height:1.08; }}
    .sub {{ color:var(--muted); font-size:1rem; max-width:1080px; }}
    .summary {{
      margin-top:18px;
      padding:14px 18px;
      background:#fff;
      border:1px solid var(--line);
      border-radius:18px;
      font-size:.92rem;
      line-height:1.45;
    }}
    .grid {{
      margin-top:18px;
      display:grid;
      grid-template-columns:repeat(3, 1fr);
      gap:18px;
    }}
    .card {{
      position:relative;
      background:#fff;
      border:1px solid var(--line);
      border-radius:22px;
      padding:18px 18px 16px;
      min-height:196px;
      box-shadow:0 8px 18px rgba(13,13,13,0.04);
    }}
    .card::before {{
      content:"";
      position:absolute;
      inset:0 0 auto 0;
      height:8px;
      border-radius:22px 22px 0 0;
      background:var(--accent);
    }}
    h3 {{ margin:10px 0 10px; font-size:1.08rem; }}
    .theme {{ margin:0; color:var(--muted); font-size:.94rem; line-height:1.42; min-height:56px; }}
    .chip {{
      display:inline-block;
      margin-top:14px;
      padding:8px 12px;
      border-radius:999px;
      background:var(--chip);
      font-size:.8rem;
      font-weight:700;
    }}
    .value {{ margin-top:14px; font-size:.9rem; line-height:1.42; }}
    .footer {{
      margin-top:18px;
      padding:10px 16px;
      background:var(--hcl-blue-dark);
      color:#fff;
      border-radius:14px;
      font-size:.82rem;
    }}
    @media print {{
      body {{ background:#fff; }}
      .slide {{ width:100%; margin:0; border-radius:0; box-shadow:none; }}
    }}
  </style>
</head>
<body>
  <main class="slide">
    <div class="topbar"></div>
    <div class="wrap">
      <div class="eyebrow">Customer Conversation</div>
      <h1>GenAI Solutions Available for Service Transformation</h1>
      <div class="sub">A customer-facing portfolio view of practical GenAI solutions that can improve support, engineering, operations, and service productivity.</div>
      <div class="summary">These solution areas are drawn from the internal GenAI solution inventory and offsite solution frameworks, and are presented here as practical transformation options for customer engagements.</div>
      <section class="grid">
        {''.join(cards_html)}
      </section>
      <div class="footer">Customer takeaway: start with one or two high-friction service workflows, prove measurable value, and then scale the model across support, engineering, and operations.</div>
    </div>
  </main>
</body>
</html>
"""
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    HTML_OUT.write_text(html, encoding="utf-8")


def main() -> None:
    build_ppt()
    build_html()
    print(f"Created {PPT_OUT}")
    print(f"Created {HTML_OUT}")


if __name__ == "__main__":
    main()
