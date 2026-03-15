from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\kumar.gn\HCLProjects\BusinessProcess")
REQ_DOC = ROOT / "docs" / "output" / "SIP_App_Technical_Requirements.md"
OUT_DIR = ROOT / "docs" / "output"
PPT_OUT = OUT_DIR / "SIP_MVP_Technical_Architecture_Diagram.pptx"
HTML_OUT = OUT_DIR / "SIP_MVP_Technical_Architecture_Diagram.html"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HCL_BLUE = RGBColor(0, 107, 182)
HCL_BLUE_DARK = RGBColor(0, 76, 129)
HCL_BLUE_SOFT = RGBColor(226, 239, 249)
INK = RGBColor(13, 13, 13)
MUTED = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(245, 247, 250)
LINE = RGBColor(210, 218, 226)
GREEN = RGBColor(38, 132, 85)
AMBER = RGBColor(218, 145, 0)
SOFT_STEEL = RGBColor(238, 242, 246)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def ensure_requirements() -> None:
    if not REQ_DOC.exists():
        raise FileNotFoundError(f"Missing requirements document: {REQ_DOC}")


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=10, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        run = p.add_run()
        run.text = item
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, fill, line=LINE, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_header(slide):
    add_panel(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.2), HCL_BLUE, HCL_BLUE, rounded=False)
    add_text(slide, Inches(0.55), Inches(0.38), Inches(2.4), Inches(0.2), "MVP TECHNICAL ARCHITECTURE", size=10, bold=True, color=HCL_BLUE)
    add_text(slide, Inches(0.55), Inches(0.7), Inches(11.8), Inches(0.45), "Support Intelligence Platform MVP: React frontend, Python APIs, connector-based intake, and governed AI assistance", size=22, bold=True, font=FONT_HEAD)
    add_text(slide, Inches(0.55), Inches(1.14), Inches(11.6), Inches(0.24), "One-line takeaway: the MVP should separate UI, orchestration, intelligence, connectors, and governance so new ticketing systems and AI capabilities can be added without redesigning the core app.", size=10, color=MUTED)


def add_column(slide, left, top, width, title, accent, items):
    add_panel(slide, left, top, width, Inches(3.52), WHITE)
    add_panel(slide, left, top, width, Inches(0.14), accent, accent, rounded=False)
    add_text(slide, left + Inches(0.14), top + Inches(0.2), width - Inches(0.28), Inches(0.38), title, size=12, bold=True)
    chip_top = top + Inches(0.72)
    for idx, item in enumerate(items):
        y = chip_top + Inches(idx * 0.62)
        add_panel(slide, left + Inches(0.14), y, width - Inches(0.28), Inches(0.38), HCL_BLUE_SOFT if accent in (HCL_BLUE, HCL_BLUE_DARK) else SOFT_STEEL)
        add_text(slide, left + Inches(0.28), y + Inches(0.11), width - Inches(0.56), Inches(0.12), item, size=8, bold=True, color=accent if accent != AMBER else INK)


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    add_header(slide)

    columns = [
        ("Enterprise Sources", HCL_BLUE_DARK, ["ServiceNow or Jira connector", "Manual intake fallback", "Change and config sources", "Monitoring context"]),
        ("API and Core Platform", HCL_BLUE, ["Python REST API", "Auth and RBAC", "Canonical ticket model", "Workflow orchestration"]),
        ("Intelligence and UI", GREEN, ["Recommendation service", "Knowledge retrieval", "Generative UI composition", "React workbench"]),
        ("Governance and Data", AMBER, ["Audit trail", "Policy controls", "Transactional store", "Search or vector index"]),
    ]

    lefts = [0.55, 3.45, 6.35, 9.25]
    for idx, (title, accent, items) in enumerate(columns):
        add_column(slide, Inches(lefts[idx]), Inches(1.72), Inches(2.62), title, accent, items)

    for x in (3.18, 6.08, 8.98):
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(3.42), Inches(x + 0.22), Inches(3.42)).line.color.rgb = HCL_BLUE

    add_panel(slide, Inches(0.7), Inches(5.52), Inches(8.55), Inches(1.05), WHITE)
    add_text(slide, Inches(0.9), Inches(5.74), Inches(1.2), Inches(0.18), "MVP flow", size=11, bold=True, color=HCL_BLUE)
    add_text(slide, Inches(1.9), Inches(5.72), Inches(7.0), Inches(0.18), "ticket arrives -> normalize and secure -> enrich and retrieve -> generate recommendations and UI composition -> engineer reviews -> approved updates write back", size=10, color=INK)

    add_panel(slide, Inches(9.5), Inches(5.52), Inches(3.1), Inches(1.05), HCL_BLUE_SOFT, HCL_BLUE_SOFT)
    add_text(slide, Inches(9.7), Inches(5.72), Inches(2.7), Inches(0.18), "MVP boundaries", size=11, bold=True, color=HCL_BLUE)
    add_bullets(
        slide,
        Inches(9.68),
        Inches(5.98),
        Inches(2.6),
        Inches(0.5),
        [
            "one production connector",
            "engineer-first workbench",
            "advisory AI only",
        ],
        size=8,
    )

    add_text(slide, Inches(0.55), Inches(7.02), Inches(12.0), Inches(0.12), "Source: docs/output/SIP_App_Technical_Requirements.md | HCLTech styling: HCL Blue #006BB6 | Aptos typography", size=8, color=MUTED, align=PP_ALIGN.RIGHT)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_OUT)


def build_html() -> None:
    html = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>SIP MVP Technical Architecture Diagram</title>
  <style>
    :root {--hcl-blue:#006bb6;--hcl-blue-dark:#004c81;--hcl-blue-soft:#e2eff9;--ink:#0d0d0d;--muted:#5c5c5c;--bg:#f5f7fa;--line:#d2dae2;--green:#268455;--amber:#da9100;}
    * { box-sizing:border-box; }
    body { margin:0; font-family:"Aptos","Segoe UI",sans-serif; background:radial-gradient(circle at top left, rgba(0,107,182,0.14), transparent 28%), linear-gradient(180deg, #eef4f8 0%, #f7f9fb 100%); color:var(--ink); }
    .canvas { width:min(1320px, calc(100vw - 32px)); margin:18px auto; background:#fff; border:1px solid var(--line); border-radius:26px; box-shadow:0 18px 40px rgba(13,13,13,0.08); overflow:hidden; }
    .topbar { height:10px; background:linear-gradient(90deg, var(--hcl-blue), #2b86c5 42%, #5ca9d8 100%); }
    .wrap { padding:26px 28px 20px; }
    .eyebrow { color:var(--hcl-blue); text-transform:uppercase; letter-spacing:.12em; font-size:.78rem; font-weight:700; }
    h1 { margin:10px 0 8px; font-family:"Aptos Display","Aptos","Segoe UI",sans-serif; font-size:2rem; line-height:1.08; }
    .sub { color:var(--muted); font-size:.95rem; max-width:1120px; }
    .grid { margin-top:22px; display:grid; grid-template-columns:repeat(4, 1fr); gap:18px; align-items:start; }
    .col { position:relative; background:#fff; border:1px solid var(--line); border-radius:20px; padding:18px 16px 16px; min-height:330px; }
    .col::before { content:""; position:absolute; inset:0 0 auto 0; height:7px; border-radius:20px 20px 0 0; }
    .blue-dark::before { background:var(--hcl-blue-dark); }
    .blue::before { background:var(--hcl-blue); }
    .green::before { background:var(--green); }
    .amber::before { background:var(--amber); }
    .col h2 { margin:6px 0 14px; font-size:1rem; }
    .chip { margin:10px 0; padding:10px 12px; background:linear-gradient(180deg, #f6fbff, #eef4f8); border:1px solid var(--line); border-radius:14px; font-size:.86rem; font-weight:700; }
    .amber .chip, .green .chip { background:linear-gradient(180deg, #f8fafb, #eef2f5); }
    .flow { margin-top:18px; display:grid; grid-template-columns:2.4fr 1fr; gap:18px; }
    .panel { background:#fff; border:1px solid var(--line); border-radius:20px; padding:16px 18px; }
    .panel h3 { margin:0 0 10px; color:var(--hcl-blue); font-size:.95rem; }
    .panel p, .panel li { margin:0; font-size:.88rem; line-height:1.45; }
    .panel ul { margin:0; padding-left:18px; }
    .footer { margin-top:14px; color:var(--muted); font-size:.75rem; text-align:right; }
    @media print {
      body { background:#fff; }
      .canvas { width:100%; margin:0; border-radius:0; box-shadow:none; }
    }
  </style>
</head>
<body>
  <main class="canvas">
    <div class="topbar"></div>
    <div class="wrap">
      <div class="eyebrow">MVP Technical Architecture</div>
      <h1>Support Intelligence Platform MVP: React frontend, Python APIs, connector-based intake, and governed AI assistance</h1>
      <div class="sub">One-line takeaway: the MVP should separate UI, orchestration, intelligence, connectors, and governance so new ticketing systems and AI capabilities can be added without redesigning the core app.</div>

      <section class="grid">
        <article class="col blue-dark">
          <h2>Enterprise Sources</h2>
          <div class="chip">ServiceNow or Jira connector</div>
          <div class="chip">Manual intake fallback</div>
          <div class="chip">Change and config sources</div>
          <div class="chip">Monitoring context</div>
        </article>
        <article class="col blue">
          <h2>API and Core Platform</h2>
          <div class="chip">Python REST API</div>
          <div class="chip">Auth and RBAC</div>
          <div class="chip">Canonical ticket model</div>
          <div class="chip">Workflow orchestration</div>
        </article>
        <article class="col green">
          <h2>Intelligence and UI</h2>
          <div class="chip">Recommendation service</div>
          <div class="chip">Knowledge retrieval</div>
          <div class="chip">Generative UI composition</div>
          <div class="chip">React workbench</div>
        </article>
        <article class="col amber">
          <h2>Governance and Data</h2>
          <div class="chip">Audit trail</div>
          <div class="chip">Policy controls</div>
          <div class="chip">Transactional store</div>
          <div class="chip">Search or vector index</div>
        </article>
      </section>

      <section class="flow">
        <article class="panel">
          <h3>MVP flow</h3>
          <p>ticket arrives -> normalize and secure -> enrich and retrieve -> generate recommendations and UI composition -> engineer reviews -> approved updates write back</p>
        </article>
        <article class="panel" style="background:var(--hcl-blue-soft);">
          <h3>MVP boundaries</h3>
          <ul>
            <li>one production connector</li>
            <li>engineer-first workbench</li>
            <li>advisory AI only</li>
          </ul>
        </article>
      </section>

      <div class="footer">Source: docs/output/SIP_App_Technical_Requirements.md | HCLTech styling: HCL Blue #006BB6 | Aptos typography</div>
    </div>
  </main>
</body>
</html>
"""
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    HTML_OUT.write_text(html, encoding="utf-8")


def main() -> None:
    ensure_requirements()
    build_ppt()
    build_html()
    print(f"Created {PPT_OUT}")
    print(f"Created {HTML_OUT}")


if __name__ == "__main__":
    main()
