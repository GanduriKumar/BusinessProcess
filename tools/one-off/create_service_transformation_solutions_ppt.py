from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\kumar.gn\HCLProjects\BusinessProcess")
OUT = ROOT / "docs" / "output" / "TechVDU_GenAI_Service_Transformation_Solutions.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HCL_BLUE = RGBColor(0, 96, 169)
HCL_DARK = RGBColor(0, 43, 92)
HCL_LIGHT = RGBColor(232, 241, 250)
HCL_ACCENT = RGBColor(0, 163, 224)
HCL_GREEN = RGBColor(0, 135, 90)
HCL_ORANGE = RGBColor(232, 119, 34)
HCL_PURPLE = RGBColor(108, 63, 165)
WHITE = RGBColor(255, 255, 255)
GRAY_50 = RGBColor(248, 250, 252)
GRAY_100 = RGBColor(241, 245, 249)
GRAY_200 = RGBColor(226, 232, 240)
GRAY_500 = RGBColor(100, 116, 139)
GRAY_700 = RGBColor(51, 65, 85)
GRAY_800 = RGBColor(30, 41, 59)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"

VALUE_ITEMS = [
    ("15+", "Production-Ready AI Solutions"),
    ("25-30%", "Average Effort Reduction in QA"),
    ("60%", "Faster Code Migration"),
    ("~90%", "Bug Fix Success Rate"),
    ("30%", "Compliance Effort Savings"),
]

CATEGORIES = [
    (
        "Quality Engineering & Testing",
        [
            ("Agentic QA Lifecycle", "Multi-agent QA orchestration from test generation through execution, integrated with delivery tooling.", "25-30% productivity gain", "Piloted"),
            ("Regression Testing Agent", "Scriptless AI-led regression execution that reduces manual effort and accelerates release cycles.", "25-30% effort reduction", "Production"),
            ("Computer Use Agent", "Vision-led application testing across changing interfaces without depending on brittle selectors.", "15-25% less validation effort", "POC"),
            ("AI Agent Evaluation", "Standardized evaluation of AI agents for quality, safety, bias, and responsible AI readiness.", "Standardized AI safety checks", "Pilot"),
        ],
        HCL_BLUE,
    ),
    (
        "Intelligent Support & Operations",
        [
            ("Universal Bug Fixer 2.0", "Workflow-driven AI agent for enterprise bug resolution from issue intake through validation and pull request creation.", "18-25% effort savings, ~90% fix rate", "Production (5 programs)"),
            ("AI Support Assistant", "RAG-powered support guidance for ticket triage, troubleshooting, and faster incident handling.", "Faster resolution, fewer reopens", "Production (3+ environments)"),
            ("Self-Healing Data Operations", "Automated anomaly detection, alert classification, and guided remediation for data operations.", "30% incident reduction in 3-4 months", "Pilot"),
            ("Delivery Excellence Agent", "AI support for risk visibility, operational monitoring, and proactive delivery governance.", "Earlier risk detection", "In development"),
        ],
        HCL_GREEN,
    ),
    (
        "Developer Productivity & Modernization",
        [
            ("Migration Automation", "Automated migration support for enterprise platforms, including mapping, script generation, and data quality steps.", "Up to 60% effort reduction", "Production"),
            ("Figma-to-Code Generator", "Turns approved design inputs into delivery-ready UI components with less manual coding effort.", "Faster UI development", "Piloted (3 adopters)"),
            ("Copilot Adoption Suite", "Structured adoption of engineering copilots to improve developer and analyst productivity at scale.", "16% dev speed, 12% QA reduction", "Active"),
            ("Image Asset Automation", "AI-powered digital asset processing that reduces repetitive designer effort and speeds turnaround.", "70% less designer effort", "Production"),
        ],
        HCL_ORANGE,
    ),
    (
        "Governance, Compliance & Business Process",
        [
            ("RACE - Accessibility Engine", "Automates accessibility testing and bug capture for web and desktop applications.", "25-30% effort savings", "Production"),
            ("Compliance Automation", "AI-driven support for regulatory monitoring, metadata extraction, and vulnerability checks.", "30% compliance effort reduction", "Production (3 areas)"),
            ("RFP Response Automation", "GenAI support for creating tailored proposal content and faster response drafting.", "Faster proposals, higher win rate", "In development"),
            ("Multilingual Transcription", "Real-time transcription and translation to support globally distributed teams and stakeholders.", "Instant multilingual access", "Prototype"),
        ],
        HCL_PURPLE,
    ),
]


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=GRAY_800, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, fill, line=GRAY_200, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_header(slide, page_num: int):
    add_panel(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), HCL_DARK, HCL_DARK, rounded=False)
    add_text(slide, Inches(0.45), Inches(0.16), Inches(8.0), Inches(0.28), "GenAI-Powered Service Transformation Solutions", size=24, bold=True, color=WHITE, font=FONT_HEAD)
    add_text(slide, Inches(0.45), Inches(0.5), Inches(7.2), Inches(0.16), "Practical AI solutions to improve service delivery, quality, and operational efficiency", size=10, color=WHITE)
    add_text(slide, Inches(12.15), Inches(0.25), Inches(0.7), Inches(0.18), f"{page_num}/2", size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def add_value_strip(slide):
    add_panel(slide, Inches(0), Inches(0.85), SLIDE_W, Inches(0.7), GRAY_50, GRAY_200, rounded=False)
    start = 0.55
    for idx, (num, label) in enumerate(VALUE_ITEMS):
        left = Inches(start + idx * 2.52)
        add_text(slide, left, Inches(1.0), Inches(0.78), Inches(0.22), num, size=17, bold=True, color=HCL_BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.75), Inches(0.96), Inches(1.45), Inches(0.28), label, size=8.5, color=GRAY_500)
        if idx < len(VALUE_ITEMS) - 1:
            add_panel(slide, left + Inches(2.18), Inches(0.98), Inches(0.01), Inches(0.34), GRAY_200, GRAY_200, rounded=False)


def add_category(slide, left, top, width, title, accent, cards):
    add_text(slide, left, top, width, Inches(0.2), title.upper(), size=9.5, bold=True, color=GRAY_700)
    add_panel(slide, left, top + Inches(0.2), width, Inches(0.01), GRAY_200, GRAY_200, rounded=False)
    card_w = (width - Inches(0.3)) / 2
    positions = [
        (left, top + Inches(0.34)),
        (left + card_w + Inches(0.3), top + Inches(0.34)),
        (left, top + Inches(1.87)),
        (left + card_w + Inches(0.3), top + Inches(1.87)),
    ]
    for (cx, cy), (card_title, desc, metric, status) in zip(positions, cards):
        add_panel(slide, cx, cy, card_w, Inches(1.35), WHITE)
        add_panel(slide, cx, cy, card_w, Inches(0.06), accent, accent, rounded=False)
        add_text(slide, cx + Inches(0.12), cy + Inches(0.11), card_w - Inches(0.24), Inches(0.24), card_title, size=10.5, bold=True, color=HCL_DARK)
        add_text(slide, cx + Inches(0.12), cy + Inches(0.35), card_w - Inches(0.24), Inches(0.5), desc, size=7.6, color=GRAY_700)
        add_panel(slide, cx + Inches(0.12), cy + Inches(0.88), card_w - Inches(0.24), Inches(0.18), HCL_LIGHT, HCL_LIGHT)
        add_text(slide, cx + Inches(0.18), cy + Inches(0.915), card_w - Inches(0.36), Inches(0.09), metric, size=7.1, bold=True, color=accent)
        add_text(slide, cx + Inches(0.12), cy + Inches(1.1), card_w - Inches(0.24), Inches(0.1), status, size=6.9, color=GRAY_500)


def add_footer(slide):
    add_panel(slide, Inches(0), Inches(7.12), SLIDE_W, Inches(0.38), HCL_DARK, HCL_DARK, rounded=False)
    add_text(slide, Inches(0.4), Inches(7.21), Inches(12.45), Inches(0.11), "All solutions are built on enterprise-grade safety controls, provider-abstracted AI, and auditable governance - ready for phased adoption.", size=8, color=WHITE)


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = WHITE
    add_header(slide1, 1)
    add_value_strip(slide1)
    add_category(slide1, Inches(0.45), Inches(1.82), Inches(6.1), CATEGORIES[0][0], CATEGORIES[0][2], CATEGORIES[0][1])
    add_category(slide1, Inches(6.78), Inches(1.82), Inches(6.1), CATEGORIES[1][0], CATEGORIES[1][2], CATEGORIES[1][1])
    add_footer(slide1)

    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE
    add_header(slide2, 2)
    add_value_strip(slide2)
    add_category(slide2, Inches(0.45), Inches(1.82), Inches(6.1), CATEGORIES[2][0], CATEGORIES[2][2], CATEGORIES[2][1])
    add_category(slide2, Inches(6.78), Inches(1.82), Inches(6.1), CATEGORIES[3][0], CATEGORIES[3][2], CATEGORIES[3][1])
    add_footer(slide2)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Created {OUT}")


if __name__ == "__main__":
    build()
