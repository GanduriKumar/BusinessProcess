from __future__ import annotations

from pathlib import Path
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\kumar.gn\HCLProjects\BusinessProcess")
SOURCE_DOC = ROOT / "docs" / "input" / "SIP_Executive_Proposal_v2.docx"
OUTPUT_DIR = ROOT / "docs" / "output"
PPT_OUTPUT = OUTPUT_DIR / "SIP_Executive_Leadership_Deck.pptx"
HTML_OUTPUT = OUTPUT_DIR / "SIP_Executive_Leadership_Deck.html"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HCL_BLUE = RGBColor(0, 107, 182)
HCL_BLUE_DARK = RGBColor(0, 76, 129)
HCL_BLUE_SOFT = RGBColor(226, 239, 249)
INK = RGBColor(13, 13, 13)
MUTED = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(245, 247, 250)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(212, 220, 228)
SOFT_STEEL = RGBColor(238, 242, 246)
RISK = RGBColor(194, 57, 43)
AMBER = RGBColor(218, 145, 0)
GREEN = RGBColor(38, 132, 85)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def ensure_source_doc() -> None:
    if not SOURCE_DOC.exists():
        raise FileNotFoundError(f"Missing source document: {SOURCE_DOC}")


def extract_doc_paragraphs(path: Path) -> list[str]:
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("word/document.xml")
    root = ET.fromstring(xml)
    paragraphs = []
    for p in root.findall(".//w:p", ns):
        parts = [t.text for t in p.findall(".//w:t", ns) if t.text]
        if parts:
            paragraphs.append("".join(parts).strip())
    return [p for p in paragraphs if p]


def source_signals(paragraphs: list[str]) -> dict[str, bool]:
    joined = "\n".join(paragraphs).lower()
    return {
        "knowledge_decay": "knowledge decay" in joined,
        "phase_zero": "phase 0" in joined,
        "human_gate": "human confirmation" in joined,
        "co_design": "co-design" in joined or "co-design sessions" in joined,
    }


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, size=16, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        p.level = 0
        run = p.add_run()
        run.text = bullet
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, fill, line=None, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def add_header(slide, kicker, title, subtitle=None):
    add_panel(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.2), HCL_BLUE, rounded=False)
    add_text(slide, Inches(0.58), Inches(0.38), Inches(2.6), Inches(0.22), kicker.upper(), size=10, bold=True, color=HCL_BLUE)
    add_text(slide, Inches(0.58), Inches(0.72), Inches(11.8), Inches(0.56), title, size=24, bold=True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, Inches(0.58), Inches(1.28), Inches(11.5), Inches(0.36), subtitle, size=11, color=MUTED)


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_footer(slide, text="Support Intelligence Platform | Executive proposal"):
    add_text(slide, Inches(0.58), Inches(7.06), Inches(12.0), Inches(0.15), text, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_panel(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.22), HCL_BLUE, rounded=False)
    add_panel(slide, Inches(0.72), Inches(0.72), Inches(2.0), Inches(0.28), HCL_BLUE, rounded=True)
    add_text(slide, Inches(0.87), Inches(0.77), Inches(1.7), Inches(0.12), "LEADERSHIP DISCUSSION", size=10, bold=True, color=WHITE)
    add_text(
        slide,
        Inches(0.72),
        Inches(1.35),
        Inches(8.8),
        Inches(1.4),
        "Reimagining Support Services with GenAI Automation and Generative UI",
        size=28,
        bold=True,
        font=FONT_HEAD,
    )
    add_text(
        slide,
        Inches(0.72),
        Inches(2.75),
        Inches(8.6),
        Inches(0.7),
        "A business-story proposal to improve support effectiveness, scale expertise, and make the service model easier to operate and experience.",
        size=15,
        color=MUTED,
    )

    themes = [
        ("Why now", "Support demand, complexity, and customer expectations are moving faster than traditional support methods."),
        ("What changes", "Use GenAI automation plus Generative UI to guide work, reduce rework, and surface the right context earlier."),
        ("What we need", "Leadership input on business priority, pilot scope, sponsorship, and readiness to move forward."),
    ]
    lefts = [0.72, 4.4, 8.08]
    for idx, (title, body) in enumerate(themes):
        add_panel(slide, Inches(lefts[idx]), Inches(4.55), Inches(3.1), Inches(1.7), CARD, LINE)
        add_panel(slide, Inches(lefts[idx]), Inches(4.55), Inches(0.12), Inches(1.7), HCL_BLUE if idx == 0 else HCL_BLUE_DARK if idx == 1 else GREEN)
        add_text(slide, Inches(lefts[idx] + 0.2), Inches(4.78), Inches(2.75), Inches(0.25), title, size=13, bold=True)
        add_text(slide, Inches(lefts[idx] + 0.2), Inches(5.12), Inches(2.75), Inches(0.86), body, size=11, color=MUTED)

    add_footer(slide, "HCLTech styling: HCL Blue #006BB6 | Aptos typography | Executive discussion draft")


def slide_name_expansion(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "What SIP Means", "Support Intelligence Platform is the working name for the proposed support modernization layer.")

    add_panel(slide, Inches(0.9), Inches(1.95), Inches(11.5), Inches(1.0), WHITE, LINE)
    add_text(slide, Inches(1.2), Inches(2.25), Inches(10.9), Inches(0.28), "SIP = Support Intelligence Platform", size=22, bold=True, color=HCL_BLUE, align=PP_ALIGN.CENTER)

    expansions = [
        ("Support", "The focus is the day-to-day support model: ticket handling, investigation, escalation, and service experience.", HCL_BLUE),
        ("Intelligence", "The differentiator is the ability to turn past tickets, context, and expertise into useful guidance at the moment of need.", HCL_BLUE_DARK),
        ("Platform", "The proposal is not a single feature. It is a reusable layer that can support engineers, managers, and clients over time.", GREEN),
    ]
    for idx, (title, body, accent) in enumerate(expansions):
        left = 0.9 + idx * 3.85
        add_panel(slide, Inches(left), Inches(3.45), Inches(3.45), Inches(2.15), CARD, LINE)
        add_panel(slide, Inches(left), Inches(3.45), Inches(3.45), Inches(0.14), accent, accent, rounded=False)
        add_text(slide, Inches(left + 0.2), Inches(3.78), Inches(2.8), Inches(0.24), title, size=16, bold=True)
        add_text(slide, Inches(left + 0.2), Inches(4.22), Inches(2.9), Inches(0.92), body, size=11, color=MUTED)

    add_panel(slide, Inches(1.2), Inches(6.1), Inches(10.8), Inches(0.52), HCL_BLUE_SOFT, HCL_BLUE_SOFT)
    add_text(slide, Inches(1.45), Inches(6.27), Inches(10.3), Inches(0.14), "Plain-language takeaway: SIP is a smarter support operating layer that helps teams use what they already know, more consistently and at scale.", size=10, bold=True, color=HCL_BLUE, align=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_current_challenge(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "The Current Challenge", "Support works today, but it is harder to scale, standardize, and sustain.")

    add_text(slide, Inches(0.72), Inches(1.72), Inches(4.8), Inches(0.24), "What leaders are seeing", size=12, bold=True, color=HCL_BLUE)
    bullets = [
        "Manual effort remains high, especially during triage, context gathering, and handoffs.",
        "Information is spread across tickets, engineer memory, tools, and client-specific documentation.",
        "Resolution quality varies because the starting context depends on who picks up the issue.",
        "Senior expertise solves problems quickly, but that expertise is difficult to reuse at scale.",
    ]
    add_bullets(slide, Inches(0.78), Inches(2.0), Inches(5.2), Inches(3.7), bullets, size=16)

    add_text(slide, Inches(6.2), Inches(1.72), Inches(5.8), Inches(0.24), "Where the friction shows up", size=12, bold=True, color=HCL_BLUE)
    stages = [
        ("Ticket lands", "The issue arrives with limited context and uneven documentation."),
        ("Engineer searches", "Time is spent hunting through history, tools, and tribal knowledge."),
        ("Escalation happens", "The next tier often repeats discovery because context is incomplete."),
        ("Knowledge fades", "The resolution is closed, but the reasoning is not easy to reuse next time."),
    ]
    top = 2.02
    for idx, (title, body) in enumerate(stages):
        y = Inches(top + (idx * 1.02))
        add_panel(slide, Inches(6.18), y, Inches(5.95), Inches(0.8), WHITE, LINE)
        add_panel(slide, Inches(6.32), y + Inches(0.17), Inches(0.45), Inches(0.45), HCL_BLUE_SOFT, HCL_BLUE_SOFT)
        add_text(slide, Inches(6.46), y + Inches(0.2), Inches(0.18), Inches(0.15), str(idx + 1), size=10, bold=True, color=HCL_BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(6.9), y + Inches(0.12), Inches(2.0), Inches(0.2), title, size=12, bold=True)
        add_text(slide, Inches(6.9), y + Inches(0.34), Inches(4.8), Inches(0.28), body, size=10, color=MUTED)

    add_footer(slide)


def slide_why_now(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "Why This Matters Now", "This is now a business issue, not just an operations issue.")

    pressure_cards = [
        ("Cost pressure", "Leaders need productivity gains without defaulting to linear headcount growth.", HCL_BLUE),
        ("Speed pressure", "Customers and account teams expect faster answers, clearer status, and less back-and-forth.", HCL_BLUE_DARK),
        ("Quality pressure", "Service consistency becomes harder when knowledge is fragmented and work is handed across tiers.", AMBER),
        ("Market pressure", "AI-enabled service models are changing what clients will expect from managed support partners.", GREEN),
    ]
    for idx, (title, body, accent) in enumerate(pressure_cards):
        left = 0.78 + (idx % 2) * 3.2
        top = 1.95 + (idx // 2) * 1.7
        add_panel(slide, Inches(left), Inches(top), Inches(2.9), Inches(1.35), CARD, LINE)
        add_panel(slide, Inches(left), Inches(top), Inches(0.11), Inches(1.35), accent)
        add_text(slide, Inches(left + 0.22), Inches(top + 0.18), Inches(2.45), Inches(0.2), title, size=12, bold=True)
        add_text(slide, Inches(left + 0.22), Inches(top + 0.44), Inches(2.45), Inches(0.6), body, size=10, color=MUTED)

    add_panel(slide, Inches(7.35), Inches(1.98), Inches(5.2), Inches(3.55), WHITE, LINE)
    add_text(slide, Inches(7.62), Inches(2.18), Inches(4.6), Inches(0.22), "If we continue as-is", size=12, bold=True, color=HCL_BLUE)
    add_text(slide, Inches(7.62), Inches(2.55), Inches(4.35), Inches(0.7), "Support remains reactive and heavily person-dependent.", size=17, bold=True)
    add_text(slide, Inches(7.62), Inches(3.18), Inches(4.45), Inches(1.25), "That creates operational drag, slows responsiveness, and makes it harder to turn years of delivery experience into a durable advantage.", size=12, color=MUTED)
    add_panel(slide, Inches(7.62), Inches(4.55), Inches(4.35), Inches(0.52), HCL_BLUE_SOFT, HCL_BLUE_SOFT)
    add_text(slide, Inches(7.82), Inches(4.71), Inches(4.0), Inches(0.16), "The timing matters because the service model is under pressure from both inside and outside the organization.", size=9, bold=True, color=HCL_BLUE)

    add_footer(slide)


def slide_opportunity(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "The Opportunity", "Move from reactive support to guided, assisted, and progressively more intelligent support.")

    add_panel(slide, Inches(0.85), Inches(1.95), Inches(3.8), Inches(3.85), CARD, LINE)
    add_text(slide, Inches(1.1), Inches(2.18), Inches(3.2), Inches(0.22), "Today", size=12, bold=True, color=RISK)
    today = [
        "The engineer gathers context manually.",
        "Knowledge retrieval depends on who knows where to look.",
        "Escalations often restart the investigation.",
        "The service experience feels fragmented to both engineers and clients.",
    ]
    add_bullets(slide, Inches(1.05), Inches(2.55), Inches(3.0), Inches(2.6), today, size=14, color=INK)

    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.95), Inches(3.85), Inches(6.05), Inches(3.85)).line.color.rgb = HCL_BLUE
    add_panel(slide, Inches(5.32), Inches(3.48), Inches(1.15), Inches(0.7), HCL_BLUE, HCL_BLUE, rounded=True)
    add_text(slide, Inches(5.52), Inches(3.71), Inches(0.76), Inches(0.12), "Shift", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_panel(slide, Inches(6.65), Inches(1.95), Inches(5.75), Inches(3.85), WHITE, LINE)
    add_text(slide, Inches(6.95), Inches(2.18), Inches(4.9), Inches(0.22), "Future state", size=12, bold=True, color=GREEN)
    future_cards = [
        ("Guided start", "Relevant past resolutions, context, and next best actions appear at ticket open."),
        ("Smarter workbench", "The interface adapts to the support tier, issue state, and user role."),
        ("Reusable expertise", "Reasoning paths, not just fixes, become part of the service knowledge base."),
        ("Better visibility", "Managers and clients see a clearer, more consistent support experience."),
    ]
    for idx, (title, body) in enumerate(future_cards):
        left = 6.95 + (idx % 2) * 2.45
        top = 2.6 + (idx // 2) * 1.25
        add_panel(slide, Inches(left), Inches(top), Inches(2.18), Inches(0.92), HCL_BLUE_SOFT if idx in (0, 3) else SOFT_STEEL, LINE)
        add_text(slide, Inches(left + 0.14), Inches(top + 0.12), Inches(1.85), Inches(0.18), title, size=10, bold=True)
        add_text(slide, Inches(left + 0.14), Inches(top + 0.34), Inches(1.85), Inches(0.4), body, size=9, color=MUTED)

    add_footer(slide)


def slide_proposal(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "What We Are Proposing", "Combine GenAI automation and Generative UI to augment support operations, not replace them.")

    add_panel(slide, Inches(0.86), Inches(2.3), Inches(3.6), Inches(2.4), CARD, LINE)
    add_text(slide, Inches(1.08), Inches(2.55), Inches(3.0), Inches(0.25), "GenAI automation", size=16, bold=True, color=HCL_BLUE)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(3.0),
        Inches(3.05),
        Inches(1.3),
        [
            "Summarizes issues and extracts clean context from messy tickets.",
            "Retrieves relevant historical cases and knowledge quickly.",
            "Recommends actions, hypotheses, and escalation guidance.",
        ],
        size=13,
    )

    add_panel(slide, Inches(4.92), Inches(2.75), Inches(0.95), Inches(0.7), HCL_BLUE, HCL_BLUE, rounded=True)
    add_text(slide, Inches(5.05), Inches(2.98), Inches(0.7), Inches(0.12), "+", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_panel(slide, Inches(6.28), Inches(2.3), Inches(3.75), Inches(2.4), CARD, LINE)
    add_text(slide, Inches(6.5), Inches(2.55), Inches(3.15), Inches(0.25), "Generative UI", size=16, bold=True, color=HCL_BLUE_DARK)
    add_bullets(
        slide,
        Inches(6.42),
        Inches(3.0),
        Inches(3.15),
        Inches(1.3),
        [
            "Builds role-based experiences for engineers, managers, and clients.",
            "Adapts the workspace to the issue context and workflow stage.",
            "Presents AI guidance in a practical, auditable operating surface.",
        ],
        size=13,
    )

    add_panel(slide, Inches(10.45), Inches(2.3), Inches(1.95), Inches(2.4), HCL_BLUE_SOFT, HCL_BLUE_SOFT)
    add_text(slide, Inches(10.67), Inches(2.6), Inches(1.45), Inches(0.45), "Targeted augmentation", size=15, bold=True, color=HCL_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10.62), Inches(3.25), Inches(1.55), Inches(0.95), "The value comes from improving the work itself: better context, better decisions, and smoother execution.", size=10, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.88), Inches(5.35), Inches(11.2), Inches(0.28), "Simple framing for leaders", size=12, bold=True, color=HCL_BLUE)
    add_panel(slide, Inches(0.88), Inches(5.7), Inches(11.1), Inches(0.72), WHITE, LINE)
    add_text(slide, Inches(1.12), Inches(5.92), Inches(10.55), Inches(0.28), "This is a support operating model upgrade powered by AI. It helps teams start faster, investigate with better context, and preserve expertise more systematically.", size=13)

    add_footer(slide)


def slide_use_cases(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "How It Would Work in Practice", "The concept is tangible because it maps directly to everyday support work.")

    use_cases = [
        (
            "Agent assist at ticket open",
            "When a ticket is opened, the system organizes the issue, surfaces likely matches, and suggests a clean first response path.",
            ["Faster triage", "Less manual searching", "Cleaner escalation context"],
            HCL_BLUE,
        ),
        (
            "Knowledge guidance during investigation",
            "As the case moves across tiers, relevant environment notes, recent changes, and next-best actions appear in the flow of work.",
            ["Better investigation quality", "More reusable reasoning", "Lower rework across tiers"],
            HCL_BLUE_DARK,
        ),
        (
            "Dynamic UI by role and workflow stage",
            "The workspace changes based on who is using it and where the issue is in the lifecycle, so each user sees the right guidance and controls.",
            ["Less clutter", "Higher adoption potential", "More consistent experience"],
            GREEN,
        ),
    ]

    for idx, (title, body, outcomes, accent) in enumerate(use_cases):
        left = 0.74 + idx * 4.15
        add_panel(slide, Inches(left), Inches(2.0), Inches(3.6), Inches(3.75), CARD, LINE)
        add_panel(slide, Inches(left), Inches(2.0), Inches(3.6), Inches(0.16), accent, accent, rounded=False)
        add_text(slide, Inches(left + 0.18), Inches(2.3), Inches(3.0), Inches(0.42), title, size=14, bold=True)
        add_text(slide, Inches(left + 0.18), Inches(2.9), Inches(3.1), Inches(1.18), body, size=11, color=MUTED)
        add_text(slide, Inches(left + 0.18), Inches(4.35), Inches(2.3), Inches(0.2), "What this unlocks", size=10, bold=True, color=accent)
        for jdx, item in enumerate(outcomes):
            add_panel(slide, Inches(left + 0.18), Inches(4.68 + (jdx * 0.32)), Inches(2.95), Inches(0.22), HCL_BLUE_SOFT if accent != GREEN else SOFT_STEEL, rounded=True)
            add_text(slide, Inches(left + 0.3), Inches(4.73 + (jdx * 0.32)), Inches(2.65), Inches(0.1), item, size=8, bold=True, color=accent if accent != GREEN else GREEN)

    add_footer(slide)


def slide_business_value(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "Business Value", "The value is operational, financial, and experiential.")

    add_panel(slide, Inches(0.86), Inches(1.95), Inches(12.0), Inches(1.0), WHITE, LINE)
    add_text(slide, Inches(1.1), Inches(2.28), Inches(11.5), Inches(0.28), "Expected outcome: support becomes easier to run, easier to improve, and easier for customers to experience.", size=17, bold=True)

    columns = [
        ("Operational value", ["Faster resolution flow", "Better context at handoff", "More consistent service execution"], HCL_BLUE),
        ("Financial value", ["Improved productivity leverage", "Reduced repeated effort", "Stronger scale potential without defaulting to headcount"], HCL_BLUE_DARK),
        ("Experience value", ["Better agent confidence", "Cleaner client communication", "Greater visibility into service health"], GREEN),
    ]
    for idx, (title, bullets, accent) in enumerate(columns):
        left = 0.9 + idx * 4.02
        add_panel(slide, Inches(left), Inches(3.35), Inches(3.7), Inches(2.55), CARD, LINE)
        add_panel(slide, Inches(left), Inches(3.35), Inches(0.12), Inches(2.55), accent)
        add_text(slide, Inches(left + 0.22), Inches(3.62), Inches(3.1), Inches(0.24), title, size=14, bold=True)
        add_bullets(slide, Inches(left + 0.18), Inches(4.05), Inches(3.05), Inches(1.4), bullets, size=13)

    add_footer(slide)


def slide_risk(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "Risk, Governance, and Adoption", "The proposal should be governed, human-in-the-loop, phased, and measurable.")

    headers = ["Leadership concern", "What it means", "Control approach"]
    widths = [3.1, 3.4, 5.1]
    x = 0.88
    for idx, header in enumerate(headers):
        add_panel(slide, Inches(x), Inches(1.9), Inches(widths[idx]), Inches(0.42), HCL_BLUE, HCL_BLUE, rounded=False)
        add_text(slide, Inches(x + 0.12), Inches(2.03), Inches(widths[idx] - 0.18), Inches(0.12), header, size=10, bold=True, color=WHITE)
        x += widths[idx]

    rows = [
        ("Accuracy and hallucination risk", "AI can produce plausible guidance that still needs judgment.", "Use confidence signals, human confirmation, and feedback loops."),
        ("Security and compliance", "Some client data may require restricted handling or alternate processing paths.", "Gate the approach through legal, privacy, and architecture review before scale."),
        ("Human oversight", "Teams must remain accountable for decisions and fixes applied in production.", "Keep the system advisory and auditable rather than autonomous."),
        ("Adoption and behavior change", "The platform will fail if engineers see it as extra work or a threat.", "Co-design with support teams and make the workflow genuinely easier to use."),
    ]
    for ridx, row in enumerate(rows):
        y = 2.32 + (ridx * 1.0)
        x = 0.88
        fills = [WHITE, WHITE, HCL_BLUE_SOFT if ridx % 2 == 0 else SOFT_STEEL]
        for cidx, value in enumerate(row):
            add_panel(slide, Inches(x), Inches(y), Inches(widths[cidx]), Inches(0.96), fills[cidx], LINE, rounded=False)
            add_text(slide, Inches(x + 0.12), Inches(y + 0.18), Inches(widths[cidx] - 0.2), Inches(0.52), value, size=10, color=INK if cidx != 1 else MUTED, bold=(cidx == 0))
            x += widths[cidx]

    add_footer(slide)


def slide_pilot(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "Pilot Recommendation", "Start with a contained support workflow, prove adoption and usefulness, then expand deliberately.")

    add_panel(slide, Inches(0.9), Inches(2.0), Inches(3.65), Inches(3.8), CARD, LINE)
    add_text(slide, Inches(1.15), Inches(2.3), Inches(2.8), Inches(0.24), "Recommended pilot shape", size=14, bold=True, color=HCL_BLUE)
    add_bullets(
        slide,
        Inches(1.08),
        Inches(2.72),
        Inches(3.0),
        Inches(1.7),
        [
            "Choose one support domain with stable ownership and enough repeatable case history.",
            "Keep the first release internal so the team can improve the workflow before exposing it to clients.",
            "Focus on a few high-value use cases rather than a broad platform rollout.",
        ],
        size=13,
    )

    add_panel(slide, Inches(4.9), Inches(2.0), Inches(3.35), Inches(3.8), CARD, LINE)
    add_text(slide, Inches(5.15), Inches(2.3), Inches(2.7), Inches(0.24), "What success should show", size=14, bold=True, color=HCL_BLUE_DARK)
    metrics = [
        "Faster turnaround on pilot cases",
        "Noticeably better productivity in day-to-day work",
        "More consistent issue handling quality",
        "Positive user adoption and confidence",
        "Clear signal that the workflow is worth scaling",
    ]
    for idx, metric in enumerate(metrics):
        add_panel(slide, Inches(5.14), Inches(2.75 + (idx * 0.46)), Inches(2.85), Inches(0.28), HCL_BLUE_SOFT if idx % 2 == 0 else SOFT_STEEL, rounded=True)
        add_text(slide, Inches(5.3), Inches(2.84 + (idx * 0.46)), Inches(2.5), Inches(0.1), metric, size=8, bold=True, color=HCL_BLUE_DARK if idx % 2 == 0 else INK)

    add_panel(slide, Inches(8.58), Inches(2.0), Inches(3.8), Inches(3.8), CARD, LINE)
    add_text(slide, Inches(8.82), Inches(2.3), Inches(3.0), Inches(0.24), "Pragmatic rollout path", size=14, bold=True, color=GREEN)
    steps = ["Validate readiness", "Run focused pilot", "Learn and tighten", "Decide on expansion"]
    for idx, step in enumerate(steps):
        y = 2.78 + idx * 0.72
        add_panel(slide, Inches(8.92), Inches(y), Inches(0.44), Inches(0.44), GREEN, GREEN)
        add_text(slide, Inches(9.02), Inches(y + 0.12), Inches(0.22), Inches(0.1), str(idx + 1), size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(9.52), Inches(y + 0.1), Inches(2.25), Inches(0.18), step, size=11, bold=True)
        if idx < len(steps) - 1:
            slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(9.14), Inches(y + 0.44), Inches(9.14), Inches(y + 0.69)).line.color.rgb = GREEN

    add_footer(slide)


def slide_ask(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "What We Need from Leadership", "Use this session to shape the proposal, sharpen the pilot, and decide whether it merits formal progression.")

    asks = [
        ("Business relevance", "Confirm whether this addresses a real delivery and support priority worth pursuing now."),
        ("Pilot sponsorship", "Help identify the right pilot area, accountable owner, and supporting functions."),
        ("Next-step alignment", "Agree whether to refine the case, validate readiness, and move toward a pilot business case."),
    ]
    for idx, (title, body) in enumerate(asks):
        left = 0.9 + idx * 4.05
        add_panel(slide, Inches(left), Inches(2.1), Inches(3.55), Inches(2.15), CARD, LINE)
        add_panel(slide, Inches(left), Inches(2.1), Inches(3.55), Inches(0.14), HCL_BLUE if idx == 0 else HCL_BLUE_DARK if idx == 1 else GREEN, rounded=False)
        add_text(slide, Inches(left + 0.2), Inches(2.45), Inches(2.9), Inches(0.24), title, size=14, bold=True)
        add_text(slide, Inches(left + 0.2), Inches(2.9), Inches(2.95), Inches(0.78), body, size=11, color=MUTED)

    add_panel(slide, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.96), HCL_BLUE, HCL_BLUE)
    add_text(slide, Inches(1.15), Inches(4.98), Inches(10.95), Inches(0.24), "Decision prompt: Is there enough business value and practical realism here to sponsor a focused pilot?", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_architecture(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "High-Level Architecture", "The solution can be understood as three simple layers working together in one governed support workflow.")

    layers = [
        (
            "Experience layer",
            "Engineer workspace, manager dashboard, and client-facing views assembled through Generative UI.",
            ["Role-based screens", "Dynamic workflow views", "Consistent user experience"],
            HCL_BLUE,
        ),
        (
            "Intelligence layer",
            "GenAI services orchestrate summarization, retrieval, reasoning, recommendation, and response generation with human oversight.",
            ["Context builder", "Knowledge retrieval", "Reasoning and guidance"],
            HCL_BLUE_DARK,
        ),
        (
            "Enterprise data layer",
            "Operational data is drawn from ticketing, knowledge, monitoring, change, and environment sources already used in support delivery.",
            ["Ticket data", "Change and monitoring signals", "Client and configuration context"],
            GREEN,
        ),
    ]

    for idx, (title, body, tags, accent) in enumerate(layers):
        top = 1.95 + (idx * 1.45)
        add_panel(slide, Inches(1.0), Inches(top), Inches(8.4), Inches(1.05), CARD, LINE)
        add_panel(slide, Inches(1.0), Inches(top), Inches(0.16), Inches(1.05), accent)
        add_text(slide, Inches(1.26), Inches(top + 0.18), Inches(2.2), Inches(0.2), title, size=14, bold=True)
        add_text(slide, Inches(3.1), Inches(top + 0.18), Inches(5.9), Inches(0.34), body, size=10, color=MUTED)
        for jdx, tag in enumerate(tags):
            add_panel(slide, Inches(3.1 + (jdx * 1.82)), Inches(top + 0.58), Inches(1.6), Inches(0.22), HCL_BLUE_SOFT if accent != GREEN else SOFT_STEEL, rounded=True)
            add_text(slide, Inches(3.2 + (jdx * 1.82)), Inches(top + 0.63), Inches(1.35), Inches(0.1), tag, size=8, bold=True, color=accent if accent != GREEN else GREEN)

    for idx in range(2):
        y1 = 3.0 + (idx * 1.45)
        y2 = 3.38 + (idx * 1.45)
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.15), Inches(y1), Inches(5.15), Inches(y2)).line.color.rgb = HCL_BLUE

    add_panel(slide, Inches(9.75), Inches(2.05), Inches(2.2), Inches(3.95), HCL_BLUE_SOFT, HCL_BLUE_SOFT)
    add_text(slide, Inches(10.0), Inches(2.35), Inches(1.7), Inches(0.3), "Governance across the stack", size=14, bold=True, color=HCL_BLUE, align=PP_ALIGN.CENTER)
    governance = [
        "Human-in-the-loop controls",
        "Confidence and approval gates",
        "Privacy and compliance guardrails",
        "Feedback loop to improve outputs over time",
    ]
    add_bullets(slide, Inches(9.95), Inches(2.95), Inches(1.55), Inches(2.05), governance, size=10, color=INK)

    add_footer(slide)


def slide_detailed_architecture(prs: Presentation):
    slide = new_slide(prs)
    add_header(slide, "Detailed Technical Architecture", "A module-level view of how the solution components work together behind the support experience.")

    columns = [
        (
            0.72,
            2.35,
            "Experience and workflow",
            [
                "Engineer copilot workspace",
                "Manager dashboard",
                "Client service views",
                "Generative UI renderer",
            ],
            HCL_BLUE,
        ),
        (
            3.38,
            2.35,
            "Orchestration and control",
            [
                "Ticket intake and classification",
                "Context builder",
                "Workflow orchestration",
                "Policy and routing layer",
            ],
            HCL_BLUE_DARK,
        ),
        (
            6.05,
            2.35,
            "Intelligence services",
            [
                "Knowledge retrieval service",
                "Prompt and reasoning service",
                "Recommendation engine",
                "Feedback and learning loop",
            ],
            GREEN,
        ),
        (
            8.72,
            2.35,
            "Enterprise data and integration",
            [
                "Ticketing and knowledge systems",
                "Monitoring and log sources",
                "Change and configuration sources",
                "Security and access controls",
            ],
            AMBER,
        ),
    ]

    for left, top, title, items, accent in columns:
        add_panel(slide, Inches(left), Inches(top), Inches(2.38), Inches(3.05), CARD, LINE)
        add_panel(slide, Inches(left), Inches(top), Inches(2.38), Inches(0.16), accent, accent, rounded=False)
        add_text(slide, Inches(left + 0.16), Inches(top + 0.28), Inches(2.0), Inches(0.42), title, size=12, bold=True)
        for idx, item in enumerate(items):
            add_panel(slide, Inches(left + 0.14), Inches(top + 0.82 + (idx * 0.5)), Inches(2.08), Inches(0.34), HCL_BLUE_SOFT if accent in (HCL_BLUE, HCL_BLUE_DARK) else SOFT_STEEL, rounded=True)
            add_text(slide, Inches(left + 0.25), Inches(top + 0.93 + (idx * 0.5)), Inches(1.82), Inches(0.1), item, size=8, bold=True, color=accent if accent != AMBER else INK)

    for x in (3.12, 5.78, 8.45):
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(3.85), Inches(x + 0.22), Inches(3.85)).line.color.rgb = HCL_BLUE

    add_panel(slide, Inches(0.9), Inches(5.82), Inches(10.1), Inches(0.72), WHITE, LINE)
    add_text(slide, Inches(1.16), Inches(6.05), Inches(9.6), Inches(0.14), "Flow: issue arrives -> orchestration assembles context -> intelligence services generate guidance -> Generative UI presents the right experience -> feedback improves future guidance.", size=10, bold=True, color=HCL_BLUE, align=PP_ALIGN.CENTER)

    add_panel(slide, Inches(11.2), Inches(2.35), Inches(1.3), Inches(4.2), HCL_BLUE_SOFT, HCL_BLUE_SOFT)
    add_text(slide, Inches(11.36), Inches(2.62), Inches(0.95), Inches(0.4), "Cross-cutting governance", size=12, bold=True, color=HCL_BLUE, align=PP_ALIGN.CENTER)
    add_bullets(
        slide,
        Inches(11.28),
        Inches(3.18),
        Inches(1.02),
        Inches(2.55),
        [
            "Identity and access",
            "Prompt and response guardrails",
            "Audit trail and observability",
            "Human approval points",
            "Privacy and compliance controls",
        ],
        size=8,
        color=INK,
    )

    add_footer(slide)


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover(prs)
    slide_name_expansion(prs)
    slide_current_challenge(prs)
    slide_why_now(prs)
    slide_opportunity(prs)
    slide_proposal(prs)
    slide_use_cases(prs)
    slide_business_value(prs)
    slide_risk(prs)
    slide_pilot(prs)
    slide_ask(prs)
    slide_architecture(prs)
    slide_detailed_architecture(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_OUTPUT)


def html_escape(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def bullet_list(items: list[str]) -> str:
    return "".join(f"<li>{html_escape(item)}</li>" for item in items)


def build_html(signals: dict[str, bool]) -> None:
    slides = [
        f"""
        <section class="slide cover">
          <div class="eyebrow">Leadership Discussion</div>
          <h1>Reimagining Support Services with GenAI Automation and Generative UI</h1>
          <p class="lede">A business-story proposal to improve support effectiveness, scale expertise, and make the service model easier to operate and experience.</p>
          <div class="hero-grid">
            <div class="hero-card"><h3>Why now</h3><p>Support demand, complexity, and customer expectations are moving faster than traditional support methods.</p></div>
            <div class="hero-card"><h3>What changes</h3><p>Use GenAI automation plus Generative UI to guide work, reduce rework, and surface the right context earlier.</p></div>
            <div class="hero-card"><h3>What we need</h3><p>Leadership input on business priority, pilot scope, sponsorship, and readiness to move forward.</p></div>
          </div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">What SIP Means</div>
          <h2>Support Intelligence Platform is the working name for the proposed support modernization layer.</h2>
          <div class="banner" style="text-align:center; font-size:1.5rem; font-weight:700;">SIP = Support Intelligence Platform</div>
          <div class="three-col">
            <div class="card">
              <h3>Support</h3>
              <p>The focus is the day-to-day support model: ticket handling, investigation, escalation, and service experience.</p>
            </div>
            <div class="card">
              <h3>Intelligence</h3>
              <p>The differentiator is the ability to turn past tickets, context, and expertise into useful guidance at the moment of need.</p>
            </div>
            <div class="card">
              <h3>Platform</h3>
              <p>The proposal is not a single feature. It is a reusable layer that can support engineers, managers, and clients over time.</p>
            </div>
          </div>
          <div class="source-note" style="margin-top:22px; text-align:center;">Plain-language takeaway: SIP is a smarter support operating layer that helps teams use what they already know, more consistently and at scale.</div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">The Current Challenge</div>
          <h2>Support works today, but it is harder to scale, standardize, and sustain.</h2>
          <div class="two-col">
            <div class="card">
              <h3>What leaders are seeing</h3>
              <ul>{bullet_list([
                  "Manual effort remains high during triage, context gathering, and handoffs.",
                  "Information is spread across tickets, engineer memory, and disconnected tools.",
                  "Resolution quality varies because the starting context depends on who picks up the issue.",
                  "Senior expertise solves problems quickly, but that expertise is difficult to reuse at scale.",
              ])}</ul>
            </div>
            <div class="card">
              <h3>Where the friction shows up</h3>
              <div class="flow">
                <div><strong>Ticket lands</strong><span>The issue arrives with limited context.</span></div>
                <div><strong>Engineer searches</strong><span>Time is spent hunting through history and tools.</span></div>
                <div><strong>Escalation happens</strong><span>The next tier often repeats discovery.</span></div>
                <div><strong>Knowledge fades</strong><span>The fix is stored, but the reasoning is not easy to reuse.</span></div>
              </div>
            </div>
          </div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">Why This Matters Now</div>
          <h2>This is now a business issue, not just an operations issue.</h2>
          <div class="split-grid">
            <div class="mini-card"><h3>Cost pressure</h3><p>Leaders need productivity gains without defaulting to linear headcount growth.</p></div>
            <div class="mini-card"><h3>Speed pressure</h3><p>Customers and account teams expect faster answers and clearer status.</p></div>
            <div class="mini-card"><h3>Quality pressure</h3><p>Consistency is harder when knowledge is fragmented and work is handed across tiers.</p></div>
            <div class="mini-card"><h3>Market pressure</h3><p>AI-enabled service models are changing what clients will expect from support partners.</p></div>
          </div>
          <div class="statement"><strong>If we continue as-is</strong><p>Support remains reactive and heavily person-dependent, which creates operational drag and slows responsiveness.</p></div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">The Opportunity</div>
          <h2>Move from reactive support to guided, assisted, and progressively more intelligent support.</h2>
          <div class="compare">
            <div class="card danger"><h3>Today</h3><ul>{bullet_list(["The engineer gathers context manually.", "Knowledge retrieval depends on who knows where to look.", "Escalations often restart the investigation.", "The service experience feels fragmented to engineers and clients."])}</ul></div>
            <div class="shift-arrow">Shift</div>
            <div class="card success">
              <h3>Future state</h3>
              <div class="future-grid">
                <div><strong>Guided start</strong><span>Relevant history and next actions appear earlier.</span></div>
                <div><strong>Smarter workbench</strong><span>The interface adapts to the support tier and issue state.</span></div>
                <div><strong>Reusable expertise</strong><span>Reasoning paths become part of the service knowledge base.</span></div>
                <div><strong>Better visibility</strong><span>Managers and clients get a clearer support experience.</span></div>
              </div>
            </div>
          </div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">What We Are Proposing</div>
          <h2>Combine GenAI automation and Generative UI to augment support operations, not replace them.</h2>
          <div class="proposal">
            <div class="card"><h3>GenAI automation</h3><ul>{bullet_list(["Summarizes issues and extracts clean context from messy tickets.", "Retrieves relevant historical cases and knowledge quickly.", "Recommends actions, hypotheses, and escalation guidance."])}</ul></div>
            <div class="plus">+</div>
            <div class="card"><h3>Generative UI</h3><ul>{bullet_list(["Builds role-based experiences for engineers, managers, and clients.", "Adapts the workspace to the issue context and workflow stage.", "Presents AI guidance in a practical, auditable operating surface."])}</ul></div>
            <div class="card accent"><h3>Targeted augmentation</h3><p>The value comes from improving the work itself: better context, better decisions, and smoother execution.</p></div>
          </div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">How It Would Work in Practice</div>
          <h2>The concept is tangible because it maps directly to everyday support work.</h2>
          <div class="three-col">
            <div class="card top-blue"><h3>Agent assist at ticket open</h3><p>Organizes the issue, surfaces likely matches, and suggests a clean first response path.</p><div class="pill-row"><span>Faster triage</span><span>Less manual searching</span><span>Cleaner escalation context</span></div></div>
            <div class="card top-dark"><h3>Knowledge guidance during investigation</h3><p>Brings in environment notes, recent changes, and next-best actions in the flow of work.</p><div class="pill-row"><span>Better investigation quality</span><span>More reusable reasoning</span><span>Lower rework</span></div></div>
            <div class="card top-green"><h3>Dynamic UI by role and workflow stage</h3><p>The workspace changes based on who is using it and where the issue is in the lifecycle.</p><div class="pill-row"><span>Less clutter</span><span>Higher adoption potential</span><span>More consistent experience</span></div></div>
          </div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">Business Value</div>
          <h2>The value is operational, financial, and experiential.</h2>
          <div class="banner">Expected outcome: support becomes easier to run, easier to improve, and easier for customers to experience.</div>
          <div class="three-col">
            <div class="card"><h3>Operational value</h3><ul>{bullet_list(["Faster resolution flow", "Better context at handoff", "More consistent service execution"])}</ul></div>
            <div class="card"><h3>Financial value</h3><ul>{bullet_list(["Improved productivity leverage", "Reduced repeated effort", "Stronger scale potential without defaulting to headcount"])}</ul></div>
            <div class="card"><h3>Experience value</h3><ul>{bullet_list(["Better agent confidence", "Cleaner client communication", "Greater visibility into service health"])}</ul></div>
          </div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">Risk, Governance, and Adoption</div>
          <h2>The proposal should be governed, human-in-the-loop, phased, and measurable.</h2>
          <table>
            <thead><tr><th>Leadership concern</th><th>What it means</th><th>Control approach</th></tr></thead>
            <tbody>
              <tr><td>Accuracy and hallucination risk</td><td>AI can produce plausible guidance that still needs judgment.</td><td>Use confidence signals, human confirmation, and feedback loops.</td></tr>
              <tr><td>Security and compliance</td><td>Some client data may require restricted handling or alternate processing paths.</td><td>Gate the approach through legal, privacy, and architecture review before scale.</td></tr>
              <tr><td>Human oversight</td><td>Teams must remain accountable for decisions and fixes applied in production.</td><td>Keep the system advisory and auditable rather than autonomous.</td></tr>
              <tr><td>Adoption and behavior change</td><td>The platform will fail if engineers see it as extra work or a threat.</td><td>Co-design with support teams and make the workflow genuinely easier to use.</td></tr>
            </tbody>
          </table>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">Pilot Recommendation</div>
          <h2>Start with a contained support workflow, prove adoption and usefulness, then expand deliberately.</h2>
          <div class="three-col">
            <div class="card"><h3>Recommended pilot shape</h3><ul>{bullet_list(["Choose one support domain with stable ownership and enough repeatable case history.", "Keep the first release internal so the team can improve the workflow before exposing it to clients.", "Focus on a few high-value use cases rather than a broad platform rollout."])}</ul></div>
            <div class="card"><h3>What success should show</h3><div class="pill-stack"><span>Faster turnaround on pilot cases</span><span>Better day-to-day productivity</span><span>More consistent handling quality</span><span>Positive user adoption</span><span>Clear signal to scale</span></div></div>
            <div class="card"><h3>Pragmatic rollout path</h3><ol class="steps"><li>Validate readiness</li><li>Run focused pilot</li><li>Learn and tighten</li><li>Decide on expansion</li></ol></div>
          </div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">What We Need from Leadership</div>
          <h2>Use this session to shape the proposal, sharpen the pilot, and decide whether it merits formal progression.</h2>
          <div class="three-col">
            <div class="card"><h3>Business relevance</h3><p>Confirm whether this addresses a real delivery and support priority worth pursuing now.</p></div>
            <div class="card"><h3>Pilot sponsorship</h3><p>Help identify the right pilot area, accountable owner, and supporting functions.</p></div>
            <div class="card"><h3>Next-step alignment</h3><p>Agree whether to refine the case, validate readiness, and move toward a pilot business case.</p></div>
          </div>
          <div class="decision-box">Decision prompt: Is there enough business value and practical realism here to sponsor a focused pilot?</div>
          <div class="source-note">Source alignment check: knowledge decay={str(signals["knowledge_decay"]).lower()}, phase-0 gating={str(signals["phase_zero"]).lower()}, human gate={str(signals["human_gate"]).lower()}, co-design={str(signals["co_design"]).lower()}.</div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">High-Level Architecture</div>
          <h2>The solution can be understood as three simple layers working together in one governed support workflow.</h2>
          <div class="two-col" style="grid-template-columns: 2.1fr 1fr;">
            <div>
              <div class="card top-blue">
                <h3>Experience layer</h3>
                <p>Engineer workspace, manager dashboard, and client-facing views assembled through Generative UI.</p>
                <div class="pill-row"><span>Role-based screens</span><span>Dynamic workflow views</span><span>Consistent user experience</span></div>
              </div>
              <div class="card top-dark" style="margin-top:16px;">
                <h3>Intelligence layer</h3>
                <p>GenAI services orchestrate summarization, retrieval, reasoning, recommendation, and response generation with human oversight.</p>
                <div class="pill-row"><span>Context builder</span><span>Knowledge retrieval</span><span>Reasoning and guidance</span></div>
              </div>
              <div class="card top-green" style="margin-top:16px;">
                <h3>Enterprise data layer</h3>
                <p>Operational data is drawn from ticketing, knowledge, monitoring, change, and environment sources already used in support delivery.</p>
                <div class="pill-row"><span>Ticket data</span><span>Change and monitoring signals</span><span>Client and configuration context</span></div>
              </div>
            </div>
            <div class="card accent">
              <h3>Governance across the stack</h3>
              <ul>{bullet_list(["Human-in-the-loop controls", "Confidence and approval gates", "Privacy and compliance guardrails", "Feedback loop to improve outputs over time"])}</ul>
            </div>
          </div>
        </section>
        """,
        f"""
        <section class="slide">
          <div class="kicker">Detailed Technical Architecture</div>
          <h2>A module-level view of how the solution components work together behind the support experience.</h2>
          <div style="display:grid; grid-template-columns:repeat(4, 1fr) 0.72fr; gap:14px; margin-top:26px;">
            <div class="card top-blue">
              <h3>Experience and workflow</h3>
              <div class="pill-stack"><span>Engineer copilot workspace</span><span>Manager dashboard</span><span>Client service views</span><span>Generative UI renderer</span></div>
            </div>
            <div class="card top-dark">
              <h3>Orchestration and control</h3>
              <div class="pill-stack"><span>Ticket intake and classification</span><span>Context builder</span><span>Workflow orchestration</span><span>Policy and routing layer</span></div>
            </div>
            <div class="card top-green">
              <h3>Intelligence services</h3>
              <div class="pill-stack"><span>Knowledge retrieval service</span><span>Prompt and reasoning service</span><span>Recommendation engine</span><span>Feedback and learning loop</span></div>
            </div>
            <div class="card" style="border-top:6px solid #da9100;">
              <h3>Enterprise data and integration</h3>
              <div class="pill-stack"><span>Ticketing and knowledge systems</span><span>Monitoring and log sources</span><span>Change and configuration sources</span><span>Security and access controls</span></div>
            </div>
            <div class="card accent">
              <h3>Governance</h3>
              <ul>{bullet_list(["Identity and access", "Guardrails", "Audit trail", "Human approval", "Privacy and compliance"])}</ul>
            </div>
          </div>
          <div class="banner">Flow: issue arrives -> orchestration assembles context -> intelligence services generate guidance -> Generative UI presents the right experience -> feedback improves future guidance.</div>
        </section>
        """,
    ]

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>SIP Executive Leadership Deck</title>
  <style>
    :root {{--hcl-blue:#006bb6;--hcl-blue-dark:#004c81;--hcl-blue-soft:#e2eff9;--ink:#0d0d0d;--muted:#5c5c5c;--bg:#f5f7fa;--card:#ffffff;--line:#d4dce4;--green:#268455;--shadow:0 18px 40px rgba(13,13,13,0.08);}}
    * {{ box-sizing: border-box; }}
    body {{ margin:0; font-family:"Aptos","Segoe UI",sans-serif; background:radial-gradient(circle at top left, rgba(0,107,182,0.14), transparent 28%), linear-gradient(180deg, #eef4f8 0%, #f7f9fb 100%); color:var(--ink); }}
    .deck {{ width:min(1200px, calc(100vw - 48px)); margin:28px auto 48px; display:grid; gap:28px; }}
    .slide {{ position:relative; min-height:675px; background:var(--card); border:1px solid var(--line); border-radius:28px; box-shadow:var(--shadow); padding:44px 48px 38px; overflow:hidden; }}
    .slide::before {{ content:""; position:absolute; inset:0 0 auto 0; height:10px; background:linear-gradient(90deg, var(--hcl-blue), #2b86c5 42%, #5ca9d8 100%); }}
    .cover {{ background:linear-gradient(135deg, rgba(226,239,249,0.9), rgba(255,255,255,0.96)), var(--card); }}
    .eyebrow,.kicker {{ color:var(--hcl-blue); text-transform:uppercase; letter-spacing:0.12em; font-size:0.78rem; font-weight:700; }}
    h1,h2 {{ margin:16px 0 10px; font-family:"Aptos Display","Aptos","Segoe UI",sans-serif; line-height:1.08; }}
    h1 {{ font-size:2.7rem; max-width:820px; }}
    h2 {{ font-size:2rem; max-width:900px; }}
    h3 {{ margin:0 0 10px; font-size:1.05rem; }}
    p,li,td,th {{ font-size:1rem; line-height:1.45; }}
    .lede {{ max-width:760px; color:var(--muted); font-size:1.15rem; }}
    .hero-grid,.three-col,.split-grid,.two-col,.compare,.proposal,.future-grid {{ display:grid; gap:18px; }}
    .hero-grid {{ grid-template-columns:repeat(3,1fr); margin-top:64px; }}
    .three-col {{ grid-template-columns:repeat(3,1fr); margin-top:30px; }}
    .split-grid {{ grid-template-columns:repeat(2,1fr); margin-top:26px; }}
    .two-col {{ grid-template-columns:1fr 1fr; margin-top:28px; }}
    .compare {{ grid-template-columns:1fr 120px 1.25fr; margin-top:28px; align-items:start; }}
    .proposal {{ grid-template-columns:1fr 90px 1fr 0.72fr; margin-top:28px; align-items:start; }}
    .card,.hero-card,.mini-card {{ background:var(--card); border:1px solid var(--line); border-radius:22px; padding:22px 24px; box-shadow:0 8px 18px rgba(13,13,13,0.04); }}
    .hero-card {{ min-height:180px; border-top:6px solid var(--hcl-blue); }}
    .mini-card {{ border-left:6px solid var(--hcl-blue); min-height:138px; }}
    .danger {{ border-top:6px solid #c2392b; }}
    .success {{ border-top:6px solid #268455; }}
    .accent {{ background:linear-gradient(180deg, var(--hcl-blue-soft), #f6fbff); }}
    .flow {{ display:grid; gap:12px; margin-top:14px; }}
    .flow div,.future-grid div {{ background:linear-gradient(180deg, #fbfdff, #f4f8fb); border:1px solid var(--line); border-radius:16px; padding:14px 16px; }}
    .flow span,.future-grid span,.source-note {{ display:block; color:var(--muted); font-size:0.94rem; margin-top:6px; }}
    .statement,.banner,.decision-box {{ margin-top:28px; background:linear-gradient(135deg, var(--hcl-blue), #2b86c5); color:#fff; border-radius:22px; padding:20px 24px; }}
    .banner {{ background:linear-gradient(135deg, #eef6fc, #e2eff9); color:var(--ink); border:1px solid var(--line); }}
    .decision-box {{ text-align:center; font-size:1.2rem; font-weight:700; }}
    .shift-arrow,.plus {{ align-self:center; justify-self:center; width:84px; height:84px; display:grid; place-items:center; border-radius:50%; background:var(--hcl-blue); color:#fff; font-weight:700; box-shadow:0 10px 22px rgba(0,107,182,0.22); }}
    .plus {{ font-size:2rem; }}
    .future-grid {{ grid-template-columns:repeat(2,1fr); margin-top:16px; }}
    .pill-row,.pill-stack {{ display:flex; flex-wrap:wrap; gap:8px; margin-top:16px; }}
    .pill-row span,.pill-stack span {{ display:inline-block; padding:8px 12px; border-radius:999px; background:var(--hcl-blue-soft); color:var(--hcl-blue-dark); font-size:0.82rem; font-weight:700; }}
    .pill-stack {{ flex-direction:column; align-items:stretch; }}
    .top-blue {{ border-top:6px solid var(--hcl-blue); }}
    .top-dark {{ border-top:6px solid var(--hcl-blue-dark); }}
    .top-green {{ border-top:6px solid #268455; }}
    table {{ width:100%; margin-top:28px; border-collapse:collapse; border:1px solid var(--line); }}
    thead th {{ background:var(--hcl-blue); color:#fff; text-align:left; padding:14px 16px; font-size:0.92rem; }}
    tbody td {{ padding:16px; vertical-align:top; border-top:1px solid var(--line); }}
    tbody tr:nth-child(even) td:last-child {{ background:#f3f8fc; }}
    .steps {{ margin:14px 0 0 18px; display:grid; gap:14px; }}
    .source-note {{ margin-top:18px; text-align:right; font-size:0.82rem; }}
    @media print {{ body {{ background:#fff; }} .deck {{ width:100%; margin:0; gap:0; }} .slide {{ min-height:7.5in; border-radius:0; box-shadow:none; page-break-after:always; margin:0; }} }}
  </style>
</head>
<body><main class="deck">{''.join(slides)}</main></body>
</html>"""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    HTML_OUTPUT.write_text(html, encoding="utf-8")


def main() -> None:
    ensure_source_doc()
    paragraphs = extract_doc_paragraphs(SOURCE_DOC)
    signals = source_signals(paragraphs)
    build_ppt()
    build_html(signals)
    print(f"Created {PPT_OUTPUT}")
    print(f"Created {HTML_OUTPUT}")


if __name__ == "__main__":
    main()
