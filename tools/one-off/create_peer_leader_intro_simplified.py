from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(r"C:\Users\kumar.gn\HCLProjects\BusinessProcess\docs\output\peer-leader-introduction-kumargn-simplified-v2.pptx")

BG = RGBColor(244, 247, 251)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
ACCENT = RGBColor(9, 58, 92)
ACCENT_SOFT = RGBColor(224, 236, 245)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(203, 213, 225)
HIGHLIGHT = RGBColor(180, 83, 9)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, title, items):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE

    add_text(slide, left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.35), title, 18, ACCENT, True)

    box = slide.shapes.add_textbox(left + Inches(0.28), top + Inches(0.62), width - Inches(0.56), height - Inches(0.82))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.bullet = True
        paragraph.level = 0
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = TEXT


def add_badge(slide, left, top, width, height, text):
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_SOFT
    badge.line.fill.background()
    add_text(slide, left, top + Inches(0.05), width, height, text, 13, ACCENT, True, PP_ALIGN.CENTER)


def build_presentation():
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)

    add_badge(slide, Inches(0.75), Inches(0.42), Inches(1.7), Inches(0.34), "PEER LEADER")
    add_text(slide, Inches(0.75), Inches(0.88), Inches(6.2), Inches(0.55), "Kumar G N", 30, TEXT, True)
    add_text(slide, Inches(0.75), Inches(1.35), Inches(7.4), Inches(0.62), "Engineering leader with deep product roots and a current focus on scaling enterprise GenAI responsibly.", 16, MUTED)
    add_text(slide, Inches(9.0), Inches(0.70), Inches(3.4), Inches(0.28), "linkedin.com/in/kumargn", 11, ACCENT, False, PP_ALIGN.RIGHT)
    add_text(slide, Inches(9.0), Inches(1.00), Inches(3.4), Inches(0.28), "medium.com/@ganduri.kumar", 11, ACCENT, False, PP_ALIGN.RIGHT)

    add_text(slide, Inches(0.75), Inches(2.0), Inches(11.6), Inches(0.36), "From embedded systems to enterprise platforms, the constant has been building teams, shaping transformation, and turning complex goals into delivery momentum.", 17, TEXT)

    add_bullets(
        slide,
        Inches(0.75),
        Inches(2.55),
        Inches(3.85),
        Inches(3.25),
        "Leadership Snapshot",
        [
            "35+ years across software engineering and product development.",
            "Experience spans embedded systems, networking, and enterprise applications.",
            "13 years with HCLTech ERS, leading modernization and delivery programs.",
        ],
    )

    add_bullets(
        slide,
        Inches(4.74),
        Inches(2.55),
        Inches(3.85),
        Inches(3.25),
        "What I Focus On",
        [
            "Driving practical GenAI adoption across the SDLC.",
            "Improving context, accuracy, consistency, and traceability in AI systems.",
            "Shaping governance patterns needed to scale enterprise AI with confidence.",
        ],
    )

    add_bullets(
        slide,
        Inches(8.73),
        Inches(2.55),
        Inches(3.85),
        Inches(3.25),
        "What Keeps Me Grounded",
        [
            "Tennis is both my reset button and competitive outlet.",
            "Family keeps life energetic, balanced, and honest.",
            "I value thoughtful collaboration, curiosity, and steady progress.",
        ],
    )

    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(6.15),
        Inches(11.83),
        Inches(0.52),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = ACCENT
    footer.line.fill.background()
    add_text(
        slide,
        Inches(1.0),
        Inches(6.27),
        Inches(11.3),
        Inches(0.22),
        "Looking forward to learning from this community and contributing where experience can help.",
        14,
        RGBColor(255, 255, 255),
        False,
        PP_ALIGN.CENTER,
    )

    presentation.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
