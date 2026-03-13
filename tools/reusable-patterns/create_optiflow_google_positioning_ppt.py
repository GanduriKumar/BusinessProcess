from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(
    r"C:\Users\kumar.gn\HCLProjects\BusinessProcess\docs\output\OptiFlow - Google Compatible Positioning Mar 2026.pptx"
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

GOOGLE_BLUE = RGBColor(66, 133, 244)
GOOGLE_RED = RGBColor(234, 67, 53)
GOOGLE_YELLOW = RGBColor(251, 188, 5)
GOOGLE_GREEN = RGBColor(52, 168, 83)
TEXT = RGBColor(32, 33, 36)
MUTED = RGBColor(95, 99, 104)
LINE = RGBColor(232, 234, 237)
CARD = RGBColor(255, 255, 255)
BG = RGBColor(248, 249, 250)
WARN = RGBColor(245, 124, 0)
BAD = RGBColor(198, 40, 40)
GOOD = RGBColor(46, 125, 50)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOOGLE_BLUE
    bar.line.fill.background()

    x = 0
    for color in (GOOGLE_BLUE, GOOGLE_RED, GOOGLE_YELLOW, GOOGLE_GREEN):
        seg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), 0, Inches(3.33325), Inches(0.18))
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.line.fill.background()
        x += 3.33325


def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.22))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = FONT_BODY
    run.font.size = Pt(8)
    run.font.color.rgb = MUTED


def add_text(slide, left, top, width, height, text, size=18, color=TEXT, bold=False, font=FONT_BODY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, size=18, color=TEXT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.bullet = True
        p.level = 0
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.name = FONT_BODY
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return tb


def add_title(slide, kicker, title, subtitle=None):
    add_top_band(slide)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(12.0), Inches(0.25), kicker, size=11, color=GOOGLE_BLUE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.82), Inches(11.8), Inches(0.95), title, size=28, bold=True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.75), Inches(11.4), Inches(0.55), subtitle, size=13, color=MUTED)


def add_card(slide, left, top, width, height, title, body, accent=GOOGLE_BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.12), height)
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()

    add_text(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.32), Inches(0.3), title, size=14, bold=True)
    add_text(slide, left + Inches(0.22), top + Inches(0.54), width - Inches(0.32), height - Inches(0.68), body, size=13, color=MUTED)


def add_status_table(slide, rows):
    left = Inches(0.65)
    top = Inches(1.8)
    widths = [Inches(1.05), Inches(3.0), Inches(1.15), Inches(6.7)]
    headers = ["Orig.", "Draft topic", "Fit", "Why / rewrite direction"]
    row_h = Inches(0.42)

    for idx, header in enumerate(headers):
        x = left + sum(widths[:idx], Inches(0))
        cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, top, widths[idx], row_h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GOOGLE_BLUE
        cell.line.color.rgb = BG
        add_text(slide, x + Inches(0.08), top + Inches(0.07), widths[idx] - Inches(0.16), Inches(0.2), header, size=11, color=CARD, bold=True)

    for r_idx, row in enumerate(rows, start=1):
        y = top + row_h * r_idx
        values = [row["slide"], row["topic"], row["fit"], row["note"]]
        for c_idx, value in enumerate(values):
            x = left + sum(widths[:c_idx], Inches(0))
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, widths[c_idx], row_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD
            cell.line.color.rgb = LINE
            color = TEXT
            if c_idx == 2:
                color = {"Red": BAD, "Yellow": WARN, "Green": GOOD}[value]
            add_text(slide, x + Inches(0.08), y + Inches(0.06), widths[c_idx] - Inches(0.14), row_h - Inches(0.08), value, size=10, color=color, bold=(c_idx == 2))


def add_provider_matrix(slide, rows):
    left = Inches(0.45)
    top = Inches(1.7)
    widths = [Inches(1.9), Inches(1.0), Inches(2.25), Inches(3.35), Inches(4.9)]
    headers = ["Provider", "Fit", "Native motion", "Where OptiFlow fits", "Changes needed in OptiFlow"]
    row_h = Inches(0.58)

    for idx, header in enumerate(headers):
        x = left + sum(widths[:idx], Inches(0))
        cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, widths[idx], row_h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GOOGLE_BLUE
        cell.line.color.rgb = BG
        add_text(slide, x + Inches(0.07), top + Inches(0.12), widths[idx] - Inches(0.14), Inches(0.22), header, size=10, color=CARD, bold=True)

    color_map = {"Green": GOOD, "Yellow": WARN, "Red": BAD}
    fill_map = {
        "Green": RGBColor(232, 245, 233),
        "Yellow": RGBColor(255, 243, 224),
        "Red": RGBColor(252, 235, 233),
    }

    for r_idx, row in enumerate(rows, start=1):
        y = top + row_h * r_idx
        values = [row["provider"], row["fit"], row["motion"], row["fitment"], row["changes"]]
        for c_idx, value in enumerate(values):
            x = left + sum(widths[:c_idx], Inches(0))
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, widths[c_idx], row_h)
            cell.fill.solid()
            if c_idx == 1:
                cell.fill.fore_color.rgb = fill_map[row["fit"]]
            else:
                cell.fill.fore_color.rgb = CARD
            cell.line.color.rgb = LINE
            color = color_map[row["fit"]] if c_idx == 1 else TEXT
            add_text(slide, x + Inches(0.07), y + Inches(0.1), widths[c_idx] - Inches(0.14), Inches(0.32), value, size=9, color=color, bold=(c_idx in (0, 1)))


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    layout = prs.slide_layouts[6]

    sources = (
        "Sources: Google Search Central AI features & your website; Google AI search guidance (May 21, 2025); "
        "Google Ads AI Overviews guidance; Google Search / Merchant Center ecommerce docs."
    )

    slide = prs.slides.add_slide(layout)
    set_bg(slide, RGBColor(255, 255, 255))
    add_top_band(slide)
    add_text(slide, Inches(0.75), Inches(0.7), Inches(12.0), Inches(0.35), "OptiFlow", size=15, color=GOOGLE_BLUE, bold=True)
    add_text(
        slide,
        Inches(0.75),
        Inches(1.2),
        Inches(11.3),
        Inches(1.1),
        "Repositioning OptiFlow for Google Search,\nAds, and Merchant Ecosystems",
        size=28,
        bold=True,
        font=FONT_HEAD,
    )
    add_text(
        slide,
        Inches(0.75),
        Inches(2.45),
        Inches(10.8),
        Inches(0.75),
        "Executive repositioning of OptiFlow for Google's ads, merchant, and AI-search ecosystem.",
        size=15,
        color=MUTED,
    )
    add_card(slide, Inches(0.78), Inches(4.25), Inches(3.8), Inches(1.35), "Bottom-line verdict", "Do not position OptiFlow as a way to pay for organic AI visibility. Position it as a paid-media, merchant, and landing-page quality accelerator.", GOOGLE_BLUE)
    add_card(slide, Inches(4.75), Inches(4.25), Inches(3.65), Inches(1.35), "Commercial fit", "Strongest fit is inside Google Ads, Merchant Center, Shopping, and AI-surface readiness workflows for paying advertisers and merchants.", GOOGLE_GREEN)
    add_card(slide, Inches(8.58), Inches(4.25), Inches(3.95), Inches(1.35), "Primary risk", "Current draft blurs ads and organic ranking, overstates outcome certainty, and relies on constructs that Google does not recognize as ranking controls.", GOOGLE_RED)
    add_footer(slide, "Prepared from official Google materials current as of March 12, 2026")

    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_title(slide, "Assessment", "Executive verdict", "OptiFlow is relevant to Google, but only after a clear repositioning.")
    add_card(slide, Inches(0.65), Inches(2.0), Inches(3.8), Inches(2.2), "What fits Google", "Improve landing pages, product data, structured markup, and merchant quality for paying customers.", GOOGLE_GREEN)
    add_card(slide, Inches(4.75), Inches(2.0), Inches(3.8), Inches(2.2), "What does not fit", "Any message that implies customers can pay for organic ranking, citation, or answer priority.", GOOGLE_RED)
    add_card(slide, Inches(8.85), Inches(2.0), Inches(3.8), Inches(2.2), "Recommended framing", "Optimization assistant for advertiser readiness, merchant integrity, and AI-surface eligibility.", GOOGLE_BLUE)
    add_bullets(
        slide,
        Inches(0.75),
        Inches(4.7),
        Inches(11.8),
        Inches(1.6),
        [
            "Google maintains a clear boundary between sponsored placements and organic inclusion.",
            "Organic AI appearance still follows standard Search quality, access, and snippet rules.",
            "The strongest fit is advertiser and merchant performance improvement, not ranking control.",
        ],
        size=17,
    )
    add_footer(slide, sources)

    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_title(slide, "Constraint", "Why the current draft misses Google's model")
    add_card(slide, Inches(0.7), Inches(1.95), Inches(3.8), Inches(3.55), "1. It blurs ads and organic", "The current draft suggests brands can be 'prioritized' in AI answers. Google does not sell organic ranking that way.", GOOGLE_RED)
    add_card(slide, Inches(4.78), Inches(1.95), Inches(3.8), Inches(3.55), "2. It leans on weak signals", "llms.txt is overstated, and GEO is treated as a Google-recognized layer. Public guidance does not support that.", GOOGLE_YELLOW)
    add_card(slide, Inches(8.86), Inches(1.95), Inches(3.8), Inches(3.55), "3. It overclaims impact", "The headline uplift numbers are too strong for a Google-facing story without product-grade evidence.", GOOGLE_RED)
    add_footer(slide, sources)

    slide = prs.slides.add_slide(layout)
    set_bg(slide, RGBColor(255, 255, 255))
    add_title(slide, "Official model", "What Google actually supports today")
    add_card(slide, Inches(0.7), Inches(1.9), Inches(3.65), Inches(1.7), "Organic AI visibility", "No separate GEO program. AI Overviews / AI Mode rely on the same core Search requirements: indexability, quality, access, snippets, and useful content.", GOOGLE_BLUE)
    add_card(slide, Inches(4.55), Inches(1.9), Inches(3.65), Inches(1.7), "Paid AI visibility", "Ads can appear in AI Overviews and AI Mode, but via ad products and auction mechanisms. Payment buys sponsored participation, not organic inclusion.", GOOGLE_GREEN)
    add_card(slide, Inches(8.4), Inches(1.9), Inches(4.2), Inches(1.7), "Commerce readiness", "Google explicitly asks merchants for accurate feeds, consistent landing pages, supported schema, crawlable product pages, and up-to-date pricing and availability.", GOOGLE_YELLOW)
    add_bullets(
        slide,
        Inches(0.85),
        Inches(4.15),
        Inches(11.5),
        Inches(2.0),
        [
            "The opportunity is real, but it sits inside advertiser and merchant enablement.",
            "Google-safe positioning improves content and data quality without selling organic advantage.",
            "OptiFlow should be described as a readiness layer for paid customers, not a shortcut to AI citations.",
        ],
        size=17,
    )
    add_footer(slide, sources)

    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_title(slide, "Rewritten thesis", "Google-compatible product story for OptiFlow")
    add_text(slide, Inches(0.75), Inches(1.85), Inches(12.0), Inches(0.45), "Recommended positioning", size=13, color=GOOGLE_BLUE, bold=True)
    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.25), Inches(11.8), Inches(1.25))
    quote.fill.solid()
    quote.fill.fore_color.rgb = CARD
    quote.line.color.rgb = LINE
    add_text(
        slide,
        Inches(1.0),
        Inches(2.6),
        Inches(11.2),
        Inches(0.7),
        "OptiFlow improves landing-page quality, structured data integrity, and merchant readiness across Google Search, Shopping, AI Overviews, and AI Mode.",
        size=22,
        bold=True,
        font=FONT_HEAD,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    add_card(slide, Inches(0.78), Inches(4.1), Inches(3.8), Inches(1.7), "What Google can sell", "Page and feed optimization, merchant diagnostics, AI-ready creative refinement, and content quality recommendations.", GOOGLE_GREEN)
    add_card(slide, Inches(4.78), Inches(4.1), Inches(3.8), Inches(1.7), "What Google should avoid", "Any promise of organic AI citation, answer priority, or ranking uplift tied to payment.", GOOGLE_RED)
    add_card(slide, Inches(8.78), Inches(4.1), Inches(3.8), Inches(1.7), "Why it works", "It improves advertiser outcomes and fits the ad-plus-commerce model Google already runs.", GOOGLE_BLUE)
    add_footer(slide, sources)

    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_title(slide, "Offer shape", "How OptiFlow should map into Google's customer workflows")
    add_card(slide, Inches(0.7), Inches(1.95), Inches(2.9), Inches(3.7), "Google Ads", "Refine campaign landing pages, tighten ad-to-page message match, improve structured page content, and surface content gaps that reduce conversion quality.", GOOGLE_BLUE)
    add_card(slide, Inches(3.85), Inches(1.95), Inches(2.9), Inches(3.7), "Merchant Center", "Detect feed/page mismatches, variant issues, missing product attributes, stale prices, weak descriptions, and unsupported imagery or disclosure gaps.", GOOGLE_GREEN)
    add_card(slide, Inches(7.0), Inches(1.95), Inches(2.9), Inches(3.7), "Shopping / AI Mode", "Prepare retailer content for richer machine readability so products are more eligible for shopping experiences and agentic commerce flows.", GOOGLE_YELLOW)
    add_card(slide, Inches(10.15), Inches(1.95), Inches(2.5), Inches(3.7), "Search readiness", "Suggest structural improvements that help publishers meet existing Search guidance without implying a separate Google GEO algorithm.", GOOGLE_RED)
    add_footer(slide, sources)

    slide = prs.slides.add_slide(layout)
    set_bg(slide, RGBColor(255, 255, 255))
    add_title(slide, "Packaging", "A Google-facing commercial model")
    add_bullets(
        slide,
        Inches(0.82),
        Inches(1.95),
        Inches(11.7),
        Inches(3.6),
        [
            "Position OptiFlow inside paid-media and merchant workflows: landing pages, feeds, assets, structured data, and diagnostics.",
            "Use language such as 'improve quality', 'increase eligibility', 'reduce mismatches', and 'strengthen machine readability'.",
            "Anchor value in measurable metrics: approval rates, feed health, page quality, conversion readiness, and post-click performance.",
            "Treat organic AI inclusion as a byproduct of better content and technical readiness, never as a paid entitlement.",
        ],
        size=18,
    )
    add_card(slide, Inches(0.85), Inches(5.65), Inches(3.65), Inches(0.9), "Good message", "Better landing pages and cleaner product data for AI-era Search and Shopping.", GOOGLE_GREEN)
    add_card(slide, Inches(4.8), Inches(5.65), Inches(3.0), Inches(0.9), "Avoid", "Pay to get cited in AI answers.", GOOGLE_RED)
    add_card(slide, Inches(8.1), Inches(5.65), Inches(4.45), Inches(0.9), "Enterprise angle", "Google enables advertisers; HCL can still sell consulting and implementation around the tool.", GOOGLE_BLUE)
    add_footer(slide, sources)

    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_title(slide, "Assessment", "Red / yellow / green review of the current draft")
    rows = [
        {"slide": "1", "topic": "Cover promise", "fit": "Yellow", "note": "Keep the problem statement. Soften 'prominently surfaced' to 'better represented and easier to understand'."},
        {"slide": "2", "topic": "GEO overview", "fit": "Red", "note": "Rebuild this slide around Search guidance, schema, crawlability, and content clarity. Remove llms.txt."},
        {"slide": "3", "topic": "Benefits", "fit": "Red", "note": "Replace hard uplift claims with directional value tied to advertiser, merchant, and page-quality outcomes."},
        {"slide": "6", "topic": "Opportunity", "fit": "Yellow", "note": "Keep the market shift. Remove any hint that brands can buy AI answer priority."},
        {"slide": "8", "topic": "Why now", "fit": "Green", "note": "Strong section. Keep it, but frame the benefit as readiness and resilience."},
        {"slide": "10", "topic": "How it works", "fit": "Yellow", "note": "Useful engine. Rename impression metrics to quality and readiness diagnostics."},
        {"slide": "14", "topic": "Business impact", "fit": "Yellow", "note": "Use lighter commercial language and separate paid outcomes from organic influence."},
        {"slide": "16", "topic": "Deployment", "fit": "Green", "note": "Good phased model. Report merchant, landing-page, and performance health metrics."},
        {"slide": "22", "topic": "Customer message", "fit": "Red", "note": "Replace 'recommended at the moment of intent' with 'more eligible, consistent, and accurately represented'."},
    ]
    add_status_table(slide, rows)
    add_footer(slide, "Red = conflicts with Google model. Yellow = usable after rewrite. Green = reusable with minimal changes.")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, RGBColor(255, 255, 255))
    add_title(slide, "Platform map", "Where OptiFlow fits across search and AI-discovery providers", "Heatmap of platform fit, native monetization motion, and the solution changes needed for each provider.")
    provider_rows = [
        {
            "provider": "Microsoft Bing / Copilot",
            "fit": "Green",
            "motion": "Ads + Webmaster + Copilot",
            "fitment": "Best fit for citation diagnostics, landing-page quality, schema, IndexNow, and advertiser optimization.",
            "changes": "Add Bing AI Performance, IndexNow, Copilot citation, and Microsoft Advertising diagnostics. Avoid any pay-to-rank language.",
        },
        {
            "provider": "Google Search / AI Mode",
            "fit": "Yellow",
            "motion": "Ads + Merchant + Search",
            "fitment": "Strong fit for merchant readiness, feed-page consistency, landing-page quality, and AI-surface eligibility.",
            "changes": "Remove llms.txt. Shift to supported schema, crawlability, snippet controls, Merchant Center health, and ad/landing-page quality.",
        },
        {
            "provider": "OpenAI ChatGPT",
            "fit": "Green",
            "motion": "Shopping + merchant",
            "fitment": "Strong commerce fit for product feeds, retailer pages, checkout readiness, and structured product content.",
            "changes": "Build merchant-feed connectors, commerce taxonomy mapping, product-page QA, and shopping-readiness scoring instead of generic GEO claims.",
        },
        {
            "provider": "Perplexity",
            "fit": "Yellow",
            "motion": "Publisher + ads",
            "fitment": "Useful for citation readiness, publisher structure, source trust, and sponsored follow-up adjacency.",
            "changes": "Add citation-quality scoring, publisher-source trust indicators, and answerability diagnostics. Keep promises advisory only.",
        },
        {
            "provider": "Brave Search",
            "fit": "Yellow",
            "motion": "Search + AI Answers",
            "fitment": "External optimization tool only; can help answerability and source clarity, but weak case as a native Brave product.",
            "changes": "Frame as independent content-clarity tooling. Remove any assumption of platform-managed ranking controls or paid inclusion paths.",
        },
        {
            "provider": "DuckDuckGo",
            "fit": "Red",
            "motion": "Private search ads",
            "fitment": "Limited direct fit. Most monetization and index dependence routes through Microsoft anyway.",
            "changes": "Treat as indirect via Microsoft Advertising and Bing optimization. Do not create a DuckDuckGo-specific product story.",
        },
        {
            "provider": "Yahoo Search",
            "fit": "Yellow",
            "motion": "Portal + Bing-backed search",
            "fitment": "Indirect opportunity through Bing-powered search inventory and advertiser workflow extensions.",
            "changes": "Bundle under a Microsoft channel strategy rather than a standalone Yahoo proposition.",
        },
        {
            "provider": "Amazon Rufus / shopping AI",
            "fit": "Yellow",
            "motion": "Commerce discovery",
            "fitment": "Relevant for commerce content, product attributes, and retailer readiness where AI-assisted product discovery matters.",
            "changes": "Extend the model from search-copy optimization to retail catalog quality, product attributes, PDP structure, and commerce conversion signals.",
        },
    ]
    add_provider_matrix(slide, provider_rows)
    add_footer(slide, "Green = strongest platform fit. Yellow = viable with channel-specific repositioning. Red = low direct fit or mostly indirect route.")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, RGBColor(255, 255, 255))
    add_title(slide, "Rewrite guide", "Language leadership should use")
    add_card(slide, Inches(0.75), Inches(1.95), Inches(5.7), Inches(3.95), "Use", "AI-ready landing pages\nMachine-readable product content\nMerchant data integrity\nStructured content quality\nSearch and Shopping readiness\nFeed-page consistency\nAdvertiser diagnostics\nEligibility and quality improvement", GOOGLE_GREEN)
    add_card(slide, Inches(6.85), Inches(1.95), Inches(5.7), Inches(3.95), "Avoid", "Pay to rank in AI search\nGuaranteed AI visibility\nPrioritized in AI responses\nGoogle GEO algorithm\nllms.txt for Google optimization\nCitation uplift guarantee\nPromoted organically through spend", GOOGLE_RED)
    add_footer(slide, sources)

    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    add_title(slide, "Solution changes", "What must change for a Google-compatible pitch", "This is the bridge from the current solution to a credible Google offer.")
    add_card(slide, Inches(0.75), Inches(1.95), Inches(3.75), Inches(3.95), "1. Change the target", "Shift from 'maximize AI visibility' to 'improve landing-page quality, merchant integrity, and AI-surface readiness'.", GOOGLE_BLUE)
    add_card(slide, Inches(4.78), Inches(1.95), Inches(3.75), Inches(3.95), "2. Replace weak signals", "Drop llms.txt. Elevate supported schema, crawlability, feed-page consistency, snippet controls, and initial HTML facts.", GOOGLE_RED)
    add_card(slide, Inches(8.81), Inches(1.95), Inches(3.75), Inches(3.95), "3. Redesign measurement", "Measure feed health, landing-page diagnostics, structured-data validity, merchant issue reduction, and conversion quality.", GOOGLE_GREEN)
    add_text(slide, Inches(0.85), Inches(6.2), Inches(11.8), Inches(0.45), "Rebuilt around diagnostics and quality improvement for advertisers and merchants, OptiFlow complements Google's business model instead of conflicting with it.", size=16, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, sources)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    build_deck()
